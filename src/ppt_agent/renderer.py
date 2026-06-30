"""Render validated Slide IR decks to editable PowerPoint files."""

from __future__ import annotations

import hashlib
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from ppt_agent.layouts import is_template_layout
from ppt_agent.models import Deck, ImageElement, ShapeElement, SlideElement, TextElement, TextStyle
from ppt_agent.theme import Theme


VISUAL_VARIANT_COUNTS = {
    "title_slide": 3,
    "two_column": 2,
    "three_column": 2,
    "four_cards": 2,
    "metric_cards": 2,
    "process_flow": 2,
    "risk_matrix": 2,
    "key_takeaway": 2,
    "closing_slide": 2,
}


def _visual_variant_for_slide(
    deck: Deck,
    deck_slide,
    slide_index: int,
    *,
    variant_count: int | None = None,
) -> int:
    """Pick a reproducible renderer-only visual variant for a template slide."""

    count = variant_count or VISUAL_VARIANT_COUNTS.get(deck_slide.layout, 1)
    if count <= 1:
        return 0

    text_seed = "|".join(
        element.text
        for element in deck_slide.elements
        if isinstance(element, TextElement)
    )
    seed = "|".join(
        [
            deck.deck_id,
            deck.title,
            deck.theme_name or "",
            deck_slide.layout,
            deck_slide.slide_id,
            deck_slide.title,
            text_seed,
        ]
    )
    digest = hashlib.sha256(seed.encode("utf-8")).hexdigest()
    return (int(digest[:8], 16) + slide_index - 1) % count


def _rgb_color(hex_color: str) -> RGBColor:
    value = hex_color.removeprefix("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _bbox_args(element: SlideElement) -> tuple:
    bbox = element.bbox
    return Inches(bbox.x), Inches(bbox.y), Inches(bbox.width), Inches(bbox.height)


def _apply_text_style(run, style: TextStyle) -> None:
    font = run.font

    if style.font_family:
        font.name = style.font_family
    if style.font_size_pt:
        font.size = Pt(style.font_size_pt)
    if style.color:
        font.color.rgb = _rgb_color(style.color)

    font.bold = style.bold
    font.italic = style.italic


def _write_text_to_shape(shape, text: str, style: TextStyle) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        _apply_text_style(run, style)


def _render_text(slide, element: TextElement, theme: Theme) -> None:
    left, top, width, height = _bbox_args(element)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_style = element.style or theme.default_text_style
    _write_text_to_shape(shape, element.text, text_style)


def _render_shape(slide, element: ShapeElement, theme: Theme) -> None:
    left, top, width, height = _bbox_args(element)
    fill_color = element.style.fill_color if element.style and element.style.fill_color else theme.colors.surface
    stroke_color = (
        element.style.stroke_color if element.style and element.style.stroke_color else theme.colors.primary
    )
    stroke_width_pt = element.style.stroke_width_pt if element.style else None

    if element.shape == "line":
        shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top + height)
        shape.line.color.rgb = _rgb_color(stroke_color)
        if stroke_width_pt:
            shape.line.width = Pt(stroke_width_pt)
        return

    shape_type = MSO_SHAPE.RECTANGLE if element.shape == "rectangle" else MSO_SHAPE.OVAL
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_color(fill_color)
    shape.line.color.rgb = _rgb_color(stroke_color)
    if stroke_width_pt:
        shape.line.width = Pt(stroke_width_pt)


def _image_path(src: str, assets_dir: str | Path | None) -> Path:
    source_path = Path(src)
    if source_path.is_absolute() or assets_dir is None:
        return source_path
    return Path(assets_dir) / source_path


def _render_image_placeholder(slide, element: ImageElement, theme: Theme) -> None:
    left, top, width, height = _bbox_args(element)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_color(theme.colors.surface)
    shape.line.color.rgb = _rgb_color(theme.colors.muted_text)

    placeholder_style = TextStyle(
        font_family=theme.default_text_style.font_family or theme.fonts.body,
        font_size_pt=theme.default_text_style.font_size_pt or 14,
        color=theme.colors.muted_text,
    )
    _write_text_to_shape(shape, element.alt_text or "Image placeholder", placeholder_style)


def _render_image(slide, element: ImageElement, theme: Theme, assets_dir: str | Path | None) -> None:
    left, top, width, height = _bbox_args(element)
    path = _image_path(element.src, assets_dir)

    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return

    _render_image_placeholder(slide, element, theme)


def _theme_text_style(
    theme: Theme,
    *,
    font_size_pt: float,
    color: str | None = None,
    bold: bool = False,
    font_family: str | None = None,
) -> TextStyle:
    return TextStyle(
        font_family=font_family or theme.fonts.body,
        font_size_pt=font_size_pt,
        color=color or theme.colors.text,
        bold=bold,
    )


def _add_textbox(
    slide,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    style: TextStyle,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _write_text_to_shape(shape, text, style)
    return shape


def _add_rect(
    slide,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    fill_color: str,
    stroke_color: str | None = None,
    stroke_width_pt: float | None = None,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_color(fill_color)
    shape.line.color.rgb = _rgb_color(stroke_color or fill_color)
    if stroke_width_pt:
        shape.line.width = Pt(stroke_width_pt)
    return shape


def _template_texts(elements: list[SlideElement]) -> list[TextElement]:
    return [element for element in elements if isinstance(element, TextElement)]


def _title_and_body_texts(deck_slide) -> tuple[str, list[str]]:
    text_elements = _template_texts(deck_slide.elements)
    if not text_elements:
        return deck_slide.title, []

    title = text_elements[0].text or deck_slide.title
    body = [element.text for element in text_elements[1:] if element.text.strip()]
    return title, body


def _text_length_score(text: str) -> int:
    cjk_chars = sum(1 for char in text if "\u4e00" <= char <= "\u9fff")
    non_cjk_chars = sum(1 for char in text if not char.isspace()) - cjk_chars
    return cjk_chars + non_cjk_chars


def _split_overloaded_body_text(text: str, slot_count: int) -> list[str]:
    # This is a render-time safety pass, not semantic rewriting. The aim is to
    # keep one oversized body from monopolizing all cards in a grid layout.
    lines = [
        line.strip(" -•\t")
        for line in text.splitlines()
        if line.strip(" -•\t")
    ]
    if len(lines) >= slot_count:
        return lines[:slot_count]

    sentence_parts = [
        part.strip(" -•\t")
        for part in re.split(r"[。；;]\s*", text)
        if part.strip(" -•\t")
    ]
    if len(sentence_parts) >= slot_count:
        return sentence_parts[:slot_count]

    comma_parts = [
        part.strip(" -•\t")
        for part in re.split(r"[，,、]\s*", text)
        if part.strip(" -•\t")
    ]
    if len(comma_parts) >= slot_count:
        return comma_parts[:slot_count]

    return lines or sentence_parts or comma_parts


def _rebalance_body_texts(body_texts: list[str], slot_count: int) -> list[str]:
    cleaned = [text for text in body_texts if text.strip()]
    if slot_count <= 1 or len(cleaned) != 1:
        return body_texts

    only_text = cleaned[0]
    if _text_length_score(only_text) < 60:
        return body_texts

    split_items = _split_overloaded_body_text(only_text, slot_count)
    if len(split_items) < 2:
        return body_texts

    return split_items[:slot_count]


def _slot_texts(body_texts: list[str], slot_count: int, fallback: str = "") -> list[str]:
    body_texts = _rebalance_body_texts(body_texts, slot_count)
    slots = body_texts[:slot_count]
    if len(body_texts) > slot_count:
        slots[-1] = slots[-1] + "\n" + "\n".join(body_texts[slot_count:])
    while len(slots) < slot_count:
        slots.append(fallback)
    return slots


def _title_slide_title_metrics(title: str) -> tuple[float, float]:
    word_count = len(title.split())
    char_count = len(title)

    if word_count > 16 or char_count > 96:
        return 24, 1.48
    if word_count > 11 or char_count > 64:
        return 28, 1.16
    return 36, 0.88


def _title_slide_subtitle_metrics(subtitle: str) -> tuple[float, float]:
    char_count = len(subtitle)
    word_count = len(subtitle.split())

    if word_count > 24 or char_count > 110:
        return 13.2, 1.05
    if word_count > 16 or char_count > 76:
        return 14.5, 0.9
    return 16, 0.72


def _split_heading_body(text: str) -> tuple[str, str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return "", ""
    if len(lines) == 1:
        return lines[0], ""
    return lines[0], "\n".join(lines[1:])


def _short_phrase(text: str, max_chars: int) -> str:
    normalized = " ".join(text.split())
    if len(normalized) <= max_chars:
        return normalized

    words: list[str] = []
    for word in normalized.split():
        candidate = " ".join([*words, word])
        if len(candidate) > max_chars:
            break
        words.append(word)

    if words:
        return " ".join(words)
    return normalized[:max_chars].rstrip()


def _safe_text(text: str, max_chars: int) -> str:
    normalized = " ".join(text.split())
    if len(normalized) <= max_chars:
        return normalized
    if max_chars <= 3:
        return normalized[:max_chars]
    return normalized[: max_chars - 3].rstrip() + "..."


def _shrink_font_size_for_text(text: str, base_size_pt: float, *, min_size_pt: float, medium: int, long: int) -> float:
    length_score = _text_length_score(text)
    if length_score > long:
        return min_size_pt
    if length_score > medium:
        return max(min_size_pt, base_size_pt - 1.8)
    return base_size_pt


def _keywords_from_title(title: str) -> list[str]:
    stop_words = {"a", "an", "and", "for", "how", "in", "of", "the", "to", "with"}
    words = [
        word.strip(" .,:;!?")
        for word in title.split()
        if word.strip(" .,:;!?") and word.strip(" .,:;!?").lower() not in stop_words
    ]
    keywords = [word for word in words if len(word) <= 14][:3]
    return keywords or ["Overview"]


def _cover_keywords(title: str, body: list[str], deck: Deck) -> list[str]:
    source = " ".join([title, *body])
    if any("\u4e00" <= char <= "\u9fff" for char in source):
        known_terms = [
            "Agent",
            "工作流",
            "边界",
            "评估",
            "风险",
            "需求",
            "产品",
        ]
        matched_terms = [term for term in known_terms if term.lower() in source.lower()]
        if matched_terms:
            return matched_terms[:3]
        candidates = [
            token.strip(" ，。,:;!?")
            for token in source.replace("/", " ").replace("｜", " ").split()
            if token.strip(" ，。,:;!?")
        ]
        compact = [token for token in candidates if 2 <= len(token) <= 8]
        if compact:
            return compact[:3]
        cjk_chars = [char for char in title if "\u4e00" <= char <= "\u9fff"]
        if len(cjk_chars) >= 6:
            return ["".join(cjk_chars[:4]), "".join(cjk_chars[4:8]), "".join(cjk_chars[8:12])][:3]
        return [_surface_label(deck, "Product", "产品"), _surface_label(deck, "Workflow", "工作流"), _surface_label(deck, "AI", "AI")]

    return _keywords_from_title(title)


def _is_sparse_card_text(text: str) -> bool:
    heading, body = _split_heading_body(text)
    word_count = len(" ".join([heading, body]).split())
    return word_count <= 7


def _body_lines(body_texts: list[str]) -> list[str]:
    return [
        line.strip()
        for text in body_texts
        for line in text.splitlines()
        if line.strip()
    ]


def _compact_lines(text: str, limit: int | None = None) -> list[str]:
    lines = [line.strip(" -\u2022\t") for line in text.splitlines() if line.strip(" -\u2022\t")]
    if not lines and text.strip():
        lines = [text.strip()]
    return lines[:limit] if limit is not None else lines


def _heading_body_with_fallback(text: str, fallback: str) -> tuple[str, str]:
    heading, body = _split_heading_body(text)
    if not heading:
        heading = fallback
    if not body:
        body = fallback
    return heading, body


def _strip_leading_index(text: str) -> str:
    return text.strip().lstrip("0123456789.、)） ").strip()


def _action_fallback(deck: Deck, heading: str) -> str:
    if _is_english_language(deck):
        if "risk" in heading.lower() or "metric" in heading.lower():
            return "Record failure samples and review checkpoints."
        if "workflow" in heading.lower() or "flow" in heading.lower():
            return "Design one confirmation point for the workflow."
        options = [
            "Define the boundary and next checkpoint.",
            "Design one confirmation point for the workflow.",
            "Record failure samples and review checkpoints.",
            "Choose one launch metric and review it weekly.",
        ]
        return options[sum(ord(char) for char in heading) % len(options)]
    if "边界" in heading:
        return "明确边界。"
    if "闭环" in heading or "工作流" in heading:
        return "设计确认点。"
    if "风险" in heading or "指标" in heading:
        return "记录失败样本。"
    options = ["明确边界。", "设计确认点。", "记录失败样本。", "建立评估指标。"]
    return options[sum(ord(char) for char in heading) % len(options)]


def _action_pairs_from_body(body_texts: list[str], deck: Deck, limit: int = 3) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    for text in body_texts:
        heading, body = _split_heading_body(text)
        heading = _strip_leading_index(heading)
        if not heading:
            continue
        if not body:
            body = _action_fallback(deck, heading)
        pairs.append((heading, body))
        if len(pairs) >= limit:
            return pairs

    lines = [_strip_leading_index(line) for line in _body_lines(body_texts)]
    lines = [line for line in lines if line]
    index = 0
    while len(pairs) < limit and index < len(lines):
        heading = lines[index]
        body = lines[index + 1] if index + 1 < len(lines) else _action_fallback(deck, heading)
        pairs.append((heading, body))
        index += 2
    return pairs


def _risk_cells_from_text(text: str, deck: Deck) -> tuple[str, str, str]:
    lines = [line.strip(" -•\t") for line in text.splitlines() if line.strip(" -•\t")]
    if len(lines) == 1 and "|" in lines[0]:
        lines = [part.strip() for part in lines[0].split("|") if part.strip()]
    cells: list[str] = []
    for line in lines[:3]:
        cleaned = line
        for prefix in ["Risk:", "Impact:", "Mitigation:", "风险：", "影响：", "缓解措施：", "缓解：", "Risk：", "Impact：", "Mitigation："]:
            if cleaned.lower().startswith(prefix.lower()):
                cleaned = cleaned[len(prefix) :].strip()
                break
        cells.append(cleaned)
    while len(cells) < 3:
        cells.append("")
    if not cells[2]:
        cells[2] = _surface_label(
            deck,
            "Set permission boundaries, human review, and operation logs.",
            "设置权限边界、人工确认和操作日志。",
        )
    return cells[0] or " ", cells[1] or " ", cells[2]


def _is_english_language(deck: Deck) -> bool:
    marker = f"{deck.title} {deck.theme_name or ''}".lower()
    if "lang:en" in marker or "language:en" in marker:
        return True

    deck_text = " ".join(
        element.text
        for slide in deck.slides
        for element in slide.elements
        if isinstance(element, TextElement)
    )
    has_cjk_text = any("\u4e00" <= char <= "\u9fff" for char in deck_text)
    return not has_cjk_text


def _surface_label(deck: Deck, english: str, chinese: str) -> str:
    return english if _is_english_language(deck) else chinese


def _render_heading_body_card(
    slide,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    theme: Theme,
    accent_color: str | None = None,
    number: int | None = None,
    label: str | None = None,
    heading_size_pt: float = 18,
    body_size_pt: float = 13.5,
) -> None:
    heading, body = _split_heading_body(text)
    # Clamp heading/body length and font size before placing text into a
    # fixed-size card. The renderer should fail soft here, not overflow.
    heading = _safe_text(_short_phrase(heading, 44), 44)
    body_limit = 68 if height <= 1.35 else 96 if height <= 2.0 else 128
    body = _safe_text(body, body_limit)
    heading_size_pt = _shrink_font_size_for_text(
        heading,
        heading_size_pt,
        min_size_pt=12.0,
        medium=24,
        long=36,
    )
    body_size_pt = _shrink_font_size_for_text(
        body,
        body_size_pt,
        min_size_pt=8.8,
        medium=48,
        long=76,
    )
    accent = accent_color or theme.colors.primary
    padding_x = 0.24

    _add_rect(
        slide,
        x=x,
        y=y,
        width=width,
        height=height,
        fill_color=theme.colors.background,
        stroke_color=theme.colors.surface,
        stroke_width_pt=1.0,
    )
    _add_rect(slide, x=x, y=y, width=width, height=0.08, fill_color=accent)

    if number is not None:
        chip_width = 0.62
        chip_height = 0.28
        chip_x = x + width - padding_x - chip_width
        chip_y = y + 0.24
        _add_rect(
            slide,
            x=chip_x,
            y=chip_y,
            width=chip_width,
            height=chip_height,
            fill_color=theme.colors.surface,
            stroke_color=theme.colors.surface,
        )
        _add_textbox(
            slide,
            x=chip_x + 0.08,
            y=chip_y + 0.03,
            width=chip_width - 0.16,
            height=0.18,
            text=f"{number:02d}",
            style=_theme_text_style(theme, font_size_pt=8.5, color=accent, bold=True),
        )

        if label:
            _add_textbox(
                slide,
                x=x + padding_x,
                y=y + 0.22,
                width=max(0.7, width - padding_x * 2 - chip_width - 0.12),
                height=0.22,
                text=label,
                style=_theme_text_style(theme, font_size_pt=7.5, color=accent, bold=True),
            )

    _add_textbox(
        slide,
        x=x + padding_x,
        y=y + 0.58,
        width=width - padding_x * 2,
        height=0.45,
        text=heading,
        style=_theme_text_style(theme, font_size_pt=heading_size_pt, color=theme.colors.text, bold=True),
    )
    if body:
        _add_textbox(
            slide,
            x=x + padding_x,
            y=y + 1.16,
            width=width - padding_x * 2,
            height=max(0.4, height - 1.38),
            text=body,
            style=_theme_text_style(theme, font_size_pt=body_size_pt, color=theme.colors.muted_text),
        )


def _render_title_slide_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    subtitle = body[0] if body else deck_slide.title
    keywords = _cover_keywords(title, body, deck)
    margin = 0.8
    title_y = 1.34
    title_font_size, title_height = _title_slide_title_metrics(title)
    subtitle_font_size, subtitle_height = _title_slide_subtitle_metrics(subtitle)
    subtitle_y = title_y + title_height + 0.34
    accent_y = subtitle_y + subtitle_height + 0.34
    right_x = 8.96
    right_width = deck.canvas_width_in - right_x - 0.72
    label = _surface_label(deck, "Product Briefing", "技术产品分享")

    _add_rect(slide, x=0, y=0, width=deck.canvas_width_in, height=deck.canvas_height_in, fill_color=theme.colors.background)
    if variant == 2:
        _add_rect(slide, x=9.3, y=0.78, width=2.55, height=5.7, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
        _add_rect(slide, x=10.0, y=1.24, width=2.25, height=0.24, fill_color=theme.colors.primary, stroke_color=theme.colors.primary)
        _add_rect(slide, x=9.65, y=2.1, width=1.1, height=1.1, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
        _add_rect(slide, x=10.7, y=3.08, width=1.28, height=1.28, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
        _add_rect(slide, x=9.88, y=5.18, width=1.72, height=0.12, fill_color=theme.colors.accent)
    elif variant == 1:
        _add_rect(slide, x=right_x + 0.18, y=0.86, width=right_width - 0.12, height=4.98, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
        _add_rect(slide, x=right_x + 0.62, y=1.42, width=right_width - 0.96, height=0.08, fill_color=theme.colors.primary)
        _add_rect(slide, x=right_x + 0.62, y=5.18, width=1.2, height=0.12, fill_color=theme.colors.accent)
    else:
        _add_rect(slide, x=right_x + 0.1, y=0.7, width=right_width, height=5.95, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
        _add_rect(slide, x=right_x + 0.48, y=1.08, width=right_width - 0.72, height=2.08, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
        _add_rect(slide, x=right_x + right_width - 0.92, y=0.36, width=0.72, height=2.1, fill_color=theme.colors.primary, stroke_color=theme.colors.primary)
        _add_rect(slide, x=right_x + 0.18, y=5.58, width=1.35, height=0.14, fill_color=theme.colors.accent)
    _add_rect(slide, x=margin, y=0.92, width=0.9, height=0.08, fill_color=theme.colors.accent)
    _add_rect(slide, x=margin, y=accent_y, width=3.6, height=0.08, fill_color=theme.colors.primary)

    _add_textbox(
        slide,
        x=margin,
        y=0.64,
        width=3.2,
        height=0.24,
        text=label,
        style=_theme_text_style(theme, font_size_pt=9.5, color=theme.colors.primary, bold=True),
    )
    _add_textbox(
        slide,
        x=margin,
        y=title_y,
        width=7.85,
        height=title_height,
        text=title,
        style=_theme_text_style(theme, font_size_pt=title_font_size, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_textbox(
        slide,
        x=margin,
        y=subtitle_y,
        width=7.58,
        height=subtitle_height,
        text=subtitle,
        style=_theme_text_style(theme, font_size_pt=subtitle_font_size, color=theme.colors.muted_text),
    )

    hero_keyword = _safe_text(keywords[0] if keywords else title, 16)
    side_label = _surface_label(deck, "Technical product perspective", "技术产品视角")
    if variant == 2:
        _add_textbox(
            slide,
            x=9.62,
            y=5.42,
            width=2.28,
            height=0.24,
            text=side_label,
            style=_theme_text_style(theme, font_size_pt=9.6, color=theme.colors.primary, bold=True),
        )
    elif variant == 1:
        _add_textbox(
            slide,
            x=right_x + 0.62,
            y=1.72,
            width=right_width - 1.08,
            height=0.32,
            text=side_label,
            style=_theme_text_style(theme, font_size_pt=10.6, color=theme.colors.muted_text),
        )
        for index, keyword in enumerate(keywords[:2]):
            chip_width = min(1.45, max(0.9, 0.54 + len(keyword) * 0.12))
            chip_x = right_x + 0.62
            chip_y = 2.48 + index * 0.58
            _add_rect(
                slide,
                x=chip_x,
                y=chip_y,
                width=chip_width,
                height=0.36,
                fill_color=theme.colors.background,
                stroke_color=theme.colors.background,
            )
            _add_textbox(
                slide,
                x=chip_x + 0.14,
                y=chip_y + 0.09,
                width=chip_width - 0.28,
                height=0.14,
                text=_safe_text(keyword, 12),
                style=_theme_text_style(theme, font_size_pt=8.4, color=theme.colors.primary, bold=True),
            )
    else:
        _add_textbox(
            slide,
            x=right_x + 0.62,
            y=1.62,
            width=right_width - 1.1,
            height=0.82,
            text=hero_keyword,
            style=_theme_text_style(theme, font_size_pt=30, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
        )
        _add_textbox(
            slide,
            x=right_x + 0.64,
            y=2.58,
            width=right_width - 1.16,
            height=0.36,
            text=side_label,
            style=_theme_text_style(theme, font_size_pt=11.5, color=theme.colors.muted_text),
        )
        for index, keyword in enumerate(keywords[1:3]):
            chip_width = min(1.35, max(0.82, 0.46 + len(keyword) * 0.12))
            chip_x = right_x + 0.64 + index * 1.48
            _add_rect(
                slide,
                x=chip_x,
                y=3.54,
                width=chip_width,
                height=0.34,
                fill_color=theme.colors.background,
                stroke_color=theme.colors.surface,
                stroke_width_pt=0.8,
            )
            _add_textbox(
                slide,
                x=chip_x + 0.12,
                y=3.64,
                width=chip_width - 0.24,
                height=0.14,
                text=_safe_text(keyword, 12),
                style=_theme_text_style(theme, font_size_pt=8.1, color=theme.colors.primary, bold=True),
            )


def _render_section_divider_template(slide, deck_slide, deck: Deck, theme: Theme) -> None:
    title, body = _title_and_body_texts(deck_slide)
    keywords = _cover_keywords(title, body, deck)
    label = _surface_label(deck, "Section", "章节过渡")

    _add_rect(slide, x=9.35, y=0.78, width=2.9, height=5.7, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
    _add_rect(slide, x=10.0, y=1.42, width=1.45, height=1.45, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
    _add_rect(slide, x=10.78, y=4.75, width=1.18, height=0.12, fill_color=theme.colors.accent)
    _add_rect(slide, x=0.75, y=1.05, width=0.12, height=5.4, fill_color=theme.colors.primary)
    _add_textbox(
        slide,
        x=1.15,
        y=1.68,
        width=2.4,
        height=0.24,
        text=label,
        style=_theme_text_style(theme, font_size_pt=9.5, color=theme.colors.primary, bold=True),
    )
    _add_textbox(
        slide,
        x=1.15,
        y=2.1,
        width=8.0,
        height=0.9,
        text=title,
        style=_theme_text_style(theme, font_size_pt=34, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    if body:
        _add_textbox(
            slide,
            x=1.17,
            y=3.18,
            width=7.65,
            height=0.7,
            text=body[0],
            style=_theme_text_style(theme, font_size_pt=18, color=theme.colors.muted_text),
        )

    for index, keyword in enumerate(keywords[:3]):
        chip_width = min(2.1, max(1.0, 0.42 + len(keyword) * 0.16))
        _add_rect(
            slide,
            x=1.17 + index * 2.2,
            y=4.62,
            width=chip_width,
            height=0.34,
            fill_color=theme.colors.surface,
            stroke_color=theme.colors.surface,
        )
        _add_textbox(
            slide,
            x=1.31 + index * 2.2,
            y=4.71,
            width=chip_width - 0.28,
            height=0.14,
            text=_safe_text(keyword, 16),
            style=_theme_text_style(theme, font_size_pt=8.2, color=theme.colors.primary, bold=True),
        )


def _render_column_template(slide, deck_slide, deck: Deck, theme: Theme, column_count: int, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    margin = 0.65
    gutter = 0.28
    top = 1.62
    card_height = 3.95 if column_count == 2 else 3.65
    column_width = (deck.canvas_width_in - (margin * 2) - gutter * (column_count - 1)) / column_count
    slots = _slot_texts(body, column_count, fallback=" ")
    if column_count == 2 and all(_is_sparse_card_text(text) for text in slots):
        top = 1.84
        card_height = 2.45
        variant = 0

    if variant == 1:
        _add_rect(slide, x=0, y=0, width=0.18, height=deck.canvas_height_in, fill_color=theme.colors.primary)
        _add_textbox(
            slide,
            x=margin,
            y=0.48,
            width=deck.canvas_width_in - margin * 2,
            height=0.55,
            text=title,
            style=_theme_text_style(theme, font_size_pt=25, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
        )
        _add_rect(slide, x=margin, y=1.17, width=1.42, height=0.06, fill_color=theme.colors.accent)

        lead_width = 4.0 if column_count == 3 else 4.8
        _render_heading_body_card(
            slide,
            x=margin,
            y=1.55,
            width=lead_width,
            height=4.35,
            text=slots[0],
            theme=theme,
            accent_color=theme.colors.primary,
            number=1,
            label=_surface_label(deck, "Point of view", "核心观点"),
            heading_size_pt=20,
            body_size_pt=14,
        )

        right_x = margin + lead_width + 0.42
        right_width = deck.canvas_width_in - right_x - margin
        supporting_slots = slots[1:] or [" "]
        row_height = (4.35 - 0.28 * (len(supporting_slots) - 1)) / len(supporting_slots)
        for index, text in enumerate(supporting_slots, start=2):
            y = 1.55 + (index - 2) * (row_height + 0.28)
            _render_heading_body_card(
                slide,
                x=right_x,
                y=y,
                width=right_width,
                height=row_height,
                text=text,
                theme=theme,
                accent_color=theme.colors.secondary if index % 2 == 0 else theme.colors.accent,
                number=index,
                label=_surface_label(deck, "Evidence", "支撑"),
                heading_size_pt=16.2,
                body_size_pt=12.2,
            )
        return

    _add_textbox(
        slide,
        x=margin,
        y=0.48,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_rect(slide, x=margin, y=1.16, width=2.2, height=0.06, fill_color=theme.colors.primary)

    for index, text in enumerate(slots):
        x = margin + index * (column_width + gutter)
        _render_heading_body_card(
            slide,
            x=x,
            y=top,
            width=column_width,
            height=card_height,
            text=text,
            theme=theme,
            accent_color=theme.colors.primary if index % 2 == 0 else theme.colors.secondary,
            number=index + 1,
            label=_surface_label(deck, "Insight", "洞察"),
            heading_size_pt=18,
            body_size_pt=13.5,
        )


def _render_four_cards_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    margin = 0.72
    gutter_x = 0.34
    gutter_y = 0.34
    card_width = (deck.canvas_width_in - margin * 2 - gutter_x) / 2
    card_height = 1.9
    top = 1.48
    slots = _slot_texts(body, 4, fallback=" ")

    if variant == 1:
        _add_textbox(
            slide,
            x=margin,
            y=0.48,
            width=deck.canvas_width_in - margin * 2,
            height=0.55,
            text=title,
            style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
        )
        _add_rect(slide, x=margin, y=1.16, width=1.4, height=0.06, fill_color=theme.colors.accent)
        lead_width = 4.75
        _render_heading_body_card(
            slide,
            x=margin,
            y=1.5,
            width=lead_width,
            height=4.35,
            text=slots[0],
            theme=theme,
            accent_color=theme.colors.primary,
            number=1,
            label=_surface_label(deck, "Anchor", "主线"),
            heading_size_pt=20,
            body_size_pt=14,
        )
        stack_x = margin + lead_width + 0.42
        stack_width = deck.canvas_width_in - stack_x - margin
        for index, text in enumerate(slots[1:], start=2):
            _render_heading_body_card(
                slide,
                x=stack_x,
                y=1.5 + (index - 2) * 1.48,
                width=stack_width,
                height=1.18,
                text=text,
                theme=theme,
                accent_color=[theme.colors.secondary, theme.colors.accent, theme.colors.primary][index - 2],
                number=index,
                label=_surface_label(deck, "Support", "支撑"),
                heading_size_pt=15.2,
                body_size_pt=11.4,
            )
        return

    _add_textbox(
        slide,
        x=margin,
        y=0.48,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_rect(slide, x=margin, y=1.16, width=2.2, height=0.06, fill_color=theme.colors.primary)

    accents = [theme.colors.primary, theme.colors.secondary, theme.colors.accent, theme.colors.primary]
    for index, text in enumerate(slots):
        row = index // 2
        column = index % 2
        _render_heading_body_card(
            slide,
            x=margin + column * (card_width + gutter_x),
            y=top + row * (card_height + gutter_y),
            width=card_width,
            height=card_height,
            text=text,
            theme=theme,
            accent_color=accents[index],
            number=index + 1,
            label=_surface_label(deck, "Action", "行动"),
            heading_size_pt=17,
            body_size_pt=12.8,
        )


def _render_metric_cards_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    margin = 0.72
    card_count = 4 if len(body) >= 4 else 3
    gutter = 0.32
    slots = _slot_texts(body, card_count, fallback=" ")

    _add_textbox(
        slide,
        x=margin,
        y=0.55,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )

    if variant == 1:
        _add_rect(slide, x=margin, y=1.17, width=1.35, height=0.06, fill_color=theme.colors.accent)
        for index, text in enumerate(slots):
            heading, body_text = _heading_body_with_fallback(
                text,
                _surface_label(deck, "Track this signal.", "持续跟踪这个信号。"),
            )
            if card_count == 4:
                card_width = (deck.canvas_width_in - margin * 2 - gutter) / 2
                card_height = 1.72
                row = index // 2
                column = index % 2
                x = margin + column * (card_width + gutter)
                y = 1.46 + row * (card_height + 0.34)
            else:
                card_width = (deck.canvas_width_in - margin * 2 - gutter) / 2
                card_height = 1.52 if index else 1.86
                x = margin if index == 0 else margin + (index - 1) * (card_width + gutter)
                y = 1.45 if index == 0 else 3.72
                if index == 0:
                    card_width = deck.canvas_width_in - margin * 2
            accent = [theme.colors.primary, theme.colors.secondary, theme.colors.accent, theme.colors.primary][index]
            _add_rect(slide, x=x, y=y, width=card_width, height=card_height, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
            _add_rect(slide, x=x, y=y, width=0.1, height=card_height, fill_color=accent)
            _add_textbox(
                slide,
                x=x + 0.32,
                y=y + 0.28,
                width=0.42,
                height=0.18,
                text=f"{index + 1:02d}",
                style=_theme_text_style(theme, font_size_pt=8.5, color=accent, bold=True),
            )
            _add_textbox(
                slide,
                x=x + 0.86,
                y=y + 0.22,
                width=card_width - 1.12,
                height=0.34,
                text=_safe_text(heading, 38),
                style=_theme_text_style(theme, font_size_pt=17 if card_count == 4 else 18, color=theme.colors.text, bold=True),
            )
            _add_textbox(
                slide,
                x=x + 0.86,
                y=y + 0.82,
                width=card_width - 1.12,
                height=0.34,
                text=_safe_text(body_text, 62),
                style=_theme_text_style(theme, font_size_pt=11.8 if card_count == 4 else 12.4, color=theme.colors.muted_text),
            )
        return

    for index, text in enumerate(slots):
        if card_count == 4:
            card_width = (deck.canvas_width_in - margin * 2 - gutter) / 2
            card_height = 1.72
            row = index // 2
            column = index % 2
            x = margin + column * (card_width + gutter)
            y = 1.46 + row * (card_height + 0.34)
            inset = 0.12
        else:
            card_width = (deck.canvas_width_in - margin * 2 - gutter * 2) / 3
            card_height = 2.75
            x = margin + index * (card_width + gutter)
            y = 1.68
            inset = 0.28
        _render_heading_body_card(
            slide,
            x=x + inset,
            y=y,
            width=card_width - inset * 2,
            height=card_height,
            text=text,
            theme=theme,
            accent_color=theme.colors.primary,
            number=index + 1,
            label=_surface_label(deck, "Priority", "重点"),
            heading_size_pt=17 if card_count == 4 else 18,
            body_size_pt=12.6 if card_count == 4 else 14,
        )


def _render_comparison_matrix_template(slide, deck_slide, deck: Deck, theme: Theme) -> None:
    title, body = _title_and_body_texts(deck_slide)
    if len(body) < 2:
        _render_column_template(slide, deck_slide, deck, theme, column_count=2, variant=0)
        return
    slots = _slot_texts(body, 3, fallback=" ")
    left_heading, left_body = _split_heading_body(slots[0])
    right_heading, right_body = _split_heading_body(slots[1])
    decision_rule = slots[2].strip()
    if not left_heading.strip() or not right_heading.strip():
        _render_column_template(slide, deck_slide, deck, theme, column_count=2, variant=0)
        return
    margin = 0.72
    table_x = margin
    table_y = 1.48
    table_width = deck.canvas_width_in - margin * 2
    dimension_width = 2.25
    column_width = (table_width - dimension_width) / 2
    header_height = 0.62
    row_height = 0.58

    _add_textbox(
        slide,
        x=margin,
        y=0.5,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_rect(slide, x=margin, y=1.17, width=2.2, height=0.06, fill_color=theme.colors.primary)

    _add_rect(
        slide,
        x=table_x,
        y=table_y,
        width=table_width,
        height=header_height,
        fill_color=theme.colors.surface,
        stroke_color=theme.colors.surface,
    )
    _add_textbox(
        slide,
        x=table_x + 0.18,
        y=table_y + 0.2,
        width=dimension_width - 0.36,
        height=0.18,
        text=_surface_label(deck, "Compare", "对比"),
        style=_theme_text_style(theme, font_size_pt=9.2, color=theme.colors.primary, bold=True),
    )
    for x, heading, accent, label in [
        (table_x + dimension_width, left_heading, theme.colors.primary, _safe_text(left_heading, 24) or " "),
        (table_x + dimension_width + column_width, right_heading, theme.colors.secondary, _safe_text(right_heading, 24) or " "),
    ]:
        _add_textbox(
            slide,
            x=x + 0.18,
            y=table_y + 0.13,
            width=0.88,
            height=0.18,
            text=label,
            style=_theme_text_style(theme, font_size_pt=11.5, color=accent, bold=True),
        )

    left_points = _compact_lines(left_body, 5)
    right_points = _compact_lines(right_body, 5)
    row_count = max(3, min(5, max(len(left_points), len(right_points))))
    while len(left_points) < row_count:
        left_points.append("")
    while len(right_points) < row_count:
        right_points.append("")

    for row_index in range(row_count):
        left_value = _safe_text(left_points[row_index], 54) or " "
        right_value = _safe_text(right_points[row_index], 54) or " "
        if not left_value.strip() and not right_value.strip():
            continue
        y = table_y + header_height + row_index * row_height
        fill = theme.colors.background if row_index % 2 == 0 else theme.colors.surface
        _add_rect(
            slide,
            x=table_x,
            y=y,
            width=table_width,
            height=row_height,
            fill_color=fill,
            stroke_color=theme.colors.surface,
            stroke_width_pt=0.6,
        )
        _add_textbox(
            slide,
            x=table_x + 0.18,
            y=y + 0.17,
            width=dimension_width - 0.36,
            height=0.18,
            text=_surface_label(deck, "Point", "判断"),
            style=_theme_text_style(theme, font_size_pt=10.0, color=theme.colors.text, bold=True),
        )
        _add_textbox(
            slide,
            x=table_x + dimension_width + 0.18,
            y=y + 0.15,
            width=column_width - 0.36,
            height=0.24,
            text=left_value,
            style=_theme_text_style(theme, font_size_pt=10.8, color=theme.colors.muted_text),
        )
        _add_textbox(
            slide,
            x=table_x + dimension_width + column_width + 0.18,
            y=y + 0.15,
            width=column_width - 0.36,
            height=0.24,
            text=right_value,
            style=_theme_text_style(theme, font_size_pt=10.8, color=theme.colors.muted_text),
        )

    _add_rect(
        slide,
        x=table_x + dimension_width,
        y=table_y,
        width=0.02,
        height=header_height + row_count * row_height,
        fill_color=theme.colors.background,
    )
    _add_rect(
        slide,
        x=table_x + dimension_width + column_width,
        y=table_y,
        width=0.02,
        height=header_height + row_count * row_height,
        fill_color=theme.colors.background,
    )

    if decision_rule:
        _add_rect(
            slide,
            x=margin + 0.6,
            y=5.95,
            width=deck.canvas_width_in - margin * 2 - 1.2,
            height=0.56,
            fill_color=theme.colors.surface,
            stroke_color=theme.colors.surface,
        )
        _add_textbox(
            slide,
            x=margin + 0.92,
            y=6.09,
            width=deck.canvas_width_in - margin * 2 - 1.84,
            height=0.28,
            text=_safe_text(decision_rule, 88),
            style=_theme_text_style(theme, font_size_pt=14, color=theme.colors.text, bold=True),
        )


def _render_process_flow_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    steps = _slot_texts(body, 3, fallback=" ") if len(body) < 3 else body[:5]
    margin = 0.7

    _add_textbox(
        slide,
        x=margin,
        y=0.5,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_rect(slide, x=margin, y=1.17, width=2.2, height=0.06, fill_color=theme.colors.primary)

    if len(steps) <= 3:
        columns_by_row = [len(steps)]
    elif len(steps) == 4:
        columns_by_row = [4]
    else:
        columns_by_row = [3, 2]

    card_height = 1.55 if len(columns_by_row) > 1 else 2.15
    row_gap = 0.52
    top = 1.72 if len(columns_by_row) > 1 else 2.05
    step_index = 0
    rows: list[list[tuple[int, float, float, float, float, str]]] = []

    for row_index, column_count in enumerate(columns_by_row):
        row_steps = steps[step_index : step_index + column_count]
        gutter = 0.34
        available_width = deck.canvas_width_in - margin * 2
        step_width = (available_width - gutter * (column_count - 1)) / column_count
        y = top + row_index * (card_height + row_gap)
        row_offset = 0.0
        if len(row_steps) == 2 and len(steps) == 5:
            step_width = (available_width - gutter * 2) / 3
            row_offset = (available_width - (step_width * 2 + gutter)) / 2

        row_positions: list[tuple[int, float, float, float, float, str]] = []
        for column_index, text in enumerate(row_steps):
            index = step_index + column_index + 1
            x = margin + row_offset + column_index * (step_width + gutter)
            row_positions.append((index, x, y, step_width, card_height, text))
        rows.append(row_positions)
        step_index += column_count

    if variant == 1:
        for row_positions in rows:
            first_step = row_positions[0]
            last_step = row_positions[-1]
            _index, x, y, width, height, _text = first_step
            _last_index, last_x, _last_y, last_width, _last_height, _last_text = last_step
            _add_rect(
                slide,
                x=x + 0.32,
                y=y + height + 0.18,
                width=(last_x + last_width) - x - 0.64,
                height=0.06,
                fill_color=theme.colors.surface,
            )

    def add_connector(x1: float, y1: float, x2: float, y2: float) -> None:
        if abs(y1 - y2) < 0.001 and x2 < x1:
            x1, x2 = x2, x1
        if abs(x1 - x2) < 0.001 and y2 < y1:
            y1, y2 = y2, y1
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        connector.line.color.rgb = _rgb_color(theme.colors.surface)
        connector.line.width = Pt(1.0)

    for row_positions in rows:
        for current, following in zip(row_positions, row_positions[1:]):
            _index, x, y, width, _height, _text = current
            _next_index, next_x, _next_y, _next_width, _next_height, _next_text = following
            connector_y = y + 0.36
            add_connector(x + width + 0.06, connector_y, next_x - 0.06, connector_y)

    if len(rows) == 2 and len(steps) == 5:
        last_top = rows[0][-1]
        first_bottom = rows[1][0]
        _index, x, y, width, height, _text = last_top
        _bottom_index, bottom_x, bottom_y, _bottom_width, _bottom_height, _bottom_text = first_bottom
        turn_x = deck.canvas_width_in - margin - 0.16
        gap_y = y + height + row_gap * 0.46
        add_connector(x + width + 0.06, gap_y, turn_x, gap_y)
        add_connector(turn_x, gap_y, turn_x, bottom_y - 0.18)
        add_connector(turn_x, bottom_y - 0.18, bottom_x + 0.28, bottom_y - 0.18)

    for row_positions in rows:
        for index, x, y, step_width, card_height, text in row_positions:
            accent = theme.colors.primary if index % 2 else theme.colors.secondary
            _add_rect(
                slide,
                x=x,
                y=y,
                width=step_width,
                height=card_height,
                fill_color=theme.colors.background,
                stroke_color=theme.colors.surface,
                stroke_width_pt=1.0,
            )
            _add_rect(
                slide,
                x=x + 0.22,
                y=y + 0.18,
                width=0.58,
                height=0.36,
                fill_color=accent,
                stroke_color=accent,
            )
            _add_textbox(
                slide,
                x=x + 0.34,
                y=y + 0.26,
                width=0.34,
                height=0.16,
                text=f"{index:02d}",
                style=_theme_text_style(theme, font_size_pt=8.2, color=theme.colors.background, bold=True),
            )
            heading, body_text = _heading_body_with_fallback(
                text,
                _surface_label(deck, "Set one checkpoint.", "设置一个校验点。"),
            )
            _add_textbox(
                slide,
                x=x + 0.24,
                y=y + 0.7,
                width=step_width - 0.48,
                height=0.3,
                text=_safe_text(heading, 30),
                style=_theme_text_style(theme, font_size_pt=14.6, color=theme.colors.text, bold=True),
            )
            _add_textbox(
                slide,
                x=x + 0.24,
                y=y + 1.08,
                width=step_width - 0.48,
                height=0.34,
                text=_safe_text(body_text, 58),
                style=_theme_text_style(theme, font_size_pt=10.8, color=theme.colors.muted_text),
            )


def _render_risk_matrix_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    risks = body[:4] if len(body) >= 3 else _slot_texts(body, 3, fallback=" ")
    if len([text for text in risks if text.strip()]) < 2:
        _render_column_template(slide, deck_slide, deck, theme, column_count=2, variant=0)
        return
    margin = 0.7
    table_x = margin
    table_y = 1.5
    table_width = deck.canvas_width_in - margin * 2
    header_height = 0.48
    row_height = 0.92 if len(risks) <= 3 else 0.8
    columns = [0.38, 0.2, 0.42] if variant == 1 else [0.32, 0.22, 0.46]
    widths = [table_width * ratio for ratio in columns]
    labels = [
        _surface_label(deck, "Issue", "风险项"),
        _surface_label(deck, "Effect", "后果"),
        _surface_label(deck, "Action", "动作"),
    ]

    _add_textbox(
        slide,
        x=margin,
        y=0.5,
        width=deck.canvas_width_in - margin * 2,
        height=0.55,
        text=title,
        style=_theme_text_style(theme, font_size_pt=26, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    _add_rect(slide, x=margin, y=1.17, width=2.2, height=0.06, fill_color=theme.colors.primary)
    _add_rect(
        slide,
        x=table_x,
        y=table_y,
        width=table_width,
        height=header_height,
        fill_color=theme.colors.secondary if variant == 1 else theme.colors.primary,
        stroke_color=theme.colors.secondary if variant == 1 else theme.colors.primary,
    )

    x = table_x
    for label, width in zip(labels, widths):
        _add_textbox(
            slide,
            x=x + 0.18,
            y=table_y + 0.13,
            width=width - 0.36,
            height=0.18,
            text=label,
            style=_theme_text_style(theme, font_size_pt=9.5, color=theme.colors.background, bold=True),
        )
        x += width

    for row_index, risk_text in enumerate(risks):
        y = table_y + header_height + row_index * row_height
        fill = theme.colors.background if row_index % 2 == 0 else theme.colors.surface
        _add_rect(
            slide,
            x=table_x,
            y=y,
            width=table_width,
            height=row_height,
            fill_color=fill,
            stroke_color=theme.colors.surface,
            stroke_width_pt=0.8,
        )
        if variant == 1:
            _add_rect(
                slide,
                x=table_x,
                y=y,
                width=0.08,
                height=row_height,
                fill_color=[theme.colors.accent, theme.colors.primary, theme.colors.secondary, theme.colors.accent][row_index],
            )
        parts = _risk_cells_from_text(risk_text, deck)
        x = table_x
        for col_index, (text, width) in enumerate(zip(parts, widths)):
            color = theme.colors.text if col_index == 0 else theme.colors.muted_text
            max_chars = 42 if col_index == 0 else 58
            _add_textbox(
                slide,
                x=x + 0.18,
                y=y + 0.18,
                width=width - 0.36,
                height=0.5,
                text=_safe_text(text, max_chars),
                style=_theme_text_style(theme, font_size_pt=11.2, color=color, bold=col_index == 0),
            )
            x += width


def _render_key_takeaway_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    takeaways = body[:4] if len(body) >= 2 else _slot_texts(body, 2, fallback=" ")
    fallback_explanation = _action_fallback(deck, title)
    main_heading, main_body = _heading_body_with_fallback(takeaways[0], fallback_explanation)
    action_items = [
        _heading_body_with_fallback(takeaway, fallback_explanation)
        for takeaway in takeaways[1:4]
    ]
    margin = 0.82

    if variant == 1:
        _add_rect(slide, x=margin, y=0.72, width=0.1, height=5.55, fill_color=theme.colors.primary)
        _add_textbox(
            slide,
            x=margin + 0.28,
            y=0.6,
            width=4.8,
            height=0.34,
            text=_surface_label(deck, "Key Takeaway", "核心结论"),
            style=_theme_text_style(theme, font_size_pt=10, color=theme.colors.primary, bold=True),
        )
        _add_textbox(
            slide,
            x=margin + 0.28,
            y=1.08,
            width=6.0,
            height=1.08,
            text=_safe_text(main_heading or title, 72),
            style=_theme_text_style(theme, font_size_pt=30, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
        )
        _add_textbox(
            slide,
            x=margin + 0.3,
            y=2.36,
            width=5.8,
            height=0.78,
            text=_safe_text(main_body, 104),
            style=_theme_text_style(theme, font_size_pt=16, color=theme.colors.muted_text),
        )
        _add_rect(slide, x=7.48, y=1.0, width=4.85, height=4.8, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
        _add_textbox(
            slide,
            x=7.84,
            y=1.3,
            width=3.9,
            height=0.24,
            text=_surface_label(deck, "Next actions", "下一步行动"),
            style=_theme_text_style(theme, font_size_pt=9.8, color=theme.colors.primary, bold=True),
        )
        for index, (action_title, action_body) in enumerate(action_items, start=1):
            y = 1.86 + (index - 1) * 1.05
            _add_rect(slide, x=7.84, y=y + 0.06, width=0.34, height=0.28, fill_color=theme.colors.background, stroke_color=theme.colors.background)
            _add_textbox(
                slide,
                x=7.93,
                y=y + 0.11,
                width=0.16,
                height=0.14,
                text=str(index),
                style=_theme_text_style(theme, font_size_pt=7.5, color=theme.colors.primary, bold=True),
            )
            _add_textbox(
                slide,
                x=8.38,
                y=y,
                width=3.44,
                height=0.2,
                text=_safe_text(action_title, 34),
                style=_theme_text_style(theme, font_size_pt=13.5, color=theme.colors.text, bold=True),
            )
            _add_textbox(
                slide,
                x=8.38,
                y=y + 0.28,
                width=3.44,
                height=0.28,
                text=_safe_text(action_body, 58),
                style=_theme_text_style(theme, font_size_pt=10.8, color=theme.colors.muted_text),
            )
        return

    _add_rect(slide, x=0, y=0, width=0.18, height=deck.canvas_height_in, fill_color=theme.colors.primary)
    _add_textbox(
        slide,
        x=margin,
        y=0.58,
        width=deck.canvas_width_in - margin * 2,
        height=0.44,
        text=_surface_label(deck, "Key Takeaway", "核心结论"),
        style=_theme_text_style(theme, font_size_pt=10, color=theme.colors.primary, bold=True),
    )
    _add_textbox(
        slide,
        x=margin,
        y=0.88,
        width=deck.canvas_width_in - margin * 2,
        height=0.26,
        text=title,
        style=_theme_text_style(theme, font_size_pt=11.5, color=theme.colors.muted_text, bold=True),
    )
    _add_textbox(
        slide,
        x=margin,
        y=1.28,
        width=deck.canvas_width_in - margin * 2,
        height=0.9,
        text=_safe_text(main_heading or title, 72),
        style=_theme_text_style(theme, font_size_pt=32, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    if main_body:
        _add_textbox(
            slide,
            x=margin,
            y=2.3,
            width=deck.canvas_width_in - margin * 2.3,
            height=0.56,
            text=_safe_text(main_body, 104),
            style=_theme_text_style(theme, font_size_pt=17, color=theme.colors.muted_text),
        )
    _add_rect(slide, x=margin, y=3.15, width=deck.canvas_width_in - margin * 2, height=0.04, fill_color=theme.colors.surface)

    for index, (action_title, action_body) in enumerate(action_items, start=1):
        y = 3.55 + (index - 1) * 0.64
        _add_rect(
            slide,
            x=margin,
            y=y,
            width=0.44,
            height=0.32,
            fill_color=theme.colors.surface,
            stroke_color=theme.colors.surface,
        )
        _add_textbox(
            slide,
            x=margin + 0.1,
            y=y + 0.07,
            width=0.24,
            height=0.14,
            text=f"{index}",
            style=_theme_text_style(theme, font_size_pt=8.2, color=theme.colors.primary, bold=True),
        )
        _add_textbox(
            slide,
            x=margin + 0.68,
            y=y - 0.01,
            width=deck.canvas_width_in - margin * 2 - 0.72,
            height=0.2,
            text=_safe_text(action_title, 44),
            style=_theme_text_style(theme, font_size_pt=14.5, color=theme.colors.text, bold=True),
        )
        _add_textbox(
            slide,
            x=margin + 0.68,
            y=y + 0.24,
            width=deck.canvas_width_in - margin * 2 - 0.72,
            height=0.22,
            text=_safe_text(action_body, 76),
            style=_theme_text_style(theme, font_size_pt=11.8, color=theme.colors.muted_text),
        )


def _render_closing_slide_template(slide, deck_slide, deck: Deck, theme: Theme, variant: int = 0) -> None:
    title, body = _title_and_body_texts(deck_slide)
    actions = _action_pairs_from_body(body, deck, limit=3)

    if variant == 1:
        _add_rect(slide, x=0.72, y=0.78, width=3.2, height=5.92, fill_color=theme.colors.surface, stroke_color=theme.colors.surface)
        _add_rect(slide, x=1.1, y=1.16, width=0.9, height=0.08, fill_color=theme.colors.accent)
        _add_textbox(
            slide,
            x=1.1,
            y=1.66,
            width=2.28,
            height=1.18,
            text=_safe_text(title, 54),
            style=_theme_text_style(theme, font_size_pt=25, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
        )
        _add_textbox(
            slide,
            x=1.12,
            y=5.74,
            width=2.2,
            height=0.22,
            text=_surface_label(deck, "Action checklist", "行动清单"),
            style=_theme_text_style(theme, font_size_pt=9.8, color=theme.colors.primary, bold=True),
        )
        if actions:
            for index, (action_heading, action_body) in enumerate(actions, start=1):
                y = 1.18 + (index - 1) * 1.52
                _add_rect(slide, x=4.6, y=y, width=7.74, height=1.04, fill_color=theme.colors.background, stroke_color=theme.colors.surface, stroke_width_pt=1.0)
                _add_rect(slide, x=4.6, y=y, width=0.1, height=1.04, fill_color=[theme.colors.primary, theme.colors.secondary, theme.colors.accent][index - 1])
                _add_textbox(
                    slide,
                    x=4.92,
                    y=y + 0.22,
                    width=0.5,
                    height=0.18,
                    text=f"{index:02d}",
                    style=_theme_text_style(theme, font_size_pt=8.2, color=theme.colors.primary, bold=True),
                )
                _add_textbox(
                    slide,
                    x=5.58,
                    y=y + 0.16,
                    width=6.3,
                    height=0.22,
                    text=_safe_text(action_heading, 42),
                    style=_theme_text_style(theme, font_size_pt=14.5, color=theme.colors.text, bold=True),
                )
                _add_textbox(
                    slide,
                    x=5.58,
                    y=y + 0.52,
                    width=6.3,
                    height=0.24,
                    text=_safe_text(action_body, 76),
                    style=_theme_text_style(theme, font_size_pt=11.4, color=theme.colors.muted_text),
                )
        elif body:
            _add_textbox(
                slide,
                x=4.6,
                y=2.0,
                width=7.2,
                height=0.7,
                text=_short_phrase(body[0], max_chars=72),
                style=_theme_text_style(theme, font_size_pt=18, color=theme.colors.muted_text),
            )
        return

    _add_rect(slide, x=4.9, y=1.6, width=3.5, height=0.08, fill_color=theme.colors.accent)
    _add_textbox(
        slide,
        x=1.45,
        y=2.05 if actions else 2.35,
        width=deck.canvas_width_in - 2.9,
        height=0.9,
        text=title,
        style=_theme_text_style(theme, font_size_pt=34, color=theme.colors.text, bold=True, font_family=theme.fonts.heading),
    )
    if actions:
        for index, (action_heading, action_body) in enumerate(actions, start=1):
            y = 3.05 + (index - 1) * 0.76
            _add_rect(
                slide,
                x=2.55,
                y=y + 0.02,
                width=0.58,
                height=0.3,
                fill_color=theme.colors.surface,
                stroke_color=theme.colors.surface,
            )
            _add_textbox(
                slide,
                x=2.62,
                y=y + 0.06,
                width=0.44,
                height=0.18,
                text=f"{index:02d}",
                style=_theme_text_style(theme, font_size_pt=7.5, color=theme.colors.primary, bold=True),
            )
            _add_textbox(
                slide,
                x=3.24,
                y=y - 0.01,
                width=deck.canvas_width_in - 6.1,
                height=0.22,
                text=_safe_text(action_heading, 42),
                style=_theme_text_style(theme, font_size_pt=14.5, color=theme.colors.text, bold=True),
            )
            _add_textbox(
                slide,
                x=3.24,
                y=y + 0.26,
                width=deck.canvas_width_in - 6.1,
                height=0.26,
                text=_safe_text(action_body, 76),
                style=_theme_text_style(theme, font_size_pt=11.6, color=theme.colors.muted_text),
            )
    elif body:
        _add_textbox(
            slide,
            x=2.2,
            y=3.35,
            width=deck.canvas_width_in - 4.4,
            height=0.7,
            text=_short_phrase(body[0], max_chars=72),
            style=_theme_text_style(theme, font_size_pt=18, color=theme.colors.muted_text),
        )


def _render_template_slide(slide, deck_slide, deck: Deck, theme: Theme, slide_index: int) -> None:
    variant = _visual_variant_for_slide(deck, deck_slide, slide_index)
    if deck_slide.layout == "title_slide":
        _render_title_slide_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "section_divider":
        _render_section_divider_template(slide, deck_slide, deck, theme)
    elif deck_slide.layout == "two_column":
        _render_column_template(slide, deck_slide, deck, theme, column_count=2, variant=variant)
    elif deck_slide.layout == "three_column":
        _render_column_template(slide, deck_slide, deck, theme, column_count=3, variant=variant)
    elif deck_slide.layout == "four_cards":
        _render_four_cards_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "metric_cards":
        _render_metric_cards_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "closing_slide":
        _render_closing_slide_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "comparison_matrix":
        _render_comparison_matrix_template(slide, deck_slide, deck, theme)
    elif deck_slide.layout == "process_flow":
        _render_process_flow_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "risk_matrix":
        _render_risk_matrix_template(slide, deck_slide, deck, theme, variant=variant)
    elif deck_slide.layout == "key_takeaway":
        _render_key_takeaway_template(slide, deck_slide, deck, theme, variant=variant)


def render_deck_to_pptx(
    deck: Deck,
    theme: Theme,
    output_path: str | Path,
    assets_dir: str | Path | None = None,
) -> Path:
    """Render a validated deck into an editable PowerPoint file."""

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(deck.canvas_width_in)
    presentation.slide_height = Inches(deck.canvas_height_in)
    blank_layout = presentation.slide_layouts[6]

    for slide_index, deck_slide in enumerate(deck.slides, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        background_fill = slide.background.fill
        background_fill.solid()
        background_fill.fore_color.rgb = _rgb_color(theme.colors.background)

        if is_template_layout(deck_slide.layout):
            _render_template_slide(slide, deck_slide, deck, theme, slide_index)
        else:
            for element in deck_slide.elements:
                if isinstance(element, TextElement):
                    _render_text(slide, element, theme)
                elif isinstance(element, ShapeElement):
                    _render_shape(slide, element, theme)
                elif isinstance(element, ImageElement):
                    _render_image(slide, element, theme, assets_dir)

    presentation.save(output)
    return output
