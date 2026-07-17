"""Deterministic renderer: DeckDesign -> native editable PPTX.

Every element becomes a real PowerPoint object (text box, autoshape,
connector, native chart, table). No screenshots, no grouped vector soup.
Deck-wide chrome (kicker, page number, footer, theme motif) is stamped here
rather than asked of the model, which anchors cross-page consistency.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from ppt_agent.v2.design import TYPE_SCALE, ColorRole, ThemeSpec, best_text_color
from ppt_agent.v2.icons import resolve_icon
from ppt_agent.v2.ir import (
    CANVAS_HEIGHT,
    CANVAS_WIDTH,
    ChartItem,
    DeckDesign,
    Frame,
    IconItem,
    ImageItem,
    LineItem,
    PageDesign,
    ShapeItem,
    TableItem,
    TextItem,
)
from ppt_agent.v2.metrics import fit_font_size


SLIDE_WIDTH_EMU = 12_192_000  # 13.333 in, 16:9
SLIDE_HEIGHT_EMU = 6_858_000
EMU_PER_UNIT = SLIDE_WIDTH_EMU / CANVAS_WIDTH

_SHAPE_MAP: dict[str, int] = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "diamond": MSO_SHAPE.DIAMOND,
    "hexagon": MSO_SHAPE.HEXAGON,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "half_moon": MSO_SHAPE.MOON,
}

_CHART_MAP: dict[str, int] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}

_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_VALIGN_MAP = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _emu(units: float) -> Emu:
    return Emu(int(round(units * EMU_PER_UNIT)))


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _resolve(theme: ThemeSpec, role: str | None, fallback: str = "text") -> str:
    return theme.palette.resolve(role or fallback)


def _set_fill_alpha(shape, alpha: float) -> None:
    """Inject <a:alpha> into the solid fill; python-pptx has no alpha API."""

    if alpha >= 0.999:
        return
    sp_pr = shape.fill._xPr  # noqa: SLF001 - no public transparency API
    srgb = sp_pr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
    if srgb is None:
        solid = sp_pr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
    if srgb is None:
        return
    alpha_el = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    srgb.append(alpha_el)


def _set_run_fonts(run, theme: ThemeSpec, *, heading: bool) -> None:
    latin = theme.fonts.heading_latin if heading else theme.fonts.body_latin
    east_asian = theme.fonts.heading_east_asian if heading else theme.fonts.body_east_asian
    run.font.name = latin
    r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 - east-asian font needs raw XML
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = r_pr.makeelement(qn("a:ea"), {})
        r_pr.append(ea)
    ea.set("typeface", east_asian)


def _fill_text_frame(
    text_frame,
    item: TextItem,
    theme: ThemeSpec,
    *,
    size_pt: float,
    color_hex: str,
) -> None:
    spec = TYPE_SCALE[item.role]
    heading = item.role in ("display", "title", "section", "h3", "kicker", "stat")
    bold = item.bold if item.bold is not None else spec.bold
    text_frame.word_wrap = True
    text_frame.vertical_anchor = _VALIGN_MAP[item.valign]
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)

    lines = item.text.split("\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _ALIGN_MAP[item.align]
        paragraph.line_spacing = spec.line_spacing
        if item.bullet == "dot" and line.strip():
            line = f"•  {line}"
        elif item.bullet == "number" and line.strip():
            line = f"{index + 1}.  {line}"
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = item.italic
        run.font.color.rgb = _rgb(color_hex)
        _set_run_fonts(run, theme, heading=heading)


def _render_text(slide, item: TextItem, theme: ThemeSpec) -> None:
    frame = item.frame
    box = slide.shapes.add_textbox(_emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h))
    spec = TYPE_SCALE[item.role]
    size_pt = fit_font_size(
        item.text,
        role=item.role,
        frame_width_units=frame.w,
        frame_height_units=frame.h,
        requested_size_pt=item.size_pt,
    )
    color_hex = _resolve(theme, item.color or spec.default_color)
    _fill_text_frame(box.text_frame, item, theme, size_pt=size_pt, color_hex=color_hex)


def _render_shape(slide, item: ShapeItem, theme: ThemeSpec) -> None:
    frame = item.frame
    shape = slide.shapes.add_shape(
        _SHAPE_MAP[item.shape], _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
    )
    if item.rotation_deg:
        shape.rotation = item.rotation_deg
    if item.shape == "pill":
        shape.adjustments[0] = 0.5
    elif item.shape == "rounded_rectangle" and theme.corner_radius > 0:
        smaller = min(frame.w, frame.h)
        shape.adjustments[0] = max(0.02, min(0.5, theme.corner_radius / smaller))
    shape.shadow.inherit = False

    if item.gradient is not None:
        shape.fill.gradient()
        stops = shape.fill.gradient_stops
        stops[0].color.rgb = _rgb(_resolve(theme, item.gradient.start))
        stops[1].color.rgb = _rgb(_resolve(theme, item.gradient.end))
        try:
            shape.fill.gradient_angle = item.gradient.angle_deg
        except ValueError:
            pass
    elif item.fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(_resolve(theme, item.fill))
        _set_fill_alpha(shape, item.fill_alpha)
    else:
        shape.fill.background()

    if item.stroke is not None:
        shape.line.color.rgb = _rgb(_resolve(theme, item.stroke))
        shape.line.width = Pt(item.stroke_width)
    else:
        shape.line.fill.background()


def _render_line(slide, item: LineItem, theme: ThemeSpec) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _emu(item.x1), _emu(item.y1), _emu(item.x2), _emu(item.y2)
    )
    connector.line.color.rgb = _rgb(_resolve(theme, item.color))
    connector.line.width = Pt(item.width)
    if item.dash:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    connector.shadow.inherit = False


def _render_icon(slide, item: IconItem, theme: ThemeSpec) -> None:
    frame = item.frame
    glyph, monochrome = resolve_icon(item.name)
    if item.background is not None and item.background_shape != "none":
        shape_kind = (
            MSO_SHAPE.OVAL if item.background_shape == "circle" else MSO_SHAPE.ROUNDED_RECTANGLE
        )
        badge = slide.shapes.add_shape(
            shape_kind, _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(_resolve(theme, item.background))
        badge.line.fill.background()
        badge.shadow.inherit = False

    box = slide.shapes.add_textbox(_emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h))
    text_frame = box.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = glyph
    # Glyph cell is square-ish; size the glyph to ~55% of the frame height.
    run.font.size = Pt(max(8.0, frame.h * 0.55 * 72 / 96))
    if monochrome:
        run.font.color.rgb = _rgb(_resolve(theme, item.color))


def _render_chart(slide, item: ChartItem, theme: ThemeSpec) -> None:
    frame = item.frame
    data = CategoryChartData()
    data.categories = item.categories
    for series in item.series:
        data.add_series(series.name, tuple(series.values))
    graphic_frame = slide.shapes.add_chart(
        _CHART_MAP[item.chart],
        _emu(frame.x),
        _emu(frame.y),
        _emu(frame.w),
        _emu(frame.h),
        data,
    )
    chart = graphic_frame.chart
    chart.has_title = item.title is not None
    if item.title:
        chart.chart_title.text_frame.text = item.title
        for run_paragraph in chart.chart_title.text_frame.paragraphs:
            for run in run_paragraph.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = _rgb(theme.palette.text)
    chart.has_legend = item.show_legend and len(item.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)

    series_colors = [
        theme.palette.primary,
        theme.palette.secondary,
        theme.palette.accent,
        theme.palette.muted,
    ]
    plot = chart.plots[0]
    if item.chart in ("pie", "doughnut"):
        # Color slices individually so single-series pies still vary.
        for point_index, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(
                series_colors[point_index % len(series_colors)]
            )
    else:
        for series_index, series in enumerate(plot.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(
                series_colors[series_index % len(series_colors)]
            )
    plot.has_data_labels = item.show_data_labels
    if item.show_data_labels:
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = _rgb(theme.palette.muted)
    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
        except ValueError:
            continue
        axis.tick_labels.font.size = Pt(10)
        axis.tick_labels.font.color.rgb = _rgb(theme.palette.muted)
        axis.format.line.color.rgb = _rgb(theme.palette.surface_alt)


def _render_table(slide, item: TableItem, theme: ThemeSpec) -> None:
    frame = item.frame
    row_count = len(item.rows) + 1
    graphic_frame = slide.shapes.add_table(
        row_count, len(item.headers), _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
    )
    table = graphic_frame.table
    table.first_row = False
    table.horz_banding = False

    def style_cell(cell, text: str, *, header: bool, banded: bool) -> None:
        cell.fill.solid()
        if header:
            cell.fill.fore_color.rgb = _rgb(theme.palette.primary)
        else:
            cell.fill.fore_color.rgb = _rgb(
                theme.palette.surface_alt if banded else theme.palette.surface
            )
        cell.margin_left = _emu(8)
        cell.margin_right = _emu(8)
        cell.margin_top = _emu(4)
        cell.margin_bottom = _emu(4)
        text_frame = cell.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(11 if header else 10.5)
        run.font.bold = header
        run.font.color.rgb = _rgb(
            theme.palette.on_primary if header else theme.palette.text
        )
        _set_run_fonts(run, theme, heading=header)

    for column, header in enumerate(item.headers):
        style_cell(table.cell(0, column), header, header=True, banded=False)
    for row_index, row in enumerate(item.rows):
        for column, value in enumerate(row):
            style_cell(
                table.cell(row_index + 1, column),
                value,
                header=False,
                banded=row_index % 2 == 1,
            )


def _render_image(slide, item: ImageItem, theme: ThemeSpec, assets_dir: Path | None) -> None:
    frame = item.frame
    path: Path | None = None
    if item.src:
        candidate = Path(item.src)
        if not candidate.is_absolute() and assets_dir is not None:
            candidate = assets_dir / candidate
        if candidate.is_file():
            path = candidate
    if path is not None:
        picture = slide.shapes.add_picture(str(path), _emu(frame.x), _emu(frame.y))
        # Contain-fit into the frame, keeping aspect ratio.
        scale = min(
            _emu(frame.w) / picture.width,
            _emu(frame.h) / picture.height,
        )
        picture.width = int(picture.width * scale)
        picture.height = int(picture.height * scale)
        picture.left = _emu(frame.x) + (_emu(frame.w) - picture.width) // 2
        picture.top = _emu(frame.y) + (_emu(frame.h) - picture.height) // 2
        return
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = _rgb(theme.palette.surface_alt)
    placeholder.line.color.rgb = _rgb(theme.palette.muted)
    placeholder.line.width = Pt(0.75)
    placeholder.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    placeholder.shadow.inherit = False
    text_frame = placeholder.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"🖼\n{item.label}"
    run.font.size = Pt(11)
    run.font.color.rgb = _rgb(theme.palette.muted)
    _set_run_fonts(run, theme, heading=False)


def _paint_background(slide, page: PageDesign, theme: ThemeSpec) -> None:
    if page.background_gradient is not None:
        backdrop = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Emu(SLIDE_WIDTH_EMU), Emu(SLIDE_HEIGHT_EMU)
        )
        backdrop.fill.gradient()
        stops = backdrop.fill.gradient_stops
        stops[0].color.rgb = _rgb(_resolve(theme, page.background_gradient.start))
        stops[1].color.rgb = _rgb(_resolve(theme, page.background_gradient.end))
        try:
            backdrop.fill.gradient_angle = page.background_gradient.angle_deg
        except ValueError:
            pass
        backdrop.line.fill.background()
        backdrop.shadow.inherit = False
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(_resolve(theme, page.background))


def _stamp_motif(slide, page: PageDesign, theme: ThemeSpec) -> None:
    if theme.motif == "none" or page.role in ("cover", "section_divider", "closing"):
        return
    if theme.motif == "corner_arc":
        arc = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, _emu(CANVAS_WIDTH - 120), _emu(-120), _emu(240), _emu(240)
        )
        arc.fill.solid()
        arc.fill.fore_color.rgb = _rgb(theme.palette.primary_soft)
        _set_fill_alpha(arc, 0.55)
        arc.line.fill.background()
        arc.shadow.inherit = False
    elif theme.motif == "side_band":
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _emu(10), _emu(CANVAS_HEIGHT))
        band.fill.solid()
        band.fill.fore_color.rgb = _rgb(theme.palette.primary)
        band.line.fill.background()
        band.shadow.inherit = False
    elif theme.motif == "top_rule":
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _emu(CANVAS_WIDTH), _emu(6))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb(theme.palette.primary)
        rule.line.fill.background()
        rule.shadow.inherit = False
    elif theme.motif == "dot_grid":
        for row in range(3):
            for column in range(6):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    _emu(CANVAS_WIDTH - 150 + column * 22),
                    _emu(CANVAS_HEIGHT - 90 + row * 22),
                    _emu(5),
                    _emu(5),
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(theme.palette.primary_soft)
                dot.line.fill.background()
                dot.shadow.inherit = False
    elif theme.motif == "diagonal":
        slash = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM, _emu(CANVAS_WIDTH - 170), _emu(-40), _emu(220), _emu(140)
        )
        slash.rotation = 12
        slash.fill.solid()
        slash.fill.fore_color.rgb = _rgb(theme.palette.primary_soft)
        _set_fill_alpha(slash, 0.5)
        slash.line.fill.background()
        slash.shadow.inherit = False


def _stamp_chrome(slide, page: PageDesign, deck: DeckDesign, theme: ThemeSpec) -> None:
    if not page.show_chrome or page.role in ("cover", "closing"):
        return
    background_hex = _resolve(theme, page.background)
    muted_role: ColorRole = (
        "muted" if best_text_color(theme.palette, background_hex) == "text" else "on_primary"
    )
    if (
        theme.chrome.show_section_kicker
        and page.section
        and page.role not in ("toc", "section_divider")
    ):
        kicker = TextItem(
            id="_chrome_kicker",
            frame=Frame(x=64, y=26, w=600, h=22),
            text=page.section.upper(),
            role="kicker",
        )
        _render_text(slide, kicker, theme)
    if theme.chrome.show_page_number:
        page_number = TextItem(
            id="_chrome_page_number",
            frame=Frame(x=CANVAS_WIDTH - 110, y=CANVAS_HEIGHT - 36, w=80, h=22),
            text=f"{page.page_number:02d}",
            role="caption",
            color=muted_role,
            align="right",
        )
        _render_text(slide, page_number, theme)
    if theme.chrome.show_footer:
        footer = TextItem(
            id="_chrome_footer",
            frame=Frame(x=64, y=CANVAS_HEIGHT - 36, w=600, h=22),
            text=deck.deck_title,
            role="caption",
            color=muted_role,
        )
        _render_text(slide, footer, theme)


def render_page(slide, page: PageDesign, deck: DeckDesign, *, assets_dir: Path | None = None) -> None:
    theme = deck.theme
    _paint_background(slide, page, theme)
    _stamp_motif(slide, page, theme)
    for element in page.elements:
        if isinstance(element, TextItem):
            _render_text(slide, element, theme)
        elif isinstance(element, ShapeItem):
            _render_shape(slide, element, theme)
        elif isinstance(element, LineItem):
            _render_line(slide, element, theme)
        elif isinstance(element, IconItem):
            _render_icon(slide, element, theme)
        elif isinstance(element, ChartItem):
            _render_chart(slide, element, theme)
        elif isinstance(element, TableItem):
            _render_table(slide, element, theme)
        elif isinstance(element, ImageItem):
            _render_image(slide, element, theme, assets_dir)
    _stamp_chrome(slide, page, deck, theme)
    if page.speaker_notes:
        slide.notes_slide.notes_text_frame.text = page.speaker_notes


def render_deck(deck: DeckDesign, output_path: str | Path, *, assets_dir: str | Path | None = None) -> Path:
    """Render the whole deck and return the written PPTX path."""

    presentation = Presentation()
    presentation.slide_width = Emu(SLIDE_WIDTH_EMU)
    presentation.slide_height = Emu(SLIDE_HEIGHT_EMU)
    blank_layout = presentation.slide_layouts[6]
    assets = Path(assets_dir) if assets_dir else None
    for page in deck.pages:
        slide = presentation.slides.add_slide(blank_layout)
        render_page(slide, page, deck, assets_dir=assets)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    return output
