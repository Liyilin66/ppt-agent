import json
import sys
import types
from pathlib import Path

from pptx import Presentation

import ppt_agent.cli as cli
import ppt_agent.pipeline as pipeline
from ppt_agent.cli import main
from ppt_agent.generation import GenerationAttempt, GenerationResult
from ppt_agent.load import load_deck
from ppt_agent.qa import analyze_deck


EXAMPLES_DIR = Path(__file__).resolve().parents[1] / "examples"


def test_render_cli_generates_pptx(tmp_path: Path) -> None:
    output_path = tmp_path / "sample_deck.pptx"

    status = main(
        [
            "render",
            str(EXAMPLES_DIR / "sample_slide_ir.json"),
            "--theme",
            str(EXAMPLES_DIR / "theme.json"),
            "--output",
            str(output_path),
        ]
    )

    assert status == 0
    assert output_path.exists()
    assert len(Presentation(output_path).slides) == 3


def test_qa_cli_writes_report_json(tmp_path: Path) -> None:
    output_path = tmp_path / "qa_report.json"

    status = main(
        [
            "qa",
            str(EXAMPLES_DIR / "sample_slide_ir.json"),
            "--theme",
            str(EXAMPLES_DIR / "theme.json"),
            "--output",
            str(output_path),
        ]
    )

    assert status == 0
    report = json.loads(output_path.read_text(encoding="utf-8"))
    assert report["deck_id"] == "sample_clean_business_deck"
    assert "score" in report
    assert "issues" in report


def test_generate_cli_without_api_key_exits_with_clear_message(
    tmp_path: Path,
    monkeypatch,
    capsys,
) -> None:
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    output_path = tmp_path / "generated_deck.json"

    status = main(
        [
            "generate",
            "--topic",
            "AI in Education",
            "--audience",
            "university students",
            "--slides",
            "8",
            "--theme",
            str(EXAMPLES_DIR / "theme.json"),
            "--output",
            str(output_path),
            "--min-qa-score",
            "85",
            "--max-attempts",
            "3",
            "--qa-output",
            str(tmp_path / "qa.json"),
            "--attempts-output",
            str(tmp_path / "attempts.json"),
        ]
    )

    captured = capsys.readouterr()
    assert status == 1
    assert "OPENAI_API_KEY is not set" in captured.out
    assert not output_path.exists()


def test_generate_cli_help_includes_quality_gate_options(capsys) -> None:
    try:
        main(["generate", "--help"])
    except SystemExit as exc:
        assert exc.code == 0

    captured = capsys.readouterr()
    assert "--min-qa-score" in captured.out
    assert "--max-attempts" in captured.out
    assert "--qa-output" in captured.out
    assert "--attempts-output" in captured.out
    assert "--requirements" in captured.out
    assert "--prompt" in captured.out


def test_generate_cli_rejects_invalid_min_qa_score() -> None:
    try:
        main(
            [
                "generate",
                "--topic",
                "AI in Education",
                "--audience",
                "university students",
                "--slides",
                "8",
                "--theme",
                str(EXAMPLES_DIR / "theme.json"),
                "--output",
                "/tmp/generated_deck.json",
                "--min-qa-score",
                "101",
            ]
        )
    except SystemExit as exc:
        assert exc.code == 2


def test_generate_cli_returns_failure_when_quality_gate_rejects(
    tmp_path: Path,
    monkeypatch,
    capsys,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    fake_module = types.ModuleType("langchain_openai")
    fake_module.ChatOpenAI = lambda model: object()
    monkeypatch.setitem(sys.modules, "langchain_openai", fake_module)

    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    qa_report = analyze_deck(deck)

    def fake_generate_deck_with_quality_gate(*args, **kwargs):
        return GenerationResult(
            deck=deck,
            qa_report=qa_report,
            attempts=[
                GenerationAttempt(
                    attempt_index=1,
                    deck=deck,
                    qa_report=qa_report,
                    accepted=False,
                )
            ],
            accepted=False,
        )

    monkeypatch.setattr(cli, "generate_deck_with_quality_gate", fake_generate_deck_with_quality_gate)
    output_path = tmp_path / "generated_deck.json"

    status = main(
        [
            "generate",
            "--topic",
            "AI in Education",
            "--audience",
            "university students",
            "--slides",
            "3",
            "--theme",
            str(EXAMPLES_DIR / "theme.json"),
            "--output",
            str(output_path),
        ]
    )

    captured = capsys.readouterr()
    assert status == 2
    assert output_path.exists()
    assert "did not meet the QA score gate" in captured.out


def test_generate_cli_passes_user_requirements(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    qa_report = analyze_deck(deck)
    captured_request = {}

    def fake_generate_deck_with_quality_gate(model, request, **kwargs):
        captured_request["request"] = request
        return GenerationResult(
            deck=deck,
            qa_report=qa_report,
            attempts=[
                GenerationAttempt(
                    attempt_index=1,
                    deck=deck,
                    qa_report=qa_report,
                    accepted=True,
                )
            ],
            accepted=True,
        )

    monkeypatch.setattr(cli, "generate_deck_with_quality_gate", fake_generate_deck_with_quality_gate)
    output_path = tmp_path / "generated_deck.json"

    status = main(
        [
            "generate",
            "--topic",
            "AI 教育",
            "--audience",
            "大学生",
            "--slides",
            "3",
            "--theme",
            str(EXAMPLES_DIR / "theme.json"),
            "--output",
            str(output_path),
            "--requirements",
            "做一份中文课堂展示，提醒学术诚信风险。",
        ]
    )

    assert status == 0
    assert captured_request["request"].language == "zh-CN"
    assert captured_request["request"].user_requirements == "做一份中文课堂展示，提醒学术诚信风险。"


def _generation_result(accepted: bool) -> GenerationResult:
    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    qa_report = analyze_deck(deck)
    return GenerationResult(
        deck=deck,
        qa_report=qa_report,
        attempts=[
            GenerationAttempt(
                attempt_index=1,
                deck=deck,
                qa_report=qa_report,
                accepted=accepted,
            )
        ],
        accepted=accepted,
    )


def _install_fake_openai(monkeypatch) -> None:
    fake_module = types.ModuleType("langchain_openai")
    fake_module.ChatOpenAI = lambda model: object()
    monkeypatch.setitem(sys.modules, "langchain_openai", fake_module)


def _build_args(output_dir: Path) -> list[str]:
    return [
        "build",
        "--topic",
        "AI in Education",
        "--audience",
        "university students",
        "--slides",
        "3",
        "--theme",
        str(EXAMPLES_DIR / "theme.json"),
        "--output-dir",
        str(output_dir),
        "--min-qa-score",
        "80",
        "--max-attempts",
        "2",
    ]


def test_build_cli_help_includes_build_options(capsys) -> None:
    try:
        main(["build", "--help"])
    except SystemExit as exc:
        assert exc.code == 0

    captured = capsys.readouterr()
    assert "--topic" in captured.out
    assert "--audience" in captured.out
    assert "--slides" in captured.out
    assert "--output-dir" in captured.out
    assert "--assets-dir" in captured.out
    assert "--patch" in captured.out
    assert "--requirements" in captured.out
    assert "--prompt" in captured.out


def test_build_cli_passes_user_requirements_to_pipeline(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    captured_request = {}

    def fake_run_build_pipeline(model, request):
        captured_request["request"] = request
        return pipeline.BuildPipelineResult(
            generation_result=_generation_result(True),
            artifacts=[],
            accepted=True,
            status_code=0,
            messages=[],
        )

    monkeypatch.setattr(cli, "run_build_pipeline", fake_run_build_pipeline)

    status = main(_build_args(tmp_path) + ["--prompt", "中文课堂展示，突出 AI 学习应用。"])

    assert status == 0
    generation_request = captured_request["request"].generation_request
    assert generation_request.language == "zh-CN"
    assert generation_request.user_requirements == "中文课堂展示，突出 AI 学习应用。"


def test_build_cli_without_api_key_exits_with_clear_message(
    tmp_path: Path,
    monkeypatch,
    capsys,
) -> None:
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)

    status = main(_build_args(tmp_path))

    captured = capsys.readouterr()
    assert status == 1
    assert "OPENAI_API_KEY is not set" in captured.out
    assert not (tmp_path / "generated_deck_ir.json").exists()


def test_build_cli_accepted_outputs_all_files(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    monkeypatch.setattr(pipeline, "generate_deck_with_quality_gate", lambda *args, **kwargs: _generation_result(True))

    status = main(_build_args(tmp_path))

    assert status == 0
    assert (tmp_path / "generated_deck_ir.json").exists()
    assert (tmp_path / "generated_qa_report.json").exists()
    assert (tmp_path / "generated_attempts.json").exists()
    pptx_path = tmp_path / "generated_deck.pptx"
    assert pptx_path.exists()
    assert len(Presentation(pptx_path).slides) == 3


def test_build_cli_with_patch_outputs_original_and_patched_files(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    monkeypatch.setattr(pipeline, "generate_deck_with_quality_gate", lambda *args, **kwargs: _generation_result(True))

    status = main(_build_args(tmp_path) + ["--patch", str(EXAMPLES_DIR / "sample_patch.json")])

    assert status == 0
    assert (tmp_path / "generated_deck_ir.json").exists()
    assert (tmp_path / "generated_qa_report.json").exists()
    assert (tmp_path / "generated_attempts.json").exists()
    assert (tmp_path / "generated_deck.pptx").exists()
    assert (tmp_path / "patchable_elements.json").exists()
    assert (tmp_path / "patched_deck_ir.json").exists()
    assert (tmp_path / "patch_report.json").exists()
    assert (tmp_path / "patched_deck.pptx").exists()


def test_build_cli_with_patch_issue_outputs_files_and_returns_2(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    monkeypatch.setattr(pipeline, "generate_deck_with_quality_gate", lambda *args, **kwargs: _generation_result(True))
    patch_path = tmp_path / "bad_patch.json"
    patch_path.write_text(
        json.dumps(
            {
                "operations": [
                    {
                        "op": "update_text",
                        "slide_id": "missing_slide",
                        "element_id": "s1_title",
                        "text": "No-op",
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    status = main(_build_args(tmp_path) + ["--patch", str(patch_path)])

    assert status == 2
    assert (tmp_path / "generated_deck_ir.json").exists()
    assert (tmp_path / "generated_qa_report.json").exists()
    assert (tmp_path / "generated_attempts.json").exists()
    assert (tmp_path / "generated_deck.pptx").exists()
    assert (tmp_path / "patchable_elements.json").exists()
    assert (tmp_path / "patched_deck_ir.json").exists()
    result = json.loads((tmp_path / "patch_report.json").read_text(encoding="utf-8"))
    assert result["issues"][0]["code"] == "SLIDE_NOT_FOUND"
    assert (tmp_path / "patched_deck.pptx").exists()


def test_build_cli_rejected_outputs_all_files_and_returns_2(
    tmp_path: Path,
    monkeypatch,
    capsys,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    _install_fake_openai(monkeypatch)
    monkeypatch.setattr(pipeline, "generate_deck_with_quality_gate", lambda *args, **kwargs: _generation_result(False))

    status = main(_build_args(tmp_path))

    captured = capsys.readouterr()
    assert status == 2
    assert (tmp_path / "generated_deck_ir.json").exists()
    assert (tmp_path / "generated_qa_report.json").exists()
    assert (tmp_path / "generated_attempts.json").exists()
    assert (tmp_path / "generated_deck.pptx").exists()
    assert "did not meet the QA score gate" in captured.out


def test_patch_cli_help(capsys) -> None:
    try:
        main(["patch", "--help"])
    except SystemExit as exc:
        assert exc.code == 0

    captured = capsys.readouterr()
    assert "--patch" in captured.out
    assert "--output" in captured.out
    assert "--result-output" in captured.out


def test_patch_cli_success_outputs_patched_deck(tmp_path: Path) -> None:
    output_path = tmp_path / "patched_deck_ir.json"

    status = main(
        [
            "patch",
            str(EXAMPLES_DIR / "sample_slide_ir.json"),
            "--patch",
            str(EXAMPLES_DIR / "sample_patch.json"),
            "--output",
            str(output_path),
        ]
    )

    assert status == 0
    patched = json.loads(output_path.read_text(encoding="utf-8"))
    assert patched["slides"][0]["elements"][0]["text"] == "Updated Q3 Operating Review"


def test_patch_cli_with_issue_returns_2_and_outputs_result(tmp_path: Path) -> None:
    patch_path = tmp_path / "bad_patch.json"
    patch_path.write_text(
        json.dumps(
            {
                "operations": [
                    {
                        "op": "update_text",
                        "slide_id": "missing_slide",
                        "element_id": "s1_title",
                        "text": "No-op",
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    output_path = tmp_path / "patched_deck_ir.json"
    result_path = tmp_path / "patch_report.json"

    status = main(
        [
            "patch",
            str(EXAMPLES_DIR / "sample_slide_ir.json"),
            "--patch",
            str(patch_path),
            "--output",
            str(output_path),
            "--result-output",
            str(result_path),
        ]
    )

    assert status == 2
    assert output_path.exists()
    result = json.loads(result_path.read_text(encoding="utf-8"))
    assert result["issues"][0]["code"] == "SLIDE_NOT_FOUND"
    assert result["accepted"] is False
    assert result["success"] is False
