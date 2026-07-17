"""End-to-end orchestrator tests with the deterministic mock client."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import pytest
from pydantic import ValidationError
from pptx import Presentation

from ppt_agent.v2.ir import Frame, PageDesign, TextItem
from ppt_agent.v2.mock import MockLLMClient
from ppt_agent.v2.orchestrator import BuildRequest, build_deck, plan_deck
from ppt_agent.v2.planning import (
    EditablePage,
    editable_plan_from_skeleton,
    skeleton_from_editable_plan,
)
from ppt_agent.v2.providers import UsageMeter


def _request(tmp_path: Path, **overrides) -> BuildRequest:
    defaults = dict(
        prompt="AI Agent 产品方案",
        page_count=20,
        output_dir=str(tmp_path / "out"),
        deck_name="test",
    )
    defaults.update(overrides)
    return BuildRequest(**defaults)


class _FailingDesignClient(MockLLMClient):
    """Mock that always fails page_design to exercise the archetype fallback."""

    async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
        if task == "page_design":
            raise RuntimeError("provider exploded")
        return await super().complete_json(task=task, **kwargs)


class _UnresolvedQADesignClient(MockLLMClient):
    """Returns a valid page payload with text collisions that repair cannot fix."""

    async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
        if task == "page_design":
            context = kwargs.get("context") or {}
            page_number = int(context.get("page_number", 1))
            return PageDesign(
                page_number=page_number,
                role="content",
                section=context.get("section_title"),
                title="QA collision",
                elements=[
                    TextItem(
                        id="first",
                        frame=Frame(x=100, y=180, w=420, h=120),
                        text="第一段文本",
                    ),
                    TextItem(
                        id="second",
                        frame=Frame(x=140, y=200, w=420, h=120),
                        text="第二段文本",
                    ),
                    TextItem(
                        id="third",
                        frame=Frame(x=720, y=440, w=280, h=80),
                        text="第三段文本",
                    ),
                ],
            ).model_dump(mode="json")
        return await super().complete_json(task=task, **kwargs)


class TestBuildDeck:
    def test_rejects_more_than_100_pages(self, tmp_path: Path) -> None:
        with pytest.raises(ValidationError):
            _request(tmp_path, page_count=101)

    def test_full_offline_build(self, tmp_path: Path) -> None:
        result = build_deck(_request(tmp_path), MockLLMClient(), progress=lambda _: None)
        assert result.status == "succeeded"
        assert result.page_count == 20
        presentation = Presentation(result.pptx_path)
        assert len(presentation.slides) == 20
        for artifact in (
            result.deck_design_path,
            result.qa_report_path,
            result.run_report_path,
        ):
            assert Path(artifact).is_file()
        report = json.loads(Path(result.run_report_path).read_text(encoding="utf-8"))
        assert len(report["outcomes"]) == 20
        qa_report = json.loads(Path(result.qa_report_path).read_text(encoding="utf-8"))
        assert qa_report["total_pages"] == 20
        assert [item["page_number"] for item in qa_report["results"]] == list(
            range(1, 21)
        )

    def test_hundred_page_build(self, tmp_path: Path) -> None:
        result = build_deck(
            _request(tmp_path, page_count=100), MockLLMClient(), progress=lambda _: None
        )
        assert result.page_count == 100
        assert len(Presentation(result.pptx_path).slides) == 100

    def test_failed_pages_fall_back_instead_of_holes(self, tmp_path: Path) -> None:
        result = build_deck(
            _request(tmp_path), _FailingDesignClient(), progress=lambda _: None
        )
        assert result.status == "succeeded_with_fallbacks"
        assert result.fallback_pages > 0
        assert result.page_count == 20  # no holes

    def test_strict_qa_gate_replaces_unresolved_pages(self, tmp_path: Path) -> None:
        result = build_deck(
            _request(tmp_path, repair_rounds=0, qa_gate="strict"),
            _UnresolvedQADesignClient(),
            progress=lambda _: None,
        )
        assert result.status == "succeeded_with_fallbacks"
        assert result.fallback_pages > 0
        assert result.pptx_path is not None
        assert Path(result.pptx_path).is_file()
        qa_report = json.loads(Path(result.qa_report_path).read_text(encoding="utf-8"))
        assert qa_report["pages_with_errors"] == 0

    def test_lenient_qa_gate_marks_but_keeps_unresolved_pages(self, tmp_path: Path) -> None:
        result = build_deck(
            _request(tmp_path, repair_rounds=0, qa_gate="lenient"),
            _UnresolvedQADesignClient(),
            progress=lambda _: None,
        )
        assert result.status == "completed_with_qa_errors"
        assert result.pptx_path is not None
        assert Path(result.pptx_path).is_file()
        qa_report = json.loads(Path(result.qa_report_path).read_text(encoding="utf-8"))
        assert qa_report["pages_with_errors"] > 0

    def test_strict_gate_does_not_render_when_errors_remain(
        self, tmp_path: Path, monkeypatch
    ) -> None:
        from ppt_agent.v2 import orchestrator

        original = orchestrator._build_anchor_pages

        def broken_anchor_pages(*args, **kwargs):
            anchors = original(*args, **kwargs)
            page_number = min(anchors)
            anchor = anchors[page_number]
            anchors[page_number] = anchor.model_copy(
                update={
                    "elements": [
                        TextItem(
                            id="anchor_a",
                            frame=Frame(x=100, y=180, w=500, h=140),
                            text="锚点页冲突 A",
                        ),
                        TextItem(
                            id="anchor_b",
                            frame=Frame(x=140, y=200, w=500, h=140),
                            text="锚点页冲突 B",
                        ),
                    ]
                }
            )
            return anchors

        monkeypatch.setattr(orchestrator, "_build_anchor_pages", broken_anchor_pages)
        result = build_deck(
            _request(tmp_path, qa_gate="strict"),
            MockLLMClient(),
            progress=lambda _: None,
        )
        assert result.status == "quality_gate_failed"
        assert result.pptx_path is None
        assert not (Path(result.deck_design_path).parent / "test.pptx").exists()
        run_report = json.loads(Path(result.run_report_path).read_text(encoding="utf-8"))
        assert run_report["quality_gate"] == {
            "mode": "strict",
            "passed": False,
            "pages_with_errors": 1,
            "pptx_generated": False,
        }

    def test_resume_reuses_page_checkpoints(self, tmp_path: Path) -> None:
        request = _request(tmp_path)
        build_deck(request, MockLLMClient(), progress=lambda _: None)

        # Poison one checkpointed page title, then resume: the poisoned value
        # must survive, proving the page was not regenerated.
        checkpoint_dir = Path(request.output_dir) / "checkpoints" / "pages"
        target = sorted(checkpoint_dir.glob("page_*.json"))[0]
        payload = json.loads(target.read_text(encoding="utf-8"))
        payload["page"]["title"] = "RESUMED_MARKER"
        target.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

        resumed = build_deck(
            request.model_copy(update={"resume": True}),
            MockLLMClient(),
            progress=lambda _: None,
        )
        design = json.loads(Path(resumed.deck_design_path).read_text(encoding="utf-8"))
        assert any(page["title"] == "RESUMED_MARKER" for page in design["pages"])

    def test_language_override_reaches_brief(self, tmp_path: Path) -> None:
        result = build_deck(
            _request(tmp_path, language="en"), MockLLMClient(), progress=lambda _: None
        )
        design = json.loads(Path(result.deck_design_path).read_text(encoding="utf-8"))
        assert design["language"] == "en"


class TestCLI:
    def test_v2_demo_command(self, tmp_path: Path, capsys) -> None:
        from ppt_agent.cli import main

        exit_code = main(
            [
                "v2",
                "demo",
                "--prompt",
                "测试主题",
                "--pages",
                "12",
                "--output-dir",
                str(tmp_path / "cli_out"),
            ]
        )
        assert exit_code == 0
        output = capsys.readouterr().out
        assert "status: succeeded" in output
        assert (tmp_path / "cli_out" / "deck.pptx").is_file()

    def test_v2_preview_command(self, tmp_path: Path, capsys) -> None:
        from ppt_agent.cli import main

        main(
            [
                "v2",
                "demo",
                "--prompt",
                "预览测试",
                "--pages",
                "8",
                "--output-dir",
                str(tmp_path / "pv"),
            ]
        )
        exit_code = main(
            [
                "v2",
                "preview",
                "--design",
                str(tmp_path / "pv" / "deck_design.json"),
                "--output",
                str(tmp_path / "pv" / "preview.html"),
            ]
        )
        assert exit_code == 0
        html = (tmp_path / "pv" / "preview.html").read_text(encoding="utf-8")
        assert html.count('class="page"') == 8


class TestPlanConfirmFlow:
    def test_plan_deck_produces_editable_briefs_with_notes(self, tmp_path: Path) -> None:
        result = plan_deck(_request(tmp_path), MockLLMClient(), progress=lambda _: None)
        assert result.skeleton.total_pages == 20
        content = result.skeleton.content_slots()
        assert content and all(slot.brief is not None for slot in content)
        assert all(slot.brief.speaker_notes for slot in content)
        editable = editable_plan_from_skeleton(result.skeleton)
        assert editable.total_pages() == 20
        # planning stage checkpoints exist for later seeding
        checkpoints = Path(_request(tmp_path).output_dir) / "checkpoints"
        for name in ("brief.json", "skeleton.json", "skeleton_with_briefs.json"):
            assert (checkpoints / name).is_file()

    def test_confirmed_plan_drives_generation_verbatim(self, tmp_path: Path) -> None:
        plan_dir = tmp_path / "plan"
        result = plan_deck(
            _request(tmp_path, output_dir=str(plan_dir)),
            MockLLMClient(),
            progress=lambda _: None,
        )
        editable = editable_plan_from_skeleton(result.skeleton)
        editable.deck_title = "用户改过的标题"
        editable.sections[0].pages[0].title = "用户改过的第一页"
        editable.sections[0].pages[0].speaker_notes = "用户确认过的口播稿。"
        editable.sections[0].pages.append(
            EditablePage(title="用户新增的一页", speaker_notes="新增页口播。")
        )
        skeleton = skeleton_from_editable_plan(editable)

        build_dir = tmp_path / "build"
        checkpoints = build_dir / "checkpoints"
        checkpoints.mkdir(parents=True)
        (checkpoints / "brief.json").write_text(
            result.brief.model_copy(update={"deck_title": editable.deck_title})
            .model_dump_json(),
            encoding="utf-8",
        )
        for name in ("skeleton.json", "skeleton_with_briefs.json"):
            (checkpoints / name).write_text(skeleton.model_dump_json(), encoding="utf-8")

        built = build_deck(
            _request(
                tmp_path,
                output_dir=str(build_dir),
                page_count=skeleton.total_pages,
                resume=True,
            ),
            MockLLMClient(),
            progress=lambda _: None,
        )
        assert built.page_count == skeleton.total_pages == 21
        design = json.loads(Path(built.deck_design_path).read_text(encoding="utf-8"))
        assert design["deck_title"] == "用户改过的标题"
        titles = [page["title"] for page in design["pages"]]
        assert "用户改过的第一页" in titles
        assert "用户新增的一页" in titles
        notes = {page["title"]: page.get("speaker_notes") for page in design["pages"]}
        assert notes["用户改过的第一页"] == "用户确认过的口播稿。"

        pptx = Presentation(built.pptx_path)
        all_notes = [
            slide.notes_slide.notes_text_frame.text
            for slide in pptx.slides
            if slide.has_notes_slide
        ]
        assert "用户确认过的口播稿。" in all_notes


class TestAnchorVariety:
    def test_all_anchor_variants_pass_qa_on_all_builtin_themes(self) -> None:
        from ppt_agent.v2.anchors import (
            CLOSING_VARIANTS,
            COVER_VARIANTS,
            DIVIDER_VARIANTS,
            build_closing_page,
            build_cover_page,
            build_section_divider,
        )
        from ppt_agent.v2.design import BUILTIN_THEMES
        from ppt_agent.v2.qa import review_page

        for theme in BUILTIN_THEMES.values():
            for variant in range(COVER_VARIANTS):
                page = build_cover_page(
                    page_number=1, deck_title="生态环境保护", subtitle="副标题",
                    language="zh-CN", theme=theme, variant=variant,
                )
                _, qa = review_page(page, theme)
                assert not qa.errors, (theme.name, "cover", variant, qa.errors)
            for variant in range(DIVIDER_VARIANTS):
                page = build_section_divider(
                    page_number=4, section_index=2, section_count=6,
                    section_title="现状与挑战", section_goal="解释为什么现在必须行动",
                    language="zh-CN", theme=theme, variant=variant,
                )
                _, qa = review_page(page, theme)
                assert not qa.errors, (theme.name, "divider", variant, qa.errors)
            for variant in range(CLOSING_VARIANTS):
                page = build_closing_page(
                    page_number=20, deck_title="生态环境保护",
                    language="zh-CN", theme=theme, variant=variant,
                )
                _, qa = review_page(page, theme)
                assert not qa.errors, (theme.name, "closing", variant, qa.errors)

    def test_fallback_variant_differs_across_deck_titles(self) -> None:
        from ppt_agent.v2.anchors import COVER_VARIANTS, anchor_variant_seed

        variants = {
            anchor_variant_seed(title) % COVER_VARIANTS
            for title in ("生态环境保护", "AI 产品方法论", "澳洲留学全攻略", "季度经营复盘", "量子计算入门")
        }
        assert len(variants) > 1

    def test_build_uses_model_designed_anchors(self, tmp_path: Path) -> None:
        result = build_deck(_request(tmp_path), MockLLMClient(), progress=lambda _: None)
        design = json.loads(Path(result.deck_design_path).read_text(encoding="utf-8"))
        cover = design["pages"][0]
        assert cover["role"] == "cover"
        assert any(element["id"] == "mock_cover_title" for element in cover["elements"])
        dividers = [page for page in design["pages"] if page["role"] == "section_divider"]
        assert dividers
        assert all(
            any(element["id"] == "mock_div_title" for element in page["elements"])
            for page in dividers
        )

    def test_anchor_design_failure_falls_back_to_variant_library(self, tmp_path: Path) -> None:
        class NoAnchorClient(MockLLMClient):
            async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
                if task == "anchor_design":
                    raise RuntimeError("anchor designer offline")
                return await super().complete_json(task=task, **kwargs)

        result = build_deck(_request(tmp_path), NoAnchorClient(), progress=lambda _: None)
        assert result.status == "succeeded"
        design = json.loads(Path(result.deck_design_path).read_text(encoding="utf-8"))
        cover = design["pages"][0]
        assert cover["role"] == "cover"
        assert all(not element["id"].startswith("mock_") for element in cover["elements"])
        report = json.loads(Path(result.run_report_path).read_text(encoding="utf-8"))
        cover_outcome = [item for item in report["outcomes"] if item["page_number"] == 1][0]
        assert cover_outcome["status"] == "anchor"

    def test_style_signature_reaches_page_prompts(self) -> None:
        from ppt_agent.v2 import prompts
        from ppt_agent.v2.design import BUILTIN_THEMES
        from ppt_agent.v2.planning import ContentBrief, PageBrief

        theme = BUILTIN_THEMES["ink"]
        brief = ContentBrief(topic="主题", deck_title="标题")
        prompt = prompts.build_page_design_user_prompt(
            brief=brief, theme=theme, deck_title="标题", section_title="章节",
            page_brief=PageBrief(title="页"), page_number=5, total_pages=20,
            neighbor_titles=[],
        )
        assert "DECK STYLE SIGNATURE" in prompt
        assert theme.style.composition in prompt
        anchor_prompt = prompts.build_anchor_design_user_prompt(
            kind="cover", brief=brief, theme=theme, deck_title="标题",
        )
        assert theme.style.cover_concept in anchor_prompt


class TestReviseDeck:
    def _built_deck(self, tmp_path: Path):
        request = _request(tmp_path, page_count=12, deck_name="deck")
        build_deck(request, MockLLMClient(), progress=lambda _: None)
        return Path(request.output_dir)

    def test_revise_single_page_keeps_others(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        output_dir = self._built_deck(tmp_path)
        before = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="第 5 页信息太密，请改成图表页",
            client=MockLLMClient(), progress=lambda _: None,
        )
        assert result.revised_pages == [5]
        assert not result.theme_changed
        after = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        page5 = [page for page in after["pages"] if page["page_number"] == 5][0]
        assert any(element["id"] == "mock_revision_tag" for element in page5["elements"])
        page6_before = [p for p in before["pages"] if p["page_number"] == 6][0]
        page6_after = [p for p in after["pages"] if p["page_number"] == 6][0]
        assert page6_before["elements"] == page6_after["elements"]
        assert "已更新第 5 页" in result.reply

    def test_global_recolor_changes_theme_without_page_redesign(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        output_dir = self._built_deck(tmp_path)
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="整体配色换成深蓝色的执行风格",
            client=MockLLMClient(), progress=lambda _: None,
        )
        assert result.theme_changed
        assert result.revised_pages == []
        after = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        assert after["theme"]["palette"]["primary"] == "#23415E"
        assert Path(result.pptx_path).is_file()

    def test_content_level_rewrite_updates_brief_and_notes(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        output_dir = self._built_deck(tmp_path)
        skeleton_payload = json.loads(
            (output_dir / "checkpoints" / "skeleton_with_briefs.json").read_text(encoding="utf-8")
        )
        target_page = next(
            slot["page_number"] for slot in skeleton_payload["slots"] if slot["kind"] == "content"
        )

        class ContentRewriteClient(MockLLMClient):
            async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
                if task == "revision_plan":
                    return {
                        "reply": "我会重写这一页的内容。",
                        "theme_instruction": None,
                        "pages": [
                            {
                                "page_number": target_page,
                                "instruction": "Rewrite for investors",
                                "new_brief": {
                                    "title": "面向投资人的增长故事",
                                    "summary": "",
                                    "points": ["市场规模", "增长引擎"],
                                    "layout_hint": "stats",
                                    "data_idea": None,
                                    "speaker_notes": "投资人版本的口播稿。",
                                },
                            }
                        ],
                    }
                return await super().complete_json(task=task, **kwargs)

        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="第一个内容页改成面向投资人",
            client=ContentRewriteClient(), progress=lambda _: None,
        )
        assert result.revised_pages == [target_page]
        after = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        page = [item for item in after["pages"] if item["page_number"] == target_page][0]
        assert page["title"] == "面向投资人的增长故事"
        assert page["speaker_notes"] == "投资人版本的口播稿。"
        skeleton = json.loads(
            (output_dir / "checkpoints" / "skeleton_with_briefs.json").read_text(encoding="utf-8")
        )
        slot = [item for item in skeleton["slots"] if item["page_number"] == target_page][0]
        assert slot["brief"]["title"] == "面向投资人的增长故事"


class TestLayoutQualityGates:
    def _theme(self):
        from ppt_agent.v2.design import BUILTIN_THEMES

        return BUILTIN_THEMES["aurora"]

    def test_text_on_chart_is_an_error(self) -> None:
        from ppt_agent.v2.ir import ChartItem, ChartSeries
        from ppt_agent.v2.qa import review_page

        page = PageDesign(
            page_number=5,
            title="质量场景",
            elements=[
                TextItem(id="t", frame=Frame(x=64, y=60, w=700, h=50), text="质量场景", role="title"),
                ChartItem(
                    id="chart",
                    frame=Frame(x=100, y=160, w=800, h=420),
                    chart="column",
                    categories=["A", "B", "C"],
                    series=[ChartSeries(name="s", values=[1, 2, 3])],
                ),
                TextItem(id="stat", frame=Frame(x=350, y=300, w=300, h=90), text="32%", role="stat"),
                TextItem(id="side", frame=Frame(x=960, y=160, w=240, h=420), text="旁注" * 20, role="body"),
            ],
        )
        _, qa = review_page(page, self._theme())
        assert any(issue.code == "heavy_overlap" for issue in qa.errors)

    def test_undersized_table_is_an_error(self) -> None:
        from ppt_agent.v2.ir import TableItem
        from ppt_agent.v2.qa import review_page

        page = PageDesign(
            page_number=7,
            title="对比",
            elements=[
                TextItem(id="t", frame=Frame(x=64, y=60, w=700, h=50), text="对比", role="title"),
                TableItem(
                    id="tbl",
                    frame=Frame(x=64, y=160, w=1000, h=88),
                    headers=["维度", "现状", "目标"],
                    rows=[["a", "b", "c"], ["d", "e", "f"], ["g", "h", "i"]],
                ),
                TextItem(id="below", frame=Frame(x=64, y=300, w=1000, h=300), text="内容" * 40, role="body"),
            ],
        )
        _, qa = review_page(page, self._theme())
        assert any(issue.code == "table_overflow" for issue in qa.errors)

    def test_sparse_layout_is_an_error(self) -> None:
        from ppt_agent.v2.qa import review_page

        page = PageDesign(
            page_number=9,
            role="stats",
            title="孤独的卡片",
            elements=[
                TextItem(id="t", frame=Frame(x=64, y=60, w=500, h=50), text="孤独的卡片", role="title"),
                TextItem(id="a", frame=Frame(x=64, y=180, w=280, h=90), text="一点点内容", role="body"),
                TextItem(id="b", frame=Frame(x=64, y=290, w=280, h=60), text="再一点", role="body_small"),
            ],
        )
        _, qa = review_page(page, self._theme())
        assert any(
            issue.code in ("layout_sparse", "bottom_half_empty") for issue in qa.errors
        )

    def test_dense_archetype_page_stays_clean(self) -> None:
        from ppt_agent.v2.fallback import design_fallback_page
        from ppt_agent.v2.planning import PageBrief
        from ppt_agent.v2.qa import review_page

        page = design_fallback_page(
            PageBrief(title="正常内容页", points=["现状", "动作", "结果", "度量"]),
            page_number=6,
            section_title="章节",
            language="zh-CN",
        )
        _, qa = review_page(page, self._theme())
        assert not qa.errors


class TestAnchorHeroGate:
    def test_pale_content_like_cover_is_rejected(self, tmp_path: Path) -> None:
        class PaleCoverClient(MockLLMClient):
            async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
                if task == "anchor_design" and (kwargs.get("context") or {}).get("kind") == "cover":
                    return {
                        "role": "cover",
                        "title": "看起来像内容页的封面",
                        "background": "background",
                        "elements": [
                            {"type": "shape", "id": "panel", "frame": {"x": 48, "y": 48, "w": 568, "h": 188}, "shape": "rounded_rectangle", "fill": "surface"},
                            {"type": "text", "id": "kicker", "frame": {"x": 80, "y": 76, "w": 480, "h": 32}, "text": "培训", "role": "kicker"},
                            {"type": "text", "id": "title", "frame": {"x": 80, "y": 116, "w": 492, "h": 70}, "text": "看起来像内容页的封面", "role": "title"},
                            {"type": "text", "id": "sub", "frame": {"x": 80, "y": 190, "w": 500, "h": 36}, "text": "副标题", "role": "subtitle"},
                            {"type": "text", "id": "detail", "frame": {"x": 88, "y": 318, "w": 468, "h": 286}, "text": "终端 / IDE 明细", "role": "body"},
                            {"type": "text", "id": "labels", "frame": {"x": 692, "y": 142, "w": 476, "h": 430}, "text": "架构调用图", "role": "body_small"},
                        ],
                    }
                return await super().complete_json(task=task, **kwargs)

        result = build_deck(_request(tmp_path), PaleCoverClient(), progress=lambda _: None)
        design = json.loads(Path(result.deck_design_path).read_text(encoding="utf-8"))
        cover = design["pages"][0]
        # Replaced by the archetype library: dark hero background, no mock panel.
        assert cover["background"] != "background" or cover.get("background_gradient")
        assert all(element["id"] != "labels" for element in cover["elements"])
        report = json.loads(Path(result.run_report_path).read_text(encoding="utf-8"))
        assert [o for o in report["outcomes"] if o["page_number"] == 1][0]["status"] == "anchor"

    def test_rejection_reasons(self) -> None:
        from ppt_agent.v2.design import BUILTIN_THEMES
        from ppt_agent.v2.orchestrator import _anchor_design_rejection

        theme = BUILTIN_THEMES["aurora"]
        good = PageDesign(
            page_number=1,
            role="cover",
            title="标题",
            background="primary",
            elements=[
                TextItem(id="t", frame=Frame(x=96, y=248, w=860, h=180), text="大标题", role="display", color="on_primary"),
            ],
        )
        assert _anchor_design_rejection(good, "cover", theme) is None
        pale = good.model_copy(update={"background": "background"})
        assert "dark hero" in _anchor_design_rejection(pale, "cover", theme)
        no_hero = good.model_copy(
            update={"elements": [TextItem(id="t", frame=Frame(x=96, y=248, w=860, h=60), text="小标题", role="body", color="on_primary")]}
        )
        assert "hero title" in _anchor_design_rejection(no_hero, "cover", theme)


class TestRevisionHonesty:
    def _built_deck(self, tmp_path: Path):
        request = _request(tmp_path, page_count=12, deck_name="deck")
        build_deck(request, MockLLMClient(), progress=lambda _: None)
        return Path(request.output_dir)

    def test_hide_page_numbers_changes_chrome_and_rerenders(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        output_dir = self._built_deck(tmp_path)
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="所有页面的页码都去掉",
            client=MockLLMClient(), progress=lambda _: None,
        )
        assert result.theme_changed
        assert "页码/页脚显示" in result.reply
        design = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        assert design["theme"]["chrome"]["show_page_number"] is False
        # The rendered PPTX no longer stamps page numbers on content pages.
        pptx = Presentation(result.pptx_path)
        texts = [
            shape.text_frame.text
            for slide in pptx.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert not any(text.strip() == "05" for text in texts)

    def test_impossible_request_admits_no_change(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        class HonestButUselessPlanClient(MockLLMClient):
            async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
                if task == "revision_plan":
                    return {
                        "reply": "这个请求超出我能修改的范围。",
                        "theme_instruction": None,
                        "pages": [],
                    }
                return await super().complete_json(task=task, **kwargs)

        output_dir = self._built_deck(tmp_path)
        pptx_path = output_dir / "deck.pptx"
        before_mtime = pptx_path.stat().st_mtime_ns
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="给每一页加一个视频背景",
            client=HonestButUselessPlanClient(), progress=lambda _: None,
        )
        assert result.revised_pages == []
        assert not result.theme_changed
        assert "没有产生实际改动" in result.reply
        assert "已更新" not in result.reply and "已调整" not in result.reply
        assert pptx_path.stat().st_mtime_ns == before_mtime  # untouched

    def test_identical_theme_revision_is_not_reported_as_change(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        class EchoThemeClient(MockLLMClient):
            async def complete_json(self, *, task: str, **kwargs: Any) -> Any:
                if task == "revision_plan":
                    return {
                        "reply": "我会调整全局风格。",
                        "theme_instruction": "Do something global",
                        "pages": [],
                    }
                if task == "theme_revise":
                    return dict((kwargs.get("context") or {}).get("theme") or {})
                return await super().complete_json(task=task, **kwargs)

        output_dir = self._built_deck(tmp_path)
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="整体再优化一下",
            client=EchoThemeClient(), progress=lambda _: None,
        )
        assert not result.theme_changed
        assert "没有产生实际改动" in result.reply

    def test_all_pages_instruction_redesigns_every_page(self, tmp_path: Path) -> None:
        from ppt_agent.v2.revise import revise_deck

        output_dir = self._built_deck(tmp_path)
        result = revise_deck(
            output_dir=output_dir, deck_name="deck",
            message="每一页左侧的装饰竖条都去掉，别弄成空白",
            client=MockLLMClient(), progress=lambda _: None,
        )
        design = json.loads((output_dir / "deck_design.json").read_text(encoding="utf-8"))
        toc_pages = [p["page_number"] for p in design["pages"] if p["role"] == "toc"]
        expected = [p["page_number"] for p in design["pages"] if p["page_number"] not in toc_pages]
        assert result.revised_pages == expected
        content_pages = [p for p in design["pages"] if p["role"] not in ("cover", "toc", "section_divider", "closing")]
        assert all(
            any(e.get("id") == "mock_revision_tag" for e in page["elements"])
            for page in content_pages
        )
        assert "已更新第" in result.reply
