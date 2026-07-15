import json
import subprocess
from pathlib import Path

from fastapi.testclient import TestClient

import ppt_agent.api as api
from ppt_agent.generation import GenerationAttempt, GenerationResult
from ppt_agent.job_store import JobStore
from ppt_agent.load import load_deck, load_patch
from ppt_agent.long_deck_orchestrator import BatchRunReport, LongDeckRunReport
from ppt_agent.long_deck_render import LongDeckRenderReport
from ppt_agent.patch import build_patchable_elements_report
from ppt_agent.pipeline import BuildArtifact, BuildPipelineResult
from ppt_agent.ppt_master_execution import PPT_MASTER_EXECUTION_PLAN_ARTIFACT
from ppt_agent.qa import analyze_deck
from ppt_agent.patch import apply_patch
from ppt_agent.ppt_master_output import (
    PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
    PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
    register_ppt_master_output_artifacts,
)
from ppt_agent.ppt_master_project import (
    PPT_MASTER_PROJECT_INSTRUCTIONS_ARTIFACT,
    PPT_MASTER_VISUAL_PROJECT_MANIFEST_ARTIFACT,
)
from ppt_agent.ppt_master_runner import PPT_MASTER_RUNNER_RESULT_ARTIFACT
from ppt_agent.v2.design import BUILTIN_THEMES
from ppt_agent.v2.ir import ChartItem, ChartSeries, DeckDesign, Frame, IconItem, PageDesign, ShapeItem, TextItem


EXAMPLES_DIR = Path(__file__).resolve().parents[1] / "examples"


class _FakeInterviewStructuredModel:
    def __init__(self, responses: list[dict]) -> None:
        self.responses = responses

    def invoke(self, prompt: str) -> dict:
        return self.responses.pop(0)


class _FakeInterviewModel:
    def __init__(self, responses: list[dict]) -> None:
        self.structured_model = _FakeInterviewStructuredModel(responses)

    def with_structured_output(self, schema):
        return self.structured_model


def _interview_decision(*, ready: bool = False) -> dict:
    brief = {
        "topic": "生态环境保护",
        "audience": "城市规划与环境工程专业学生" if ready else None,
        "slide_count": 36 if ready else None,
        "language": "zh-CN",
        "purpose": "专业课程分享" if ready else None,
        "tone": "专业、清晰",
        "visual_direction": "浅色背景，使用地图和指标图表" if ready else None,
        "content_focus": ["问题诊断", "治理方案"] if ready else [],
        "constraints": ["可编辑 PPTX"] if ready else [],
        "user_requirements": "生成 36 页中文可编辑生态环境保护 PPT。" if ready else None,
    }
    if ready:
        return {
            "status": "ready",
            "assistant_message": "需求已经足够具体，可以开始生成。",
            "question": None,
            "options": [],
            "brief": brief,
            "missing_fields": [],
            "confidence": 0.94,
            "auto_start": False,
        }
    return {
        "status": "clarifying",
        "assistant_message": "我先确认这份演示的主要用途。",
        "question": "你希望观众看完后获得什么？",
        "options": [
            {"option_id": "learn", "label": "理解基础知识", "description": "建立认知框架"},
            {"option_id": "decide", "label": "支持方案决策", "description": "比较可选路径"},
            {"option_id": "act", "label": "推动具体行动", "description": "形成执行计划"},
        ],
        "brief": brief,
        "missing_fields": ["audience", "purpose", "slide_count", "visual_direction"],
        "confidence": 0.42,
        "auto_start": False,
    }


def _client(tmp_path: Path) -> TestClient:
    return TestClient(api.create_app(data_dir=tmp_path))


def _job_payload() -> dict:
    return {
        "topic": "AI in Education",
        "audience": "university students",
        "slides": 3,
        "theme_path": str(EXAMPLES_DIR / "theme.json"),
        "min_qa_score": 80,
        "max_attempts": 2,
    }


def _long_deck_payload() -> dict:
    return {
        "topic": "AI 产品经理如何设计 Agent 产品",
        "audience": "准备进入 AI 产品岗位的 IT 硕士学生",
        "slide_count": 30,
        "language": "zh-CN",
        "deck_type": "technical_product_share",
        "user_requirements": "讲技术边界、用户需求分析、工作流设计、评估指标和落地风险。",
        "batch_size": 2,
        "max_batch_attempts": 1,
    }


def _mock_expected_ppt_master_root(tmp_path: Path) -> Path:
    root = tmp_path / "ppt-master"
    skill_dir = root / "skills" / "ppt-master"
    skill_dir.mkdir(parents=True, exist_ok=True)
    (skill_dir / "SKILL.md").write_text("# PPT Master Skill\n", encoding="utf-8")
    (skill_dir / "scripts").mkdir(parents=True, exist_ok=True)
    (root / "requirements.txt").write_text("# test requirements\n", encoding="utf-8")
    (root / "README.md").write_text("# ppt-master\n", encoding="utf-8")
    subprocess.run(["git", "init"], cwd=root, check=True, capture_output=True)
    subprocess.run(
        ["git", "remote", "add", "origin", "https://github.com/hugohe3/ppt-master.git"],
        cwd=root,
        check=True,
        capture_output=True,
    )
    return root


def _build_mock_ppt_master_output(output_dir: Path, *, slide_count: int = 3) -> Path:
    from pptx import Presentation

    output_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    while len(presentation.slides) < slide_count:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(50, 50, 400, 80)
        textbox.text_frame.text = f"Slide {len(presentation.slides)}"
    if len(presentation.slides) > slide_count:
        while len(presentation.slides) > slide_count:
            slide_id = presentation.slides._sldIdLst[-1].rId  # type: ignore[attr-defined]
            presentation.part.drop_rel(slide_id)
            del presentation.slides._sldIdLst[-1]  # type: ignore[attr-defined]
    presentation.save(output_dir / "generated_by_ppt_master.pptx")
    (output_dir / "generation_notes.md").write_text(
        "# Notes\n\nMode: Native DrawingML shapes (directly editable)\n",
        encoding="utf-8",
    )
    (output_dir / "agent_pm_recovery_ppt169_20260701" / "svg_output").mkdir(parents=True, exist_ok=True)
    return output_dir


def _build_mock_ppt_master_package(job_dir: Path) -> Path:
    package_dir = job_dir / "ppt_master_package"
    package_dir.mkdir(parents=True, exist_ok=True)
    (package_dir / "source.md").write_text("# Presentation Request\n", encoding="utf-8")
    (package_dir / "run_prompt.md").write_text("# PPT Master Local Job Prompt\n", encoding="utf-8")
    (package_dir / "README.md").write_text("# PPT Master Job Package\n", encoding="utf-8")
    (package_dir / "manifest.json").write_text('{"is_available": true}\n', encoding="utf-8")
    return package_dir


def _generation_result(accepted: bool = True) -> GenerationResult:
    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    qa_report = analyze_deck(deck)
    return GenerationResult(
        deck=deck,
        qa_report=qa_report,
        attempts=[
            GenerationAttempt(
                attempt_index=1,
                deck=deck,
                qa_report=qa_report,
                accepted=accepted,
            )
        ],
        accepted=accepted,
    )


def _fake_pipeline_result(output_dir: Path, accepted: bool = True) -> BuildPipelineResult:
    output_dir.mkdir(parents=True, exist_ok=True)
    generation_result = _generation_result(accepted=accepted)

    deck_path = output_dir / "generated_deck_ir.json"
    patchable_elements_path = output_dir / "patchable_elements.json"
    qa_path = output_dir / "generated_qa_report.json"
    attempts_path = output_dir / "generated_attempts.json"
    pptx_path = output_dir / "generated_deck.pptx"

    deck_path.write_text(generation_result.deck.model_dump_json(indent=2), encoding="utf-8")
    patchable_elements_path.write_text(
        build_patchable_elements_report(generation_result.deck).model_dump_json(indent=2),
        encoding="utf-8",
    )
    qa_path.write_text(generation_result.qa_report.model_dump_json(indent=2), encoding="utf-8")
    attempts_path.write_text(generation_result.model_dump_json(indent=2), encoding="utf-8")
    pptx_path.write_bytes(b"fake pptx")

    return BuildPipelineResult(
        generation_result=generation_result,
        artifacts=[
            BuildArtifact(name="generated_deck_ir", kind="json", path=deck_path),
            BuildArtifact(name="patchable_elements", kind="json", path=patchable_elements_path),
            BuildArtifact(name="generated_qa_report", kind="json", path=qa_path),
            BuildArtifact(name="generated_attempts", kind="json", path=attempts_path),
            BuildArtifact(name="generated_deck", kind="pptx", path=pptx_path),
        ],
        accepted=accepted,
        status_code=0 if accepted else 2,
        messages=[] if accepted else ["Generated Deck IR did not meet the QA score gate."],
    )


def _fake_pipeline_result_with_patch(output_dir: Path, accepted: bool = True) -> BuildPipelineResult:
    result = _fake_pipeline_result(output_dir, accepted=accepted)
    deck = result.generation_result.deck
    patch = load_patch(EXAMPLES_DIR / "sample_patch.json")
    patch_result = apply_patch(deck, patch)

    patched_deck_path = output_dir / "patched_deck_ir.json"
    patch_report_path = output_dir / "patch_report.json"
    patched_pptx_path = output_dir / "patched_deck.pptx"

    patched_deck_path.write_text(patch_result.deck.model_dump_json(indent=2), encoding="utf-8")
    patch_report_path.write_text(
        patch_result.model_copy(
            update={
                "input_patch_path": str(EXAMPLES_DIR / "sample_patch.json"),
                "output_pptx_path": str(patched_pptx_path),
            }
        ).model_dump_json(indent=2),
        encoding="utf-8",
    )
    patched_pptx_path.write_bytes(b"fake patched pptx")

    return result.model_copy(
        update={
            "patch_result": patch_result,
            "artifacts": [
                *result.artifacts,
                BuildArtifact(name="patched_deck_ir", kind="json", path=patched_deck_path),
                BuildArtifact(name="patch_report", kind="json", path=patch_report_path),
                BuildArtifact(name="patched_deck", kind="pptx", path=patched_pptx_path),
            ],
        }
    )


def _install_fake_backend(monkeypatch, accepted: bool = True) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_build_pipeline(model, request, **kwargs):
        return _fake_pipeline_result(request.output_dir, accepted=accepted)

    monkeypatch.setattr(api, "run_build_pipeline", fake_run_build_pipeline)


def _fake_long_deck_run_report(request) -> LongDeckRunReport:
    output_dir = request.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)
    batch_dir = output_dir / "batches"
    batch_dir.mkdir(parents=True, exist_ok=True)

    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    plan_path = output_dir / "generated_long_deck_plan.json"
    merged_path = output_dir / "generated_long_deck_ir.json"
    qa_path = output_dir / "generated_long_deck_qa.json"
    quality_gate_path = output_dir / "generated_long_deck_quality_gate.json"
    run_report_path = output_dir / "long_deck_run_report.json"
    batch_status_path = batch_dir / "batch_01_status.json"
    batch_deck_path = batch_dir / "batch_01_deck_ir.json"
    batch_qa_path = batch_dir / "batch_01_qa_report.json"
    batch_attempts_path = batch_dir / "batch_01_attempts.json"

    plan_path.write_text('{"sections":[],"batches":[]}\n', encoding="utf-8")
    merged_path.write_text(deck.model_dump_json(indent=2), encoding="utf-8")
    qa_path.write_text('{"score":0.82,"passed":true}\n', encoding="utf-8")
    quality_gate_path.write_text('{"status":"passed","score":82,"issues":[],"blocked_codes":[],"blocked_slide_ids":[],"blocked_element_ids":[],"message":"passed"}\n', encoding="utf-8")
    batch_deck_path.write_text(deck.model_dump_json(indent=2), encoding="utf-8")
    batch_qa_path.write_text('{"score":82}\n', encoding="utf-8")
    batch_attempts_path.write_text('{"attempts":[]}\n', encoding="utf-8")
    batch_status_path.write_text('{"status":"succeeded"}\n', encoding="utf-8")

    batch_report = BatchRunReport(
        batch_id="batch_01",
        start_slide=1,
        end_slide=request.batch_size,
        status="succeeded",
        deck_ir_path=batch_deck_path,
        qa_report_path=batch_qa_path,
        attempts_path=batch_attempts_path,
        status_path=batch_status_path,
    )
    report = LongDeckRunReport(
        run_id="test-long-deck-run",
        slide_count=request.slide_count,
        batch_size=request.batch_size,
        total_batches=15,
        completed_batches=["batch_01"],
        failed_batches=[],
        status="succeeded",
        batch_reports=[batch_report],
        merged_deck_ir_path=merged_path,
        long_deck_qa_path=qa_path,
        long_deck_quality_gate_path=quality_gate_path,
        long_deck_plan_path=plan_path,
        run_report_path=run_report_path,
    )
    run_report_path.write_text(report.model_dump_json(indent=2), encoding="utf-8")
    return report


def _inject_recovery_source_noise(deck_ir_path: Path) -> None:
    payload = json.loads(deck_ir_path.read_text(encoding="utf-8"))
    payload["slides"][0]["title"] = "slide_id: leaked internal field"
    payload["slides"][0]["elements"][0]["text"] = "risk: 把这一点转化为明确的下一步行动"
    payload["slides"][0]["elements"][1]["text"] = (
        "Impact：先列出 Agent 不允许自动执行的动作\n"
        "判断点 1：用户任务要先被拆成可执行工作流"
    )
    payload["slides"][1]["elements"][0]["text"] = "Mitigation：用灰度和回滚控制上线质量"
    payload["slides"][1]["elements"][1]["text"] = "Option A: element_id should not leak"
    deck_ir_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _install_fake_long_deck_backend(
    monkeypatch,
    captured: dict | None = None,
    *,
    ppt_master_root: Path | None = None,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    if ppt_master_root is not None:
        monkeypatch.setenv("PPT_MASTER_DIR", str(ppt_master_root))
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_long_deck_batch_generation(request, model, *, progress_logger=None, cancel_checker=None):
        if captured is not None:
            captured["request"] = request
            captured.setdefault("requests", []).append(request)
            captured["model"] = model
            captured["cancel_checker"] = cancel_checker
        if progress_logger is not None:
            progress_logger("Starting long deck run: 30 slides, batch_size=2, total_batches=15")
            progress_logger("Starting batch_01 slides 1-2")
            progress_logger("Completed batch_01 in 0.1s")
            progress_logger("Merging 15 batches")
            progress_logger("Running long deck QA")
            progress_logger("Running long deck hard quality gate")
            progress_logger("Long deck run succeeded")
        return _fake_long_deck_run_report(request)

    def fake_render_long_deck_ir_to_pptx(
        input_deck_ir_path,
        output_pptx_path,
        report_path,
        *,
        theme_path,
        assets_dir=None,
    ):
        if captured is not None:
            captured["render_input"] = input_deck_ir_path
            captured["render_output"] = output_pptx_path
        Path(output_pptx_path).write_bytes(b"fake long deck pptx")
        report = LongDeckRenderReport(
            status="succeeded",
            input_deck_ir_path=Path(input_deck_ir_path),
            output_pptx_path=Path(output_pptx_path),
            slide_count=30,
            generated_at="2026-06-30T00:00:00+00:00",
        )
        Path(report_path).write_text(report.model_dump_json(indent=2), encoding="utf-8")
        return report

    monkeypatch.setattr(api, "run_long_deck_batch_generation", fake_run_long_deck_batch_generation)
    monkeypatch.setattr(api, "render_long_deck_ir_to_pptx", fake_render_long_deck_ir_to_pptx)


def _install_fake_v2_long_deck_backend(
    monkeypatch,
    captured: dict | None = None,
    *,
    status: str = "succeeded",
) -> None:
    monkeypatch.setenv("PPT_AGENT_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_v2_model_client", lambda: object())

    def fake_build_v2_deck(request, client, *, search_provider=None, progress=print):
        if captured is not None:
            captured["request"] = request
            captured.setdefault("requests", []).append(request)
            captured["client"] = client
        output_dir = Path(request.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        design_path = output_dir / f"{request.deck_name}_design.json"
        qa_path = output_dir / f"{request.deck_name}_qa_report.json"
        run_report_path = output_dir / f"{request.deck_name}_run_report.json"
        pptx_path = output_dir / f"{request.deck_name}.pptx"
        design_path.write_text(
            json.dumps({"deck_title": "Test", "pages": request.page_count}),
            encoding="utf-8",
        )
        pages_with_errors = 1 if status == "quality_gate_failed" else 0
        qa_path.write_text(
            json.dumps(
                {
                    "total_pages": request.page_count,
                    "pages_with_errors": pages_with_errors,
                    "pages_with_warnings": 0,
                    "auto_fix_count": 0,
                    "repaired_pages": [],
                    "fallback_pages": [],
                    "results": [],
                }
            ),
            encoding="utf-8",
        )
        run_report_path.write_text(
            json.dumps({"status": status, "page_count": request.page_count}),
            encoding="utf-8",
        )
        result_pptx_path = None
        if status != "quality_gate_failed":
            pptx_path.write_bytes(b"fake v2 pptx")
            result_pptx_path = str(pptx_path)
        progress("[stage] intake finished in 0.1s")
        progress("[brief] 'Test' | language=zh-CN")
        progress("[theme] technical (motif: grid)")
        progress(f"[outline] 8 sections / {request.page_count} pages")
        progress("[stage] page_briefs finished in 0.1s")
        progress(f"[design] {max(1, request.page_count - 10)}/{max(1, request.page_count - 10)} content pages done")
        progress("[stage] page_designs finished in 0.1s")
        progress("[stage] assemble_qa finished in 0.1s")
        if result_pptx_path is not None:
            progress("[stage] render finished in 0.1s")
        return api.V2BuildResult(
            status=status,
            pptx_path=result_pptx_path,
            deck_design_path=str(design_path),
            qa_report_path=str(qa_path),
            run_report_path=str(run_report_path),
            page_count=request.page_count,
            model_pages=request.page_count,
            repaired_pages=0,
            fallback_pages=0,
            usage={"estimated_cost_usd": 1.25},
            stage_seconds={"render": 0.1},
        )

    monkeypatch.setattr(api, "build_v2_deck", fake_build_v2_deck)


def _install_fake_long_deck_quality_gate_failure(monkeypatch, captured: dict | None = None) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_long_deck_batch_generation(request, model, *, progress_logger=None, cancel_checker=None):
        if captured is not None:
            captured["request"] = request
        if progress_logger is not None:
            progress_logger("Starting long deck run: 30 slides, batch_size=2, total_batches=15")
            progress_logger("Starting batch_01 slides 1-2")
            progress_logger("Completed batch_01 in 0.1s")
            progress_logger("Merging 15 batches")
            progress_logger("Running long deck QA")
            progress_logger("Running long deck hard quality gate")
            progress_logger("Long deck run failed_quality_gate")
        report = _fake_long_deck_run_report(request)
        if report.merged_deck_ir_path is not None:
            _inject_recovery_source_noise(report.merged_deck_ir_path)
        quality_gate_path = request.output_dir / "generated_long_deck_quality_gate.json"
        quality_gate_path.write_text(
            json.dumps(
                {
                    "status": "failed_quality_gate",
                    "score": 68,
                    "issues": [
                        {
                            "severity": "error",
                            "slide_id": "slide_021",
                            "element_id": "s021_e01",
                            "code": "instruction_leakage",
                            "message": "meta leakage",
                        }
                    ],
                    "blocked_codes": ["instruction_leakage"],
                    "blocked_slide_ids": ["slide_021"],
                    "blocked_element_ids": ["s021_e01"],
                    "message": "Long-deck hard quality gate failed because audience-visible instruction leakage or matrix placeholder content was detected.",
                },
                ensure_ascii=False,
            )
            + "\n",
            encoding="utf-8",
        )
        report = report.model_copy(
            update={
                "status": "failed_quality_gate",
                "long_deck_quality_gate_path": quality_gate_path,
                "error_message": "Long-deck hard quality gate failed because audience-visible instruction leakage or matrix placeholder content was detected.",
                "error_type": "failed_quality_gate",
                "suggestion": "Fix instruction leakage or placeholder matrix content before rendering PPTX.",
            }
        )
        report.run_report_path.write_text(report.model_dump_json(indent=2), encoding="utf-8")
        return report

    def fail_if_render_called(*args, **kwargs):
        raise AssertionError("render_long_deck_ir_to_pptx should not be called after quality gate failure")

    monkeypatch.setattr(api, "run_long_deck_batch_generation", fake_run_long_deck_batch_generation)
    monkeypatch.setattr(api, "render_long_deck_ir_to_pptx", fail_if_render_called)


def _install_fake_long_deck_timeout_before_merge(monkeypatch, captured: dict | None = None) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_long_deck_batch_generation(request, model, *, progress_logger=None, cancel_checker=None):
        if captured is not None:
            captured["request"] = request
        if progress_logger is not None:
            progress_logger("Starting long deck run: 30 slides, batch_size=2, total_batches=15")
            for batch_number in range(1, 13):
                progress_logger(f"Starting batch_{batch_number:02d} slides 1-2")
                progress_logger(f"Completed batch_{batch_number:02d} in 0.1s")
            progress_logger("Starting batch_13 slides 25-26")
        raise TimeoutError(
            f"Job timed out after {api._long_deck_job_timeout_seconds()} seconds while running stage "
            "'generating_batch_13_of_15'."
        )

    monkeypatch.setattr(api, "run_long_deck_batch_generation", fake_run_long_deck_batch_generation)


def test_health_endpoint(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/health")

    assert response.status_code == 200
    assert response.json() == {"status": "ok"}


def test_index_page_returns_html_with_generate_button(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    assert "text/html" in response.headers["content-type"]
    assert 'id="generateLongDeckButton"' in response.text
    assert "开始生成 PPT" in response.text


def test_index_page_contains_chinese_job_labels(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    for text in [
        "创建演示",
        "主题",
        "目标观众",
        "页数",
        "生成进度",
        "生成文件",
        "资料覆盖度",
        "叙事健康度",
        "可编辑成片",
        "演示历史",
        "搜索主题、观众或任务 ID",
    ]:
        assert text in response.text


def test_index_page_contains_long_deck_product_workspace(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    for text in [
        "创建演示",
        "和 Agent 一起定义演示",
        "告诉我你想做什么演示",
        "不确定，暂时跳过",
        "问题数量动态调整",
        "像聊天一样说出想法",
        "Agent 理解充分后会直接准备生成",
        "Agent 已经理解，可以开始生成",
        "继续调整",
        "Agent 理解",
        "Live generation studio",
        "生成工作台",
        "实时画布",
        "页面会在这里实时出现",
        "Agent 正在做什么",
        "演示预览",
        "章节页数分配",
        "质量与成本",
        "开始生成 PPT",
        "取消任务",
        "继续/重试演示",
        "PPT Master 渲染包",
        "PPT Master 执行桥",
        "PPT Master Visual Project",
        "PPT Master 本地导出",
        "PPT Master 生成结果",
        "PPT Master Source Markdown",
        "PPT Master Run Prompt",
        "PPT Master Package Manifest",
        "PPT Master Package README",
        "PPT Master Execution Plan",
        "PPT Master Visual Project Manifest",
        "PPT Master Project Instructions",
        "PPT Master Runner Result",
        "PPT Master Generated PPTX",
        "PPT Master Generation Notes",
        "PPT Master Output Manifest",
        "准备 PPT Master 执行计划",
        "准备 PPT Master Visual Project",
        "运行 PPT Master 本地导出",
        "execution status",
        "bootstrap status",
        "expected pptx",
        "runner status",
        "需要外部 AI 生成 project",
        "原因",
        "建议",
        "package_mode",
        "quality_gate",
        "Recovery package 已生成",
        "系统会从已完成 batch 后继续",
        "当前阶段不会自动运行 ppt-master",
    ]:
        assert text in response.text
    assert 'id="longDeckForm"' in response.text
    assert 'id="interviewComposer"' in response.text
    assert 'id="interviewQuestionPanel"' in response.text
    assert 'id="interviewOptions"' in response.text
    assert 'id="generationConfirmation"' in response.text
    assert 'id="confirmGenerationButton"' in response.text
    assert 'id="continueInterviewButton"' in response.text
    assert "也可以在下方直接输入你的想法" not in response.text
    assert "不确定，暂时跳过" in response.text
    assert "interviewOptions.after(interviewComposer)" in response.text
    assert 'interviewComposer.setAttribute("aria-busy", String(isBusy))' in response.text
    assert "interviewInput.disabled = isBusy" not in response.text
    assert "正在快速整理这一轮需求" in response.text
    assert "Number(decision.brief.slide_count) <= 10" in response.text
    assert "longDeckForm.hidden = true" in response.text
    assert 'id="manualBriefButton"' not in response.text
    assert "manualBriefVisible" not in response.text
    assert 'id="briefStatus"' in response.text
    assert "/api/presentation-interviews" in response.text
    assert 'id="long_topic" name="topic" type="hidden"' in response.text
    assert 'id="long_audience" name="audience" type="hidden"' in response.text
    assert 'id="long_slide_count" name="slide_count" type="hidden"' in response.text
    assert 'id="long_user_requirements" name="user_requirements" type="hidden"' in response.text
    assert '<label>主题<input id="long_topic"' not in response.text
    assert '<label>目标观众<input id="long_audience"' not in response.text
    assert "PPT 详细要求<textarea" not in response.text
    assert "高级生成设置" not in response.text
    assert 'id="long_batch_size"' not in response.text
    assert 'id="long_max_batch_attempts"' not in response.text
    assert '/api/jobs/${id}/preview-slides/${slideNumber}' in response.text
    assert "manifest.highlight_slide_numbers" in response.text
    assert "正在展示视觉高光页" in response.text
    assert 'id="previewSlide1"' in response.text
    assert "<iframe" in response.text
    assert "ppt_agent_long_deck_form_draft" in response.text
    assert "schedulePoll(id, 1000)" in response.text
    assert "setInterval(renderElapsedClock, 250)" in response.text
    assert 'value="AI 产品经理如何设计 Agent 产品"' not in response.text
    assert "选择 30、50 或 100 页" not in response.text
    assert "普通 1-10 页生成器" not in response.text
    assert 'id="jobForm"' not in response.text
    assert 'id="min_qa_score"' not in response.text
    assert 'id="max_attempts"' not in response.text
    assert 'id="patch_path"' not in response.text
    assert 'submitJob("/api/jobs", buildShortDeckPayload())' in response.text
    assert 'submitJob("/api/long-deck-jobs", buildLongDeckPayload())' in response.text
    assert "pageCount <= 3" in response.text
    assert 'deck_type: "visual_design_v2"' in response.text
    assert 'data-view-target="create"' in response.text
    assert 'data-view-target="studio"' in response.text
    assert 'data-view-target="preview"' in response.text
    assert 'data-view-target="history"' in response.text
    assert 'data-view-target="delivery"' in response.text
    assert 'id="liveSlideThumbnails"' in response.text
    assert 'id="liveSlidePreview"' in response.text
    assert "renderLiveSlideWorkspace(id, manifest, available)" in response.text
    assert "IntersectionObserver" in response.text
    assert "observeLiveThumbnailPreviews();" in response.text
    assert "previewStartIndex" not in response.text
    assert "live-slide-thumbnail-summary" not in response.text
    assert "/api/presentations" in response.text


def test_index_page_keeps_only_unified_presentation_fields(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    for field_name in ["topic", "audience", "slide_count", "user_requirements"]:
        assert f'name="{field_name}"' in response.text


def test_index_page_contains_progress_stage_fields(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    assert "currentStage" in response.text
    assert "elapsedSeconds" in response.text
    assert "generate_deck_chunk_" in response.text
    assert "正在快速解析需求" in response.text
    assert "正在快速规划大纲" in response.text
    assert "正在生成 Deck：第" in response.text
    assert "已生成，但未通过 QA" in response.text
    assert "已生成，但 Patch 需要修正" in response.text
    assert "已生成，但 QA 和 Patch 仍需修正" in response.text
    assert "任务运行时间较长，请检查后端日志" in response.text
    assert "generating_batch_" in response.text
    assert "正在生成长 PPT：batch" in response.text
    assert "正在执行长 PPT质量门禁" in response.text
    assert "正在渲染长 PPT PPTX" in response.text
    assert "质量门禁失败" in response.text
    assert "currentBatch" in response.text
    assert "totalBatches" in response.text
    assert "completedBatches" in response.text
    assert "failedBatches" in response.text
    assert "取消请求已发送" in response.text


def test_index_page_patch_path_is_optional_json_placeholder(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/")

    assert response.status_code == 200
    assert 'id="patch_path"' not in response.text
    assert "sample_patch.js\"" not in response.text
    assert "sample_patch.js<" not in response.text


def test_create_job_without_api_key_returns_clear_error(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)

    response = _client(tmp_path).post("/api/jobs", json=_job_payload())

    assert response.status_code == 503
    assert "OPENAI_API_KEY is not set" in response.json()["detail"]


def test_requirements_interview_without_api_key_returns_clear_error(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)

    response = _client(tmp_path).post(
        "/api/presentation-interviews",
        json={"message": "我想做一份生态环境保护演示。"},
    )

    assert response.status_code == 503
    assert "OPENAI_API_KEY is not set" in response.json()["detail"]


def test_requirements_interview_model_uses_low_latency_gpt55_defaults(monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.delenv("OPENAI_MODEL", raising=False)
    monkeypatch.delenv("PPT_AGENT_INTERVIEW_REASONING_EFFORT", raising=False)
    monkeypatch.delenv("PPT_AGENT_INTERVIEW_MAX_TOKENS", raising=False)

    model = api._create_interview_chat_model()

    assert model.model_name == "gpt-5.5"
    assert model.reasoning_effort == "low"
    assert model.max_tokens == 1800
    assert model.max_retries == 1


def test_requirements_interview_persists_adaptive_question_and_options(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    model = _FakeInterviewModel([_interview_decision(), _interview_decision(ready=True)])
    monkeypatch.setattr(api, "_create_interview_chat_model", lambda: model)
    client = _client(tmp_path)

    started = client.post(
        "/api/presentation-interviews",
        json={"message": "我想做一份生态环境保护 PPT，但还没有想清楚。"},
    )

    assert started.status_code == 201
    first = started.json()
    assert first["status"] == "clarifying"
    assert first["turn_count"] == 1
    assert first["decision"]["question"] == "你希望观众看完后获得什么？"
    assert len(first["decision"]["options"]) == 3
    assert client.app.state.job_store.get_presentation_interview(first["interview_id"]) is not None

    continued = client.post(
        f"/api/presentation-interviews/{first['interview_id']}/messages",
        json={"message": "推动具体行动", "selected_option_id": "act"},
    )

    assert continued.status_code == 200
    final = continued.json()
    assert final["status"] == "ready"
    assert final["turn_count"] == 2
    assert final["decision"]["brief"]["slide_count"] == 36
    restored = client.get(f"/api/presentation-interviews/{first['interview_id']}")
    assert restored.json() == final


def test_ready_requirements_interview_can_continue_with_natural_language_adjustments(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    model = _FakeInterviewModel([_interview_decision(ready=True), _interview_decision(ready=True)])
    monkeypatch.setattr(api, "_create_interview_chat_model", lambda: model)
    client = _client(tmp_path)

    started = client.post(
        "/api/presentation-interviews",
        json={"message": "生成一份 36 页生态环境保护课程演示。"},
    )
    continued = client.post(
        f"/api/presentation-interviews/{started.json()['interview_id']}/messages",
        json={"message": "改成 24 页，并增加真实案例。"},
    )

    assert started.status_code == 201
    assert started.json()["status"] == "ready"
    assert continued.status_code == 200
    assert continued.json()["status"] == "ready"
    assert continued.json()["turn_count"] == 2
    assert continued.json()["messages"][-2]["content"] == "改成 24 页，并增加真实案例。"


def test_create_long_deck_job_without_api_key_returns_clear_error(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)

    response = _client(tmp_path).post("/api/long-deck-jobs", json=_long_deck_payload())

    assert response.status_code == 503
    assert "OPENAI_API_KEY is not set" in response.json()["detail"]


def test_create_job_success_returns_job_id(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)

    response = _client(tmp_path).post("/api/jobs", json=_job_payload())

    assert response.status_code == 202
    body = response.json()
    assert body["job_id"]
    assert body["status"] == "pending"


def test_create_job_persists_request_in_presentation_history(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    client = _client(tmp_path)
    payload = {
        **_job_payload(),
        "user_requirements": "突出可编辑图表和课堂表达。",
        "interview_id": "interview-history-link",
    }

    job_id = client.post("/api/jobs", json=payload).json()["job_id"]
    history = client.get("/api/presentations").json()

    item = next(item for item in history["items"] if item["job_id"] == job_id)
    assert item["topic"] == payload["topic"]
    assert item["audience"] == payload["audience"]
    assert item["slide_count"] == payload["slides"]
    assert item["user_requirements"] == payload["user_requirements"]
    assert item["pptx_download_url"].startswith("/api/artifacts/")
    request = client.app.state.job_store.get_presentation_request(job_id)
    assert request is not None
    assert request.interview_id == "interview-history-link"


def test_presentation_history_supports_status_and_text_filters(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    matching = store.create_job(job_type="long_deck_v2")
    store.save_presentation_request(
        matching.job_id,
        topic="未来智慧校园",
        audience="高校管理者",
        user_requirements="技术产品蓝图",
        slide_count=100,
    )
    store.update_job(matching.job_id, status="succeeded", accepted=True, qa_score=96)
    other = store.create_job(job_type="short_deck")
    store.save_presentation_request(
        other.job_id,
        topic="零售经营周报",
        audience="门店经理",
        user_requirements="经营复盘",
        slide_count=8,
    )
    store.update_job(other.job_id, status="failed", accepted=False)
    client = TestClient(api.create_app(data_dir=tmp_path, store=store))

    response = client.get("/api/presentations", params={"status": "succeeded", "query": "高校"})

    assert response.status_code == 200
    body = response.json()
    assert body["total"] == 1
    assert body["items"][0]["job_id"] == matching.job_id


def test_existing_long_deck_request_is_backfilled_into_history(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    job = store.create_job(job_type="long_deck")
    job_dir = tmp_path / "jobs" / job.job_id
    job_dir.mkdir(parents=True)
    (job_dir / "long_deck_request.json").write_text(
        json.dumps(_long_deck_payload(), ensure_ascii=False),
        encoding="utf-8",
    )

    body = _client(tmp_path).get("/api/presentations").json()

    item = next(item for item in body["items"] if item["job_id"] == job.job_id)
    assert item["topic"] == _long_deck_payload()["topic"]
    assert item["slide_count"] == 30


def test_create_long_deck_job_rejects_slide_count_outside_supported_range(tmp_path: Path) -> None:
    client = _client(tmp_path)

    too_short = client.post(
        "/api/long-deck-jobs",
        json={**_long_deck_payload(), "slide_count": 3},
    )
    too_long = client.post(
        "/api/long-deck-jobs",
        json={**_long_deck_payload(), "slide_count": 101},
    )

    assert too_short.status_code == 422
    assert too_long.status_code == 422


def test_create_arbitrary_page_long_deck_uses_v2_pipeline(tmp_path: Path, monkeypatch) -> None:
    captured: dict = {}
    _install_fake_v2_long_deck_backend(monkeypatch, captured)

    response = _client(tmp_path).post(
        "/api/long-deck-jobs",
        json={**_long_deck_payload(), "slide_count": 75},
    )

    assert response.status_code == 202
    assert captured["request"].page_count == 75


def test_visual_design_web_requests_use_v2_for_four_to_thirty_pages(tmp_path: Path, monkeypatch) -> None:
    captured: dict = {}
    _install_fake_v2_long_deck_backend(monkeypatch, captured)
    client = _client(tmp_path)

    for page_count in (4, 30):
        response = client.post(
            "/api/long-deck-jobs",
            json={
                **_long_deck_payload(),
                "slide_count": page_count,
                "deck_type": "visual_design_v2",
            },
        )

        assert response.status_code == 202
        assert captured["request"].page_count == page_count
        assert client.get(f"/api/jobs/{response.json()['job_id']}").json()["job_type"] == "long_deck_v2"


def test_create_100_page_long_deck_uses_v2_pipeline_and_registers_artifacts(
    tmp_path: Path,
    monkeypatch,
) -> None:
    captured: dict = {}
    _install_fake_v2_long_deck_backend(monkeypatch, captured)
    client = _client(tmp_path)

    response = client.post(
        "/api/long-deck-jobs",
        json={**_long_deck_payload(), "slide_count": 100},
    )

    assert response.status_code == 202
    job_id = response.json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    artifact_names = {artifact["name"] for artifact in artifacts}
    assert body["status"] == "succeeded"
    assert body["accepted"] is True
    assert body["job_type"] == "long_deck_v2"
    assert body["current_stage"] == "v2_completed"
    assert body["total_batches"] == 100
    assert body["completed_batches"] == 100
    assert body["qa_score"] == 100
    assert body["ppt_master_package"] is None
    assert captured["request"].page_count == 100
    assert captured["request"].concurrency == api.DEFAULT_V2_CONCURRENCY
    assert captured["request"].qa_gate == "strict"
    assert {
        "long_deck_request",
        "generated_long_deck_v2",
        "generated_long_deck_v2_design",
        "generated_long_deck_v2_qa_report",
        "generated_long_deck_v2_run_report",
    } <= artifact_names
    assert not any(name.startswith("page_") for name in artifact_names)


def test_v2_quality_gate_failure_does_not_register_pptx(tmp_path: Path, monkeypatch) -> None:
    _install_fake_v2_long_deck_backend(monkeypatch, status="quality_gate_failed")
    client = _client(tmp_path)

    job_id = client.post(
        "/api/long-deck-jobs",
        json={**_long_deck_payload(), "slide_count": 50},
    ).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifact_names = {
        artifact["name"]
        for artifact in client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    }

    assert body["status"] == "failed_quality_gate"
    assert body["accepted"] is False
    assert body["current_stage"] == "v2_quality_gate_failed"
    assert body["qa_score"] == 98
    assert "generated_long_deck_v2" not in artifact_names
    assert "generated_long_deck_v2_design" in artifact_names
    assert "generated_long_deck_v2_qa_report" in artifact_names


def test_create_long_deck_job_defaults_batch_size_to_two(tmp_path: Path, monkeypatch) -> None:
    captured: dict = {}
    _install_fake_long_deck_backend(monkeypatch, captured)
    payload = _long_deck_payload()
    payload.pop("batch_size")

    response = _client(tmp_path).post("/api/long-deck-jobs", json=payload)

    assert response.status_code == 202
    assert captured["request"].batch_size == 2
    assert captured["request"].slide_count == 30


def test_create_job_accepts_user_requirements(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())
    captured_request = {}

    def fake_run_build_pipeline(model, request, **kwargs):
        captured_request["request"] = request
        return _fake_pipeline_result(request.output_dir, accepted=True)

    monkeypatch.setattr(api, "run_build_pipeline", fake_run_build_pipeline)
    payload = {
        **_job_payload(),
        "user_requirements": "做一份中文课堂展示，提醒学术诚信风险。",
    }

    response = _client(tmp_path).post("/api/jobs", json=payload)

    assert response.status_code == 202
    generation_request = captured_request["request"].generation_request
    assert generation_request.language == "zh-CN"
    assert generation_request.user_requirements == "做一份中文课堂展示，提醒学术诚信风险。"
    assert generation_request.use_llm_brief is False
    assert generation_request.use_llm_plan is False


def test_long_deck_job_writes_ir_pptx_and_registers_artifacts(tmp_path: Path, monkeypatch) -> None:
    captured: dict = {}
    ppt_master_root = _mock_expected_ppt_master_root(tmp_path)
    _install_fake_long_deck_backend(monkeypatch, captured, ppt_master_root=ppt_master_root)
    client = _client(tmp_path)

    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    artifact_names = {artifact["name"] for artifact in artifacts}

    assert body["status"] == "succeeded"
    assert body["accepted"] is True
    assert body["qa_score"] == 82
    assert body["current_stage"] == "completed"
    assert body["job_type"] == "long_deck"
    assert body["total_batches"] == 15
    assert body["completed_batches"] == 1
    assert body["failed_batches"] == 0
    assert body["current_batch"] == "batch_01"
    assert body["cancel_requested"] is False
    assert body["ppt_master_package"]["generated"] is True
    assert body["ppt_master_package"]["available"] is True
    assert body["ppt_master_package"]["is_expected_repo"] is True
    assert body["ppt_master_package"]["package_mode"] == "normal"
    assert body["ppt_master_package"]["reason"] == "normal_generated"
    assert body["ppt_master_package"]["source_quality_gate_status"] == "passed"
    assert body["ppt_master_package"]["warning"] is None
    assert body["ppt_master_package"]["ppt_master_root"] == str(ppt_master_root.resolve())
    assert body["ppt_master_package"]["missing_paths"] == []
    assert body["ppt_master_package"]["source_artifact_id"]
    assert body["ppt_master_package"]["run_prompt_artifact_id"]
    assert body["ppt_master_package"]["manifest_artifact_id"]
    assert body["ppt_master_package"]["readme_artifact_id"]
    assert captured["render_output"].name == "generated_long_deck.pptx"
    assert {
        "generated_long_deck_plan",
        "generated_long_deck_ir",
        "generated_long_deck_qa",
        "generated_long_deck_quality_gate",
        "generated_long_deck",
        "ppt_master_source",
        "ppt_master_run_prompt",
        "ppt_master_package_manifest",
        "ppt_master_package_README",
        "long_deck_run_report",
        "long_deck_render_report",
        "long_deck_request",
        "batch_01_status",
        "batch_01_deck_ir",
        "batch_01_qa_report",
        "batch_01_attempts",
    } <= artifact_names

    deck_ir_artifact = next(artifact for artifact in artifacts if artifact["name"] == "generated_long_deck_ir")
    ppt_master_artifact = next(artifact for artifact in artifacts if artifact["name"] == "ppt_master_source")
    ppt_master_prompt_artifact = next(artifact for artifact in artifacts if artifact["name"] == "ppt_master_run_prompt")
    ppt_master_manifest_artifact = next(
        artifact for artifact in artifacts if artifact["name"] == "ppt_master_package_manifest"
    )
    ppt_master_readme_artifact = next(
        artifact for artifact in artifacts if artifact["name"] == "ppt_master_package_README"
    )
    assert client.get(deck_ir_artifact["download_url"]).status_code == 200
    assert ppt_master_artifact["kind"] == "md"
    assert client.get(ppt_master_artifact["download_url"]).status_code == 200
    assert ppt_master_prompt_artifact["kind"] == "md"
    assert client.get(ppt_master_prompt_artifact["download_url"]).status_code == 200
    assert ppt_master_manifest_artifact["kind"] == "json"
    assert ppt_master_readme_artifact["kind"] == "md"
    assert client.get(ppt_master_readme_artifact["download_url"]).status_code == 200
    manifest_response = client.get(ppt_master_manifest_artifact["download_url"])
    assert manifest_response.status_code == 200
    manifest = manifest_response.json()
    assert manifest["is_available"] is True
    assert manifest["is_expected_repo"] is True
    assert manifest["package_mode"] == "normal"


def test_long_deck_resume_endpoint_reuses_original_output_dir(tmp_path: Path, monkeypatch) -> None:
    captured: dict = {}
    _install_fake_long_deck_backend(monkeypatch, captured)
    client = _client(tmp_path)

    original_job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    resume_response = client.post(f"/api/long-deck-jobs/{original_job_id}/resume")

    assert resume_response.status_code == 202
    resume_job_id = resume_response.json()["job_id"]
    assert resume_job_id != original_job_id
    assert captured["requests"][-1].resume is True
    assert captured["requests"][-1].output_dir == tmp_path / "jobs" / original_job_id
    body = client.get(f"/api/jobs/{resume_job_id}").json()
    assert body["status"] == "succeeded"
    artifacts = client.get(f"/api/jobs/{resume_job_id}/artifacts").json()["artifacts"]
    assert {artifact["name"] for artifact in artifacts} >= {
        "generated_long_deck_ir",
        "generated_long_deck",
        "ppt_master_source",
        "ppt_master_run_prompt",
        "ppt_master_package_manifest",
        "ppt_master_package_README",
        "long_deck_run_report",
    }


def test_v2_long_deck_resume_reuses_checkpoints_and_original_output_dir(
    tmp_path: Path,
    monkeypatch,
) -> None:
    captured: dict = {}
    _install_fake_v2_long_deck_backend(monkeypatch, captured)
    client = _client(tmp_path)
    payload = {**_long_deck_payload(), "slide_count": 100}

    original_job_id = client.post("/api/long-deck-jobs", json=payload).json()["job_id"]
    response = client.post(f"/api/long-deck-jobs/{original_job_id}/resume")

    assert response.status_code == 202
    resume_job_id = response.json()["job_id"]
    assert resume_job_id != original_job_id
    assert captured["requests"][-1].resume is True
    assert captured["requests"][-1].output_dir == str(tmp_path / "jobs" / original_job_id)
    body = client.get(f"/api/jobs/{resume_job_id}").json()
    assert body["status"] == "succeeded"
    assert body["job_type"] == "long_deck_v2"
    assert body["completed_batches"] == 100


def test_long_deck_job_quality_gate_failure_keeps_ir_artifacts_and_skips_render(
    tmp_path: Path,
    monkeypatch,
) -> None:
    captured: dict = {}
    _install_fake_long_deck_quality_gate_failure(monkeypatch, captured)
    client = _client(tmp_path)

    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    artifact_names = {artifact["name"] for artifact in artifacts}

    assert body["status"] == "failed_quality_gate"
    assert body["accepted"] is False
    assert body["current_stage"] == "failed_quality_gate"
    assert "quality gate failed" in body["error_message"].lower()
    assert body["ppt_master_package"]["generated"] is True
    assert body["ppt_master_package"]["package_mode"] == "recovery"
    assert body["ppt_master_package"]["reason"] == "quality_gate_failed_recovery_generated"
    assert body["ppt_master_package"]["source_quality_gate_status"] == "failed_quality_gate"
    assert body["ppt_master_package"]["source_artifact_id"]
    assert "不会生成旧 renderer PPTX" in body["ppt_master_package"]["message"]
    assert "recovery package" in body["ppt_master_package"]["message"]
    assert "failed the hard quality gate" in body["ppt_master_package"]["warning"]
    assert {
        "generated_long_deck_plan",
        "generated_long_deck_ir",
        "generated_long_deck_qa",
        "generated_long_deck_quality_gate",
        "ppt_master_source",
        "ppt_master_run_prompt",
        "ppt_master_package_manifest",
        "ppt_master_package_README",
        "long_deck_run_report",
        "long_deck_request",
        "batch_01_status",
        "batch_01_deck_ir",
        "batch_01_qa_report",
        "batch_01_attempts",
    } <= artifact_names
    assert "generated_long_deck" not in artifact_names
    assert "long_deck_render_report" not in artifact_names
    source_artifact = next(artifact for artifact in artifacts if artifact["name"] == "ppt_master_source")
    manifest_artifact = next(artifact for artifact in artifacts if artifact["name"] == "ppt_master_package_manifest")
    source_markdown = client.get(source_artifact["download_url"]).text
    lowered_source = source_markdown.lower()
    assert "risk:" not in lowered_source
    assert "risk：" not in lowered_source
    assert "impact:" not in lowered_source
    assert "impact：" not in lowered_source
    assert "mitigation:" not in lowered_source
    assert "mitigation：" not in lowered_source
    assert "判断点 1" not in source_markdown
    assert "Option A" not in source_markdown
    assert "把这一点转化为明确的下一步行动" not in source_markdown
    assert "先列出 Agent 不允许自动执行的动作" not in source_markdown
    assert "bbox" not in source_markdown
    assert "element_id" not in source_markdown
    assert "slide_id" not in source_markdown
    manifest = client.get(manifest_artifact["download_url"]).json()
    assert manifest["package_mode"] == "recovery"
    assert manifest["source_quality_gate_status"] == "failed_quality_gate"
    assert manifest["source_quality_gate_report_path"].endswith("generated_long_deck_quality_gate.json")
    assert manifest["warning"] == (
        "This package was generated from a Deck IR that failed the hard quality gate. "
        "It is intended for PPT Master recovery rendering, not direct renderer output."
    )


def test_cancel_endpoint_marks_running_long_deck_job(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="running", current_stage="generating_batch_01_of_15")

    response = client.post(f"/api/jobs/{job.job_id}/cancel")

    assert response.status_code == 200
    body = response.json()
    assert body["cancel_requested"] is True
    assert body["current_stage"] == "cancel_requested"
    assert store.is_cancel_requested(job.job_id) is True


def test_cancel_endpoint_rejects_short_deck_job(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="short_deck")
    store.update_job(job.job_id, status="running", current_stage="generate_deck")

    response = client.post(f"/api/jobs/{job.job_id}/cancel")

    assert response.status_code == 400
    assert "Only long deck jobs" in response.json()["detail"]


def test_long_deck_job_status_keeps_batch_stage_on_failure(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_long_deck_batch_generation(request, model, *, progress_logger=None, cancel_checker=None):
        if progress_logger is not None:
            progress_logger("Starting batch_01 slides 1-2")
        raise RuntimeError("provider stopped")

    monkeypatch.setattr(api, "run_long_deck_batch_generation", fake_run_long_deck_batch_generation)
    client = _client(tmp_path)

    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]

    assert body["status"] == "failed"
    assert body["current_stage"] == "generating_batch_01_of_15"
    assert "provider stopped" in body["error_message"]
    assert body["ppt_master_package"]["generated"] is False
    assert body["ppt_master_package"]["reason"] == "batch_generation_failed_before_merge"
    assert body["ppt_master_package"]["available"] is None
    assert body["ppt_master_package"]["is_expected_repo"] is None
    assert body["ppt_master_package"]["ppt_master_root"] is None
    assert "Resume the job" in body["ppt_master_package"]["message"]
    assert "ppt_master_source" not in {artifact["name"] for artifact in artifacts}


def test_long_deck_job_timeout_before_merge_keeps_ppt_master_state_unknown(tmp_path: Path, monkeypatch) -> None:
    _install_fake_long_deck_timeout_before_merge(monkeypatch)
    client = _client(tmp_path)

    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]

    assert body["status"] == "failed"
    assert body["current_stage"] == "generating_batch_13_of_15"
    assert "timed out after 3600 seconds" in body["error_message"]
    assert body["ppt_master_package"]["generated"] is False
    assert body["ppt_master_package"]["package_mode"] is None
    assert body["ppt_master_package"]["reason"] == "job_timeout_before_merge"
    assert body["ppt_master_package"]["available"] is None
    assert body["ppt_master_package"]["is_expected_repo"] is None
    assert body["ppt_master_package"]["ppt_master_root"] is None
    assert body["ppt_master_package"]["missing_paths"] == []
    assert "Resume the job" in body["ppt_master_package"]["message"]
    assert "PPT Master package will be generated after merge" in body["ppt_master_package"]["message"]
    assert "ppt_master_source" not in {artifact["name"] for artifact in artifacts}


def test_job_status_can_be_queried(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]

    response = client.get(f"/api/jobs/{job_id}")

    assert response.status_code == 200
    body = response.json()
    assert body["job_id"] == job_id
    assert body["status"] == "succeeded"
    assert body["accepted"] is True
    assert isinstance(body["qa_score"], int)
    assert body["current_stage"] == "complete_job"
    assert body["last_updated_at"] == body["updated_at"]
    assert body["elapsed_seconds"] >= 0
    assert body["ppt_master_package"] is None
    assert body["ppt_master_execution"] is None
    assert body["ppt_master_output"] is None
    assert body["ppt_master_runner"] is None


def test_latest_long_deck_job_can_be_queried(tmp_path: Path, monkeypatch) -> None:
    _install_fake_long_deck_backend(monkeypatch)
    client = _client(tmp_path)
    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]

    response = client.get("/api/jobs/latest?job_type=long_deck")

    assert response.status_code == 200
    body = response.json()
    assert body["job_id"] == job_id
    assert body["job_type"] == "long_deck"


def test_latest_long_deck_job_excludes_short_deck_jobs(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    _install_fake_long_deck_backend(monkeypatch)
    client = _client(tmp_path)
    client.post("/api/jobs", json=_job_payload())
    long_job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]

    response = client.get("/api/jobs/latest?job_type=long_deck")

    assert response.status_code == 200
    body = response.json()
    assert body["job_id"] == long_job_id
    assert body["job_type"] == "long_deck"


def test_old_long_deck_job_without_ppt_master_package_still_returns_status(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="succeeded", accepted=True, current_stage="completed")

    response = client.get(f"/api/jobs/{job.job_id}")

    assert response.status_code == 200
    body = response.json()
    assert body["ppt_master_package"]["generated"] is False
    assert "not been generated" in body["ppt_master_package"]["message"]
    assert body["ppt_master_execution"]["status"] == "not_prepared"
    assert "has not been prepared" in body["ppt_master_execution"]["message"]
    assert body["ppt_master_visual_project"]["status"] == "not_bootstrapped"
    assert "has not been bootstrapped" in body["ppt_master_visual_project"]["message"]
    assert body["ppt_master_output"]["detected"] is False
    assert body["ppt_master_output"]["message"] == "No PPT Master output has been registered for this job."
    assert body["ppt_master_runner"]["status"] == "not_run"
    assert "has not been run" in body["ppt_master_runner"]["message"]


def test_long_deck_job_status_reports_registered_ppt_master_output(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="succeeded", accepted=True, current_stage="completed")
    output_dir = tmp_path / "jobs" / job.job_id / "ppt_master_output"
    _build_mock_ppt_master_output(output_dir, slide_count=30)
    register_ppt_master_output_artifacts(store, job_id=job.job_id, output_dir=output_dir)

    response = client.get(f"/api/jobs/{job.job_id}")

    assert response.status_code == 200
    body = response.json()
    assert body["ppt_master_output"]["detected"] is True
    assert body["ppt_master_output"]["slide_count"] == 30
    assert body["ppt_master_output"]["generation_status"] == "succeeded"
    assert body["ppt_master_output"]["pptx_artifact_id"]
    assert body["ppt_master_output"]["notes_artifact_id"]
    assert body["ppt_master_output"]["manifest_artifact_id"]
    assert body["ppt_master_output"]["output_dir"] == str(output_dir.resolve())


def test_long_deck_job_status_auto_registers_existing_ppt_master_output(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="succeeded", accepted=True, current_stage="completed")
    output_dir = tmp_path / "jobs" / job.job_id / "ppt_master_output"
    _build_mock_ppt_master_output(output_dir, slide_count=30)

    response = client.get(f"/api/jobs/{job.job_id}")
    artifacts = client.get(f"/api/jobs/{job.job_id}/artifacts").json()["artifacts"]

    assert response.status_code == 200
    body = response.json()
    assert body["ppt_master_output"]["detected"] is True
    assert {artifact["name"] for artifact in artifacts} >= {
        PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
        PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
        PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    }


def test_prepare_ppt_master_execution_endpoint_returns_waiting_status(
    tmp_path: Path,
    monkeypatch,
) -> None:
    root = _mock_expected_ppt_master_root(tmp_path)
    monkeypatch.setenv("PPT_MASTER_DIR", str(root))
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")
    _build_mock_ppt_master_package(tmp_path / "jobs" / job.job_id)

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/prepare-ppt-master-execution")
    status_response = client.get(f"/api/jobs/{job.job_id}")
    artifacts = client.get(f"/api/jobs/{job.job_id}/artifacts").json()["artifacts"]

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "waiting_for_external_ppt_master_run"
    assert body["plan_artifact_id"]
    assert body["expected_pptx_path"].endswith("generated_by_ppt_master.pptx")
    assert any("run_prompt.md" in step for step in body["suggested_steps"])
    assert any("register_ppt_master_output.py" in step for step in body["suggested_steps"])
    status_body = status_response.json()
    assert status_body["ppt_master_execution"]["status"] == "waiting_for_external_ppt_master_run"
    assert status_body["ppt_master_execution"]["plan_artifact_id"]
    assert PPT_MASTER_EXECUTION_PLAN_ARTIFACT in {artifact["name"] for artifact in artifacts}


def test_prepare_ppt_master_execution_endpoint_detects_output_and_registers_it(
    tmp_path: Path,
    monkeypatch,
) -> None:
    root = _mock_expected_ppt_master_root(tmp_path)
    monkeypatch.setenv("PPT_MASTER_DIR", str(root))
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")
    job_dir = tmp_path / "jobs" / job.job_id
    _build_mock_ppt_master_package(job_dir)
    _build_mock_ppt_master_output(job_dir / "ppt_master_output", slide_count=30)

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/prepare-ppt-master-execution")
    status_response = client.get(f"/api/jobs/{job.job_id}")
    artifacts = client.get(f"/api/jobs/{job.job_id}/artifacts").json()["artifacts"]

    assert response.status_code == 200
    assert response.json()["status"] == "output_detected"
    status_body = status_response.json()
    assert status_body["ppt_master_execution"]["status"] == "output_detected"
    assert status_body["ppt_master_output"]["detected"] is True
    assert status_body["ppt_master_output"]["slide_count"] == 30
    assert {artifact["name"] for artifact in artifacts} >= {
        PPT_MASTER_EXECUTION_PLAN_ARTIFACT,
        PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
        PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
        PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    }


def test_prepare_ppt_master_execution_endpoint_reports_missing_package(tmp_path: Path) -> None:
    root = _mock_expected_ppt_master_root(tmp_path)
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/prepare-ppt-master-execution")

    assert response.status_code == 200
    assert response.json()["status"] == "missing_package"
    assert (tmp_path / "jobs" / job.job_id / "ppt_master_execution_plan.json").exists()
    assert root.exists()


def test_prepare_ppt_master_execution_endpoint_reports_unavailable_ppt_master(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("PPT_MASTER_DIR", str(tmp_path / "missing-ppt-master"))
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")
    _build_mock_ppt_master_package(tmp_path / "jobs" / job.job_id)

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/prepare-ppt-master-execution")

    assert response.status_code == 200
    assert response.json()["status"] == "ppt_master_unavailable"


def test_prepare_ppt_master_execution_endpoint_rejects_short_deck_job(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="short_deck")

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/prepare-ppt-master-execution")

    assert response.status_code == 400
    assert "Only long deck jobs" in response.json()["detail"]


def test_bootstrap_ppt_master_project_endpoint_creates_scaffold_and_registers_artifacts(
    tmp_path: Path,
    monkeypatch,
) -> None:
    root = _mock_expected_ppt_master_root(tmp_path)
    monkeypatch.setenv("PPT_MASTER_DIR", str(root))
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")
    job_dir = tmp_path / "jobs" / job.job_id
    _build_mock_ppt_master_package(job_dir)

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/bootstrap-ppt-master-project")
    status_response = client.get(f"/api/jobs/{job.job_id}")
    artifacts = client.get(f"/api/jobs/{job.job_id}/artifacts").json()["artifacts"]

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "created"
    assert body["manifest_artifact_id"]
    assert body["instructions_artifact_id"]
    assert body["project_dir"].endswith("ppt_master_visual_project")
    assert body["project_source_path"].endswith("inputs/source.md")
    assert body["project_prompt_path"].endswith("inputs/run_prompt.md")
    assert body["expected_svg_output_dir"].endswith("svg_output")
    assert body["expected_svg_final_dir"].endswith("svg_final")
    assert body["expected_pptx_path"].endswith("generated_by_ppt_master.pptx")
    assert any("PROJECT_INSTRUCTIONS.md" in step for step in body["next_steps"])
    status_body = status_response.json()
    assert status_body["ppt_master_visual_project"]["status"] == "created"
    assert status_body["ppt_master_visual_project"]["instructions_artifact_id"]
    assert {artifact["name"] for artifact in artifacts} >= {
        PPT_MASTER_VISUAL_PROJECT_MANIFEST_ARTIFACT,
        PPT_MASTER_PROJECT_INSTRUCTIONS_ARTIFACT,
    }


def test_bootstrap_ppt_master_project_endpoint_rejects_short_deck_job(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="short_deck")

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/bootstrap-ppt-master-project")

    assert response.status_code == 400
    assert "Only long deck jobs" in response.json()["detail"]


def test_run_ppt_master_local_export_endpoint_reports_external_generation_needed(
    tmp_path: Path,
    monkeypatch,
) -> None:
    root = _mock_expected_ppt_master_root(tmp_path)
    monkeypatch.setenv("PPT_MASTER_DIR", str(root))
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="long_deck")
    store.update_job(job.job_id, status="failed_quality_gate", current_stage="failed_quality_gate")
    _build_mock_ppt_master_package(tmp_path / "jobs" / job.job_id)

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/run-ppt-master-local-export")
    status_response = client.get(f"/api/jobs/{job.job_id}")
    artifacts = client.get(f"/api/jobs/{job.job_id}/artifacts").json()["artifacts"]

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "requires_external_ai_generation"
    assert body["requires_external_ai_generation"] is True
    assert body["result_artifact_id"]
    assert "visual project" in body["message"]
    status_body = status_response.json()
    assert status_body["ppt_master_runner"]["status"] == "requires_external_ai_generation"
    assert status_body["ppt_master_runner"]["result_artifact_id"]
    assert PPT_MASTER_RUNNER_RESULT_ARTIFACT in {artifact["name"] for artifact in artifacts}


def test_run_ppt_master_local_export_endpoint_rejects_short_deck_job(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job(job_type="short_deck")

    response = client.post(f"/api/long-deck-jobs/{job.job_id}/run-ppt-master-local-export")

    assert response.status_code == 400
    assert "Only long deck jobs" in response.json()["detail"]


def test_job_qa_gate_failure_is_completed_with_artifacts_not_runtime_failed(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch, accepted=False)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]

    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]

    assert body["status"] == "succeeded"
    assert body["accepted"] is False
    assert "QA score gate" in body["error_message"]
    assert {artifact["name"] for artifact in artifacts} >= {
        "generated_deck_ir",
        "patchable_elements",
        "generated_qa_report",
        "generated_attempts",
        "generated_deck",
    }


def test_job_failure_marks_failed_instead_of_running(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_build_pipeline(model, request, **kwargs):
        raise RuntimeError("provider failed with api_key=sk-testsecret123456789")

    monkeypatch.setattr(api, "run_build_pipeline", fake_run_build_pipeline)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]

    body = client.get(f"/api/jobs/{job_id}").json()

    assert body["status"] == "failed"
    assert body["accepted"] is False
    assert "provider failed" in body["error_message"]
    assert "sk-testsecret" not in body["error_message"]


def test_job_timeout_marks_failed_with_clear_error(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_build_pipeline(model, request, **kwargs):
        raise TimeoutError("Job timed out while running stage 'generate_deck' after 600 seconds.")

    monkeypatch.setattr(api, "run_build_pipeline", fake_run_build_pipeline)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]

    body = client.get(f"/api/jobs/{job_id}").json()

    assert body["status"] == "failed"
    assert "timed out" in body["error_message"]
    assert "generate_deck" in body["error_message"]


def test_long_deck_job_timeout_seconds_can_be_configured(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("LONG_DECK_JOB_TIMEOUT_SECONDS", "7200")
    _install_fake_long_deck_timeout_before_merge(monkeypatch)
    client = _client(tmp_path)

    job_id = client.post("/api/long-deck-jobs", json=_long_deck_payload()).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()

    assert api._long_deck_job_timeout_seconds() == 7200
    assert "timed out after 7200 seconds" in body["error_message"]
    assert body["ppt_master_package"]["reason"] == "job_timeout_before_merge"


def test_job_status_expires_stale_running_job(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setattr(api, "JOB_TIMEOUT_SECONDS", -1)
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job()
    store.update_job(job.job_id, status="running", current_stage="generate_deck")

    body = client.get(f"/api/jobs/{job.job_id}").json()

    assert body["status"] == "failed"
    assert "timed out" in body["error_message"]
    assert "generate_deck" in body["error_message"]


def test_job_with_missing_patch_path_fails_fast(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())
    client = _client(tmp_path)
    payload = {
        **_job_payload(),
        "patch_path": str(tmp_path / "missing_patch.json"),
    }

    job_id = client.post("/api/jobs", json=payload).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()

    assert body["status"] == "failed"
    assert "Patch file not found" in body["error_message"]
    assert body["current_stage"] == "apply_patch"


def test_artifacts_can_be_listed(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]

    response = client.get(f"/api/jobs/{job_id}/artifacts")

    assert response.status_code == 200
    artifacts = response.json()["artifacts"]
    assert {artifact["name"] for artifact in artifacts} == {
        "generated_deck_ir",
        "patchable_elements",
        "generated_qa_report",
        "generated_attempts",
        "generated_deck",
    }
    assert all(artifact["download_url"].startswith("/api/artifacts/") for artifact in artifacts)


def test_long_deck_svg_previews_are_served_from_job_output(tmp_path: Path) -> None:
    client = _client(tmp_path)
    job = client.app.state.job_store.create_job(job_type="long_deck")
    svg_dir = (
        client.app.state.jobs_root
        / job.job_id
        / "ppt_master_output"
        / "visual_project"
        / "svg_final"
    )
    svg_dir.mkdir(parents=True)
    for number in range(1, 4):
        (svg_dir / f"slide_{number:02d}.svg").write_text(
            f'<svg xmlns="http://www.w3.org/2000/svg"><text>Slide {number}</text></svg>',
            encoding="utf-8",
        )

    response = client.get(f"/api/jobs/{job.job_id}/preview-slides/2")
    manifest = client.get(f"/api/jobs/{job.job_id}/preview-slides")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith("image/svg+xml")
    assert b"Slide 2" in response.content
    assert manifest.json()["available_slide_numbers"] == [1, 2, 3]
    assert manifest.json()["preview_kind"] == "ppt_master_svg"
    assert client.get(f"/api/jobs/{job.job_id}/preview-slides/4").status_code == 404


def test_v2_checkpoint_previews_are_available_before_deck_completion(tmp_path: Path) -> None:
    client = _client(tmp_path)
    job = client.app.state.job_store.create_job(job_type="long_deck_v2")
    client.app.state.job_store.update_long_deck_progress(job.job_id, total_batches=100)
    checkpoint_root = client.app.state.jobs_root / job.job_id / "checkpoints"
    pages_dir = checkpoint_root / "pages"
    pages_dir.mkdir(parents=True)
    (checkpoint_root / "theme.json").write_text(
        BUILTIN_THEMES["aurora"].model_dump_json(indent=2),
        encoding="utf-8",
    )
    (checkpoint_root / "skeleton.json").write_text(
        json.dumps({"deck_title": "渐进预览测试", "subtitle": "Checkpoint", "language": "zh-CN"}),
        encoding="utf-8",
    )
    page = PageDesign(
        page_number=5,
        title="第五页已生成",
        elements=[
            TextItem(
                id="title",
                frame=Frame(x=80, y=90, w=800, h=100),
                text="第五页已生成",
                role="title",
            )
        ],
    )
    (pages_dir / "page_005.json").write_text(
        json.dumps({"page": page.model_dump(mode="json"), "qa": {}, "outcome": "generated"}),
        encoding="utf-8",
    )

    manifest = client.get(f"/api/jobs/{job.job_id}/preview-slides")
    preview = client.get(f"/api/jobs/{job.job_id}/preview-slides/5")

    assert manifest.status_code == 200
    assert manifest.json()["available_slide_numbers"] == [5]
    assert manifest.json()["total_requested"] == 100
    assert manifest.json()["preview_kind"] == "v2_html"
    assert preview.status_code == 200
    assert preview.headers["content-type"].startswith("text/html")
    assert "第五页已生成" in preview.text
    assert client.get(f"/api/jobs/{job.job_id}/preview-slides/1").status_code == 404


def test_v2_final_design_previews_include_all_pages(tmp_path: Path) -> None:
    client = _client(tmp_path)
    job = client.app.state.job_store.create_job(job_type="long_deck_v2")
    job_dir = client.app.state.jobs_root / job.job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    pages = [
        PageDesign(
            page_number=number,
            title=f"第 {number} 页",
            elements=[
                TextItem(
                    id=f"title-{number}",
                    frame=Frame(x=80, y=90, w=800, h=100),
                    text=f"第 {number} 页",
                    role="title",
                )
            ],
        )
        for number in range(1, 4)
    ]
    deck = DeckDesign(deck_title="完整预览", theme=BUILTIN_THEMES["aurora"], pages=pages)
    (job_dir / "generated_long_deck_v2_design.json").write_text(
        deck.model_dump_json(indent=2),
        encoding="utf-8",
    )

    manifest = client.get(f"/api/jobs/{job.job_id}/preview-slides")

    assert manifest.json()["available_slide_numbers"] == [1, 2, 3]
    assert manifest.json()["highlight_slide_numbers"] == [1, 2, 3]
    assert client.get(f"/api/jobs/{job.job_id}/preview-slides/3").status_code == 200


def test_v2_preview_manifest_selects_visually_rich_highlight_pages(tmp_path: Path) -> None:
    client = _client(tmp_path)
    job = client.app.state.job_store.create_job(job_type="long_deck_v2")
    job_dir = client.app.state.jobs_root / job.job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    def title(number: int) -> TextItem:
        return TextItem(
            id=f"title-{number}",
            frame=Frame(x=80, y=70, w=800, h=80),
            text=f"第 {number} 页",
            role="title",
        )

    pages = [
        PageDesign(page_number=1, role="cover", title="封面", elements=[title(1)]),
        PageDesign(page_number=2, role="content", title="纯文字", elements=[title(2)]),
        PageDesign(
            page_number=3,
            role="stats",
            title="指标高光",
            elements=[
                title(3),
                ChartItem(
                    id="chart-3",
                    frame=Frame(x=100, y=190, w=650, h=380),
                    chart="bar",
                    categories=["A", "B", "C"],
                    series=[ChartSeries(name="趋势", values=[3, 7, 5])],
                ),
                IconItem(id="icon-3", frame=Frame(x=850, y=220, w=90, h=90), name="chart"),
            ],
        ),
        PageDesign(
            page_number=4,
            role="timeline",
            title="流程高光",
            elements=[
                title(4),
                ShapeItem(id="step-1", frame=Frame(x=120, y=240, w=240, h=150)),
                ShapeItem(id="step-2", frame=Frame(x=440, y=240, w=240, h=150)),
                IconItem(id="icon-4", frame=Frame(x=760, y=260, w=80, h=80), name="arrow-right"),
            ],
        ),
        PageDesign(
            page_number=5,
            role="comparison",
            title="对比高光",
            elements=[
                title(5),
                ShapeItem(id="left", frame=Frame(x=100, y=210, w=470, h=330)),
                ShapeItem(id="right", frame=Frame(x=710, y=210, w=470, h=330)),
            ],
        ),
        PageDesign(page_number=6, role="closing", title="结束页", elements=[title(6)]),
    ]
    deck = DeckDesign(deck_title="高光页测试", theme=BUILTIN_THEMES["aurora"], pages=pages)
    (job_dir / "generated_long_deck_v2_design.json").write_text(
        deck.model_dump_json(indent=2),
        encoding="utf-8",
    )

    manifest = client.get(f"/api/jobs/{job.job_id}/preview-slides").json()

    assert manifest["available_slide_numbers"] == [1, 2, 3, 4, 5, 6]
    assert manifest["highlight_slide_numbers"] == [3, 4, 5]


def test_artifacts_include_patch_outputs_when_patch_succeeds(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setattr(api, "_create_chat_model", lambda: object())

    def fake_run_build_pipeline(model, request, **kwargs):
        return _fake_pipeline_result_with_patch(request.output_dir, accepted=True)

    monkeypatch.setattr(api, "run_build_pipeline", fake_run_build_pipeline)
    client = _client(tmp_path)
    payload = {
        **_job_payload(),
        "patch_path": str(EXAMPLES_DIR / "sample_patch.json"),
    }

    job_id = client.post("/api/jobs", json=payload).json()["job_id"]
    body = client.get(f"/api/jobs/{job_id}").json()
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]

    assert body["status"] == "succeeded"
    assert body["accepted"] is True
    assert {artifact["name"] for artifact in artifacts} >= {
        "generated_deck_ir",
        "patchable_elements",
        "patch_report",
        "patched_deck",
        "patched_deck_ir",
    }


def test_registered_artifact_can_be_downloaded(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]
    artifacts = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    deck_artifact = next(artifact for artifact in artifacts if artifact["name"] == "generated_deck_ir")

    response = client.get(deck_artifact["download_url"])

    assert response.status_code == 200
    assert json.loads(response.content)["deck_id"] == "sample_clean_business_deck"


def test_artifact_download_rejects_unregistered_file(tmp_path: Path, monkeypatch) -> None:
    _install_fake_backend(monkeypatch)
    client = _client(tmp_path)
    job_id = client.post("/api/jobs", json=_job_payload()).json()["job_id"]
    unregistered = tmp_path / "jobs" / job_id / "unregistered.json"
    unregistered.write_text("{}", encoding="utf-8")

    response = client.get("/api/artifacts/unregistered.json")

    assert response.status_code == 404


def test_artifact_download_rejects_registered_path_outside_jobs_root(tmp_path: Path) -> None:
    store = JobStore(tmp_path / "jobs.sqlite3")
    app = api.create_app(data_dir=tmp_path, store=store)
    client = TestClient(app)
    job = store.create_job()
    outside = tmp_path / "outside.json"
    outside.write_text("{}", encoding="utf-8")
    artifact = store.add_artifact(job.job_id, name="outside", kind="json", path=outside)

    response = client.get(f"/api/artifacts/{artifact.artifact_id}")

    assert response.status_code == 404


def test_missing_job_returns_404(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/api/jobs/missing-job")

    assert response.status_code == 404


def test_missing_artifact_returns_404(tmp_path: Path) -> None:
    response = _client(tmp_path).get("/api/artifacts/missing-artifact")

    assert response.status_code == 404
