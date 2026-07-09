"""Tests for v2 QA rules, deterministic fixes, and the PPTX renderer."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ppt_agent.v2.design import get_builtin_theme
from ppt_agent.v2.fallback import design_fallback_page
from ppt_agent.v2.ir import (
    DeckDesign,
    Frame,
    PageDesign,
    ShapeItem,
    TextItem,
)
from ppt_agent.v2.metrics import fit_font_size
from ppt_agent.v2.planning import PageBrief
from ppt_agent.v2.qa import review_page
from ppt_agent.v2.render import render_deck


THEME = get_builtin_theme("aurora")


def _page(elements, page_number: int = 1, **kwargs) -> PageDesign:
    return PageDesign(page_number=page_number, role="content", elements=elements, **kwargs)


class TestMetrics:
    def test_fit_font_size_shrinks_long_text(self) -> None:
        short = fit_font_size("短", role="body", frame_width_units=300, frame_height_units=60)
        long = fit_font_size("很长的中文内容" * 20, role="body", frame_width_units=300, frame_height_units=60)
        assert long < short
        assert long >= 10  # role minimum


class TestQA:
    def test_contrast_autofix_recolors_text_on_dark_card(self) -> None:
        page = _page(
            [
                ShapeItem(id="card", frame=Frame(x=100, y=200, w=400, h=200), fill="primary"),
                TextItem(
                    id="label",
                    frame=Frame(x=120, y=220, w=360, h=60),
                    text="dark on dark",
                    color="text",
                ),
            ]
        )
        fixed, result = review_page(page, THEME)
        assert fixed.elements[1].color == "on_primary"
        assert any("recolored" in fix for fix in result.auto_fixes)

    def test_extreme_overflow_is_an_error(self) -> None:
        page = _page(
            [
                ShapeItem(id="card", frame=Frame(x=100, y=200, w=400, h=200)),
                TextItem(
                    id="wall",
                    frame=Frame(x=120, y=220, w=200, h=40),
                    text="这是一段远远超过框体容量的文本内容。" * 15,
                ),
            ]
        )
        _, result = review_page(page, THEME)
        assert any(issue.code == "text_overflow" for issue in result.errors)

    def test_overlapping_text_frames_are_an_error(self) -> None:
        page = _page(
            [
                TextItem(id="a", frame=Frame(x=100, y=200, w=300, h=80), text="one"),
                TextItem(id="b", frame=Frame(x=150, y=220, w=300, h=80), text="two"),
                TextItem(id="c", frame=Frame(x=700, y=500, w=200, h=60), text="three"),
            ]
        )
        _, result = review_page(page, THEME)
        assert any(issue.code == "text_overlap" for issue in result.errors)

    def test_chrome_zone_intrusion_is_shifted(self) -> None:
        page = _page(
            [
                TextItem(id="high", frame=Frame(x=100, y=2, w=300, h=40), text="too high"),
                TextItem(id="fine", frame=Frame(x=100, y=300, w=300, h=40), text="ok"),
                TextItem(id="third", frame=Frame(x=600, y=400, w=300, h=40), text="ok2"),
            ]
        )
        fixed, result = review_page(page, THEME)
        assert fixed.elements[0].frame.y >= 56
        assert any("safe zone" in fix for fix in result.auto_fixes)

    def test_fallback_pages_pass_qa(self) -> None:
        brief = PageBrief(
            title="一个正常长度的页面标题",
            summary="页面摘要内容。",
            points=["要点一的内容", "要点二的内容", "要点三的内容", "要点四的内容"],
        )
        for hint in ("cards", "two_column", "stats", "timeline", "chart", "table", "quote", "list"):
            page = design_fallback_page(
                brief.model_copy(update={"layout_hint": hint}),
                page_number=5,
                section_title="章节",
                language="zh-CN",
            )
            _, result = review_page(page, THEME)
            assert result.passed, f"{hint}: {[i.message for i in result.errors]}"


class TestRenderer:
    def test_renders_editable_pptx_with_notes(self, tmp_path: Path) -> None:
        pages = [
            design_fallback_page(
                PageBrief(
                    title=f"页面 {number}",
                    summary="备注内容",
                    points=["要点一", "要点二", "要点三"],
                    layout_hint="auto",
                ),
                page_number=number,
                section_title="章节",
                language="zh-CN",
            )
            for number in range(1, 6)
        ]
        deck = DeckDesign(deck_title="渲染测试", theme=THEME, pages=pages)
        output = render_deck(deck, tmp_path / "out.pptx")

        presentation = Presentation(str(output))
        assert len(presentation.slides) == 5
        first = presentation.slides[0]
        texts = [
            shape.text_frame.text
            for shape in first.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        assert any("页面 1" in text for text in texts), "title must be a native text box"
        assert first.notes_slide.notes_text_frame.text.strip()

    def test_gradient_cover_and_chart_render(self, tmp_path: Path) -> None:
        from ppt_agent.v2.anchors import build_cover_page
        from ppt_agent.v2.ir import ChartItem, ChartSeries

        cover = build_cover_page(
            page_number=1,
            deck_title="封面",
            subtitle="副标题",
            language="zh-CN",
            theme=THEME,
        )
        chart_page = _page(
            [
                ChartItem(
                    id="c",
                    frame=Frame(x=100, y=160, w=800, h=400),
                    chart="pie",
                    categories=["甲", "乙", "丙"],
                    series=[ChartSeries(name="占比", values=[3, 2, 1])],
                )
            ],
            page_number=2,
        )
        deck = DeckDesign(deck_title="混合", theme=THEME, pages=[cover, chart_page])
        output = render_deck(deck, tmp_path / "mixed.pptx")
        presentation = Presentation(str(output))
        assert any(shape.has_chart for shape in presentation.slides[1].shapes)
