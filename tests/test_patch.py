from pathlib import Path

from pptx import Presentation

from ppt_agent.load import load_deck, load_patch, load_theme
from ppt_agent.models import ShapeElement, TextElement
from ppt_agent.patch import SlidePatch, apply_patch
from ppt_agent.renderer import render_deck_to_pptx


EXAMPLES_DIR = Path(__file__).resolve().parents[1] / "examples"


def _sample_deck():
    return load_deck(EXAMPLES_DIR / "sample_slide_ir.json")


def _element(deck, slide_id: str, element_id: str):
    slide = next(slide for slide in deck.slides if slide.slide_id == slide_id)
    return next(element for element in slide.elements if element.element_id == element_id)


def test_sample_patch_loads_from_json() -> None:
    patch = load_patch(EXAMPLES_DIR / "sample_patch.json")

    assert patch.patch_id == "sample_patch_001"
    assert len(patch.operations) == 3
    assert patch.operations[0].op == "update_text"


def test_apply_patch_returns_new_deck_without_mutating_original() -> None:
    deck = _sample_deck()
    patch = load_patch(EXAMPLES_DIR / "sample_patch.json")
    original_title = _element(deck, "slide_001", "s1_title").text

    result = apply_patch(deck, patch)

    assert result.deck is not deck
    assert result.applied_count == 3
    assert result.issues == []
    assert _element(deck, "slide_001", "s1_title").text == original_title
    assert _element(result.deck, "slide_001", "s1_title").text == "Updated Q3 Operating Review"


def test_update_text_modifies_text_element() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "update_text",
                    "slide_id": "slide_001",
                    "element_id": "s1_title",
                    "text": "New title",
                }
            ]
        }
    )

    result = apply_patch(deck, patch)
    element = _element(result.deck, "slide_001", "s1_title")

    assert result.applied_count == 1
    assert isinstance(element, TextElement)
    assert element.text == "New title"


def test_move_element_modifies_bbox_position() -> None:
    deck = _sample_deck()
    original = _element(deck, "slide_001", "s1_accent_bar")
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "move_element",
                    "slide_id": "slide_001",
                    "element_id": "s1_accent_bar",
                    "dx": 0.5,
                    "dy": 0.25,
                }
            ]
        }
    )

    result = apply_patch(deck, patch)
    moved = _element(result.deck, "slide_001", "s1_accent_bar")

    assert result.applied_count == 1
    assert moved.bbox.x == original.bbox.x + 0.5
    assert moved.bbox.y == original.bbox.y + 0.25


def test_update_shape_style_modifies_shape_element_style() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "update_shape_style",
                    "slide_id": "slide_002",
                    "element_id": "s2_side_panel",
                    "fill_color": "#ECFDF5",
                    "stroke_color": "#10B981",
                    "stroke_width_pt": 2.0,
                }
            ]
        }
    )

    result = apply_patch(deck, patch)
    element = _element(result.deck, "slide_002", "s2_side_panel")

    assert result.applied_count == 1
    assert isinstance(element, ShapeElement)
    assert element.style is not None
    assert element.style.fill_color == "#ECFDF5"
    assert element.style.stroke_color == "#10B981"
    assert element.style.stroke_width_pt == 2.0


def test_missing_slide_id_records_issue() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "update_text",
                    "slide_id": "missing_slide",
                    "element_id": "s1_title",
                    "text": "No-op",
                }
            ]
        }
    )

    result = apply_patch(deck, patch)

    assert result.applied_count == 0
    assert result.issues[0].code == "SLIDE_NOT_FOUND"
    assert "missing_slide" in result.issues[0].message


def test_missing_element_id_records_issue() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "update_text",
                    "slide_id": "slide_001",
                    "element_id": "missing_element",
                    "text": "No-op",
                }
            ]
        }
    )

    result = apply_patch(deck, patch)

    assert result.applied_count == 0
    assert result.issues[0].code == "ELEMENT_NOT_FOUND"
    assert "missing_element" in result.issues[0].message


def test_update_text_on_shape_element_records_type_issue() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "update_text",
                    "slide_id": "slide_001",
                    "element_id": "s1_accent_bar",
                    "text": "Wrong target",
                }
            ]
        }
    )

    result = apply_patch(deck, patch)

    assert result.applied_count == 0
    assert result.issues[0].code == "ELEMENT_TYPE_MISMATCH"
    assert "requires a TextElement" in result.issues[0].message


def test_patch_that_moves_bbox_out_of_bounds_records_validation_issue() -> None:
    deck = _sample_deck()
    patch = SlidePatch.model_validate(
        {
            "operations": [
                {
                    "op": "move_element",
                    "slide_id": "slide_001",
                    "element_id": "s1_title",
                    "dx": 100.0,
                    "dy": 0.0,
                }
            ]
        }
    )

    result = apply_patch(deck, patch)

    assert result.applied_count == 0
    assert result.issues[0].code == "DECK_VALIDATION_FAILED"
    assert "exceeds canvas width" in result.issues[0].message
    assert _element(result.deck, "slide_001", "s1_title").bbox.x == _element(deck, "slide_001", "s1_title").bbox.x


def test_patched_deck_can_render_to_pptx(tmp_path: Path) -> None:
    deck = _sample_deck()
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    patch = load_patch(EXAMPLES_DIR / "sample_patch.json")
    result = apply_patch(deck, patch)

    output_path = render_deck_to_pptx(result.deck, theme, tmp_path / "patched_deck.pptx", assets_dir=EXAMPLES_DIR)

    assert output_path.exists()
    presentation = Presentation(output_path)
    assert len(presentation.slides) == len(result.deck.slides)
