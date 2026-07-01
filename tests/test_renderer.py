import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ppt_agent.layouts import TEMPLATE_LAYOUTS
from ppt_agent.load import load_deck, load_theme
from ppt_agent.models import Deck
from ppt_agent.long_deck_render import render_long_deck_ir_to_pptx, sanitize_deck_ir_for_render
from ppt_agent.renderer import VISUAL_VARIANT_COUNTS, _visual_variant_for_slide, render_deck_to_pptx


EXAMPLES_DIR = Path(__file__).resolve().parents[1] / "examples"
INTERNAL_SURFACE_TERMS = {"TEMPLATE-GUIDED", "Editable PPTX", "COLUMN", "CARD"}


def test_render_deck_to_pptx_generates_editable_powerpoint(tmp_path: Path) -> None:
    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    output_path = tmp_path / "sample_deck.pptx"

    rendered_path = render_deck_to_pptx(deck, theme, output_path, assets_dir=EXAMPLES_DIR)

    assert rendered_path == output_path
    assert output_path.exists()
    assert output_path.stat().st_size > 0

    presentation = Presentation(output_path)
    assert len(presentation.slides) == len(deck.slides)
    assert presentation.slide_width == Inches(deck.canvas_width_in)
    assert presentation.slide_height == Inches(deck.canvas_height_in)

    first_slide = presentation.slides[0]
    editable_text_shapes = [
        shape
        for shape in first_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]

    assert editable_text_shapes
    assert any("Q3 Operating Review" in shape.text for shape in editable_text_shapes)


def test_render_deck_to_pptx_creates_missing_image_placeholder(tmp_path: Path) -> None:
    deck = load_deck(EXAMPLES_DIR / "sample_slide_ir.json")
    theme = load_theme(EXAMPLES_DIR / "theme.json")

    output_path = render_deck_to_pptx(
        deck,
        theme,
        tmp_path / "sample_deck.pptx",
        assets_dir=EXAMPLES_DIR,
    )

    presentation = Presentation(output_path)
    third_slide = presentation.slides[2]

    assert any(
        shape.has_text_frame and "Placeholder" in shape.text
        for shape in third_slide.shapes
    )


def _template_slide_payload(layout: str, index: int) -> dict:
    body_texts = {
        "title_slide": ["A concise subtitle for the title slide"],
        "section_divider": ["This section frames the next part of the story."],
        "two_column": ["Left column content\n- Point A\n- Point B", "Right column content\n- Point C\n- Point D"],
        "three_column": ["First pillar", "Second pillar", "Third pillar"],
        "four_cards": [
            "Discover\nMap user needs",
            "Prioritize\nChoose high-value work",
            "Prototype\nBuild the first flow",
            "Measure\nTrack quality signals",
        ],
        "metric_cards": ["Revenue\n$4.2M", "Growth\n18%", "Retention\n91%"],
        "closing_slide": ["Thank you. Questions and next steps."],
        "comparison_matrix": [
            "Normal AI\nResponds to direct prompts\nNeeds manual handoff",
            "AI Agent\nPlans multi-step work\nUses tools with guardrails",
            "Choose Agent when workflow ownership matters",
        ],
        "process_flow": [
            "Discover\nMap user needs",
            "Design\nDefine workflow",
            "Build\nCreate prototype",
            "Measure\nTrack quality",
        ],
        "risk_matrix": [
            "Hallucination\nHigh user trust impact\nUse retrieval checks",
            "Data leakage\nCompliance impact\nLimit sensitive inputs",
            "Over-automation\nOperational impact\nKeep human review",
        ],
        "key_takeaway": [
            "AI agents need product judgment\nStart from owned workflows, not model demos",
            "Pick one high-value user journey",
            "Add evals before scaling",
            "Keep human checkpoints",
        ],
    }[layout]

    elements = [
        {
            "element_id": f"s{index}_title",
            "type": "text",
            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
            "text": f"{layout} title",
        }
    ]
    for body_index, text in enumerate(body_texts, start=1):
        elements.append(
            {
                "element_id": f"s{index}_body_{body_index}",
                "type": "text",
                "bbox": {"x": 0.8 + body_index, "y": 1.5, "width": 3.0, "height": 1.0},
                "text": text,
            }
        )

    return {
        "slide_id": f"template_{index:03d}",
        "title": f"{layout} title",
        "layout": layout,
        "elements": elements,
    }


def test_template_layouts_render_to_editable_powerpoint(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "template_layout_demo",
            "title": "Template Layout Demo",
            "theme_name": "clean_business",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                _template_slide_payload(layout, index)
                for index, layout in enumerate(TEMPLATE_LAYOUTS, start=1)
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "template_layouts.pptx")

    presentation = Presentation(output_path)
    assert len(presentation.slides) == len(TEMPLATE_LAYOUTS)

    for layout, rendered_slide in zip(TEMPLATE_LAYOUTS, presentation.slides):
        editable_text_shapes = [
            shape
            for shape in rendered_slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        assert editable_text_shapes
        assert any(f"{layout} title" in shape.text for shape in editable_text_shapes)


def _shape_with_text(slide, text: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip() == text:
            return shape
    raise AssertionError(f"Could not find editable text shape: {text}")


def _visible_texts(presentation: Presentation) -> list[str]:
    return [
        shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]


def _large_blank_rect_signature(slide) -> tuple[tuple[int, int, int, int], ...]:
    return tuple(
        sorted(
            (
                int(shape.left),
                int(shape.top),
                int(shape.width),
                int(shape.height),
            )
            for shape in slide.shapes
            if not getattr(shape, "text", "").strip()
            and shape.shape_type != MSO_SHAPE_TYPE.LINE
            and shape.width >= Inches(1.0)
            and shape.height >= Inches(0.8)
        )
    )


def _overlap_area(shape_a, shape_b) -> int:
    left = max(shape_a.left, shape_b.left)
    right = min(shape_a.left + shape_a.width, shape_b.left + shape_b.width)
    top = max(shape_a.top, shape_b.top)
    bottom = min(shape_a.top + shape_a.height, shape_b.top + shape_b.height)
    if right <= left or bottom <= top:
        return 0
    return int((right - left) * (bottom - top))


def test_title_slide_long_title_keeps_subtitle_below_title_bbox(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    long_title = "AI 产品经理如何设计可落地的 Agent 工作流"
    subtitle = "从用户需求、工具调用、状态管理到风险控制的技术产品分享"
    deck = Deck.model_validate(
        {
            "deck_id": "title_spacing_demo",
            "title": "Title Spacing Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Title Spacing Demo",
                    "layout": "title_slide",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": long_title,
                        },
                        {
                            "element_id": "subtitle",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 1.3, "width": 6.0, "height": 0.5},
                            "text": subtitle,
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "title_spacing.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    title_shape = _shape_with_text(rendered_slide, long_title)
    subtitle_shape = _shape_with_text(rendered_slide, subtitle)
    visible_texts = [
        shape.text.strip()
        for shape in rendered_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]

    assert subtitle_shape.top >= title_shape.top + title_shape.height + Inches(0.25)
    assert title_shape.left < Inches(1.0)
    assert title_shape.width >= Inches(7.5)
    assert visible_texts.count(long_title) == 1
    assert visible_texts.count(subtitle) == 1
    assert "Focus" not in visible_texts
    assert "Editable PPTX" not in visible_texts
    assert "TEMPLATE-GUIDED" not in visible_texts


def test_title_slide_does_not_duplicate_subtitle_in_side_panel(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    subtitle = "Building practical fluency, ethical judgment, and responsible AI habits"
    deck = Deck.model_validate(
        {
            "deck_id": "title_side_summary_demo",
            "title": "Title Side Summary Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Title Side Summary Demo",
                    "layout": "title_slide",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "AI Education for Practical Student Readiness",
                        },
                        {
                            "element_id": "subtitle",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 1.3, "width": 6.0, "height": 0.5},
                            "text": subtitle,
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "title_side_summary.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    subtitle_shapes = [
        shape
        for shape in rendered_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip() == subtitle
    ]

    assert len(subtitle_shapes) == 1
    assert subtitle_shapes[0].left < Inches(8.0)


def test_title_slide_keeps_side_panel_minimal_without_numbered_keyword_stack(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    title = "AI 产品经理如何设计 Agent 工作流"
    deck = Deck.model_validate(
        {
            "deck_id": "title_minimal_side_panel_demo",
            "title": "Title Minimal Side Panel Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Title Minimal Side Panel Demo",
                    "layout": "title_slide",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": title,
                        },
                        {
                            "element_id": "subtitle",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 1.3, "width": 6.0, "height": 0.5},
                            "text": "从边界、工作流、评估指标到落地风险",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "title_minimal_side_panel.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert visible_texts.count(title) == 1
    assert "技术产品视角" in visible_texts
    assert "01  工作流" not in visible_texts
    assert "02  边界" not in visible_texts


def test_rendered_slide_background_uses_pale_blue_green_theme(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "background_demo",
            "title": "Background Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [_template_slide_payload("two_column", 1)],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "background_demo.pptx")
    rendered_slide = Presentation(output_path).slides[0]

    assert str(rendered_slide.background.fill.fore_color.rgb) == "F3FBF8"


def test_four_cards_renders_heading_and_body_as_editable_text(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "four_cards_demo",
            "title": "Four Cards Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [_template_slide_payload("four_cards", 1)],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "four_cards.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    editable_texts = [
        shape.text.strip()
        for shape in rendered_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]

    assert "Discover" in editable_texts
    assert "Map user needs" in editable_texts
    assert "01" in editable_texts
    assert {"Action", "Anchor"} & set(editable_texts)
    assert not (INTERNAL_SURFACE_TERMS & set(editable_texts))
    assert "Prioritize" in editable_texts
    assert "Choose high-value work" in editable_texts


def test_four_cards_rebalances_one_overloaded_body_across_cards(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    overloaded_text = (
        "边界清单：先列禁止自动执行动作\n"
        "权限确认：高风险工具必须人工确认\n"
        "失败回退：异常时回到人工流程\n"
        "上线指标：用接管率和错误成本判断"
    )
    deck = Deck.model_validate(
        {
            "deck_id": "four_cards_rebalance_demo",
            "title": "Four Cards Rebalance Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "边界卡片",
                    "layout": "four_cards",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "边界卡片",
                        },
                        {
                            "element_id": "body",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.6},
                            "text": overloaded_text,
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "four_cards_rebalanced.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert overloaded_text not in visible_texts
    assert "边界清单：先列禁止自动执行动作" in visible_texts
    assert "权限确认：高风险工具必须人工确认" in visible_texts
    assert "失败回退：异常时回到人工流程" in visible_texts
    assert "上线指标：用接管率和错误成本判断" in visible_texts


def test_chinese_deck_uses_chinese_card_labels(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "chinese_labels_demo",
            "title": "AI 学习应用",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "AI 学习应用",
                    "layout": "four_cards",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "AI 学习应用",
                        },
                        {
                            "element_id": "card_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.5, "width": 3.0, "height": 1.0},
                            "text": "预习\n快速理解背景",
                        },
                        {
                            "element_id": "card_2",
                            "type": "text",
                            "bbox": {"x": 4.0, "y": 1.5, "width": 3.0, "height": 1.0},
                            "text": "练习\n获得即时反馈",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "chinese_labels.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "行动" in visible_texts
    assert "Action" not in visible_texts
    assert "Insight" not in visible_texts
    assert "Priority" not in visible_texts
    assert "Step" not in visible_texts


def test_sparse_two_column_layout_uses_compact_cards(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "sparse_two_column_demo",
            "title": "Sparse Two Column Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Sparse Two Column Demo",
                    "layout": "two_column",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "Two Focus Areas",
                        },
                        {
                            "element_id": "left",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.5, "width": 3.0, "height": 1.0},
                            "text": "Practice\nBuild fluency",
                        },
                        {
                            "element_id": "right",
                            "type": "text",
                            "bbox": {"x": 4.0, "y": 1.5, "width": 3.0, "height": 1.0},
                            "text": "Reflect\nUse judgment",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "sparse_two_column.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    card_backgrounds = [
        shape
        for shape in rendered_slide.shapes
        if not getattr(shape, "text", "").strip()
        and shape.width > Inches(5)
        and shape.height > Inches(1)
    ]

    assert len(card_backgrounds) == 2
    assert max(shape.height for shape in card_backgrounds) <= Inches(2.5)


def test_closing_slide_renders_multiple_next_steps_as_editable_text(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "closing_actions_demo",
            "title": "Closing Actions Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Closing Actions Demo",
                    "layout": "closing_slide",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "Next Steps",
                        },
                        {
                            "element_id": "actions",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.5, "width": 5.0, "height": 1.0},
                            "text": "Pilot the workflow\nGather feedback\nScale the pattern",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "closing_actions.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    editable_texts = [
        shape.text.strip()
        for shape in rendered_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]

    assert "01" in editable_texts
    assert "02" in editable_texts
    assert "03" in editable_texts
    for number in ["01", "02", "03"]:
        number_shape = _shape_with_text(rendered_slide, number)
        assert number_shape.width >= Inches(0.4)
    assert "Pilot the workflow" in editable_texts
    assert "Gather feedback" in editable_texts
    assert "Scale the pattern" in editable_texts


def test_template_rendering_does_not_expose_internal_surface_terms(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "surface_cleanup_demo",
            "title": "Surface Cleanup Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                _template_slide_payload(layout, index)
                for index, layout in enumerate(TEMPLATE_LAYOUTS, start=1)
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "surface_cleanup.pptx")
    visible_text = "\n".join(_visible_texts(Presentation(output_path)))

    for term in INTERNAL_SURFACE_TERMS:
        assert term not in visible_text


def test_professional_layouts_render_distinct_editable_templates(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    professional_layouts = [
        "comparison_matrix",
        "process_flow",
        "risk_matrix",
        "key_takeaway",
    ]
    deck = Deck.model_validate(
        {
            "deck_id": "professional_layout_demo",
            "title": "Professional Layout Demo",
            "theme_name": "clean_business",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                _template_slide_payload(layout, index)
                for index, layout in enumerate(professional_layouts, start=1)
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "professional_layouts.pptx")
    presentation = Presentation(output_path)

    assert len(presentation.slides) == 4
    visible_texts = _visible_texts(presentation)
    assert "Normal AI" in visible_texts
    assert "AI Agent" in visible_texts
    assert "Option A" not in visible_texts
    assert "Option B" not in visible_texts
    assert "01" in visible_texts
    assert "Issue" in visible_texts
    assert "Action" in visible_texts
    assert "Key Takeaway" in visible_texts
    for layout in professional_layouts:
        assert any(f"{layout} title" in text for text in visible_texts)


@pytest.mark.parametrize("step_count", [3, 4, 5])
def test_process_flow_renders_readable_step_counts(tmp_path: Path, step_count: int) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    elements = [
        {
            "element_id": "title",
            "type": "text",
            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
            "text": f"{step_count} Step Workflow",
        }
    ]
    for index in range(1, step_count + 1):
        elements.append(
            {
                "element_id": f"step_{index}",
                "type": "text",
                "bbox": {"x": 0.8, "y": 1.2 + index * 0.4, "width": 3.0, "height": 0.6},
                "text": f"Step {index}\nComplete focused action {index} without excessive detail",
            }
        )
    deck = Deck.model_validate(
        {
            "deck_id": f"process_flow_{step_count}_demo",
            "title": "Process Flow Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Process Flow Demo",
                    "layout": "process_flow",
                    "elements": elements,
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / f"process_{step_count}.pptx")
    presentation = Presentation(output_path)
    rendered_slide = presentation.slides[0]
    visible_texts = _visible_texts(presentation)
    step_numbers = {f"{index:02d}" for index in range(1, step_count + 1)}
    card_backgrounds = [
        shape
        for shape in rendered_slide.shapes
        if not getattr(shape, "text", "").strip()
        and shape.width >= Inches(2.5 if step_count == 4 else (3.5 if step_count == 5 else 3.0))
        and shape.height >= Inches(1.4)
    ]

    assert step_numbers <= set(visible_texts)
    assert len(card_backgrounds) >= step_count
    if step_count >= 4:
        for card in card_backgrounds:
            assert card.left >= 0
            assert card.top >= 0
            assert card.left + card.width <= presentation.slide_width
            assert card.top + card.height <= presentation.slide_height

        connectors = [
            shape
            for shape in rendered_slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.LINE
        ]
        core_text_shapes = [
            shape
            for shape in rendered_slide.shapes
            if getattr(shape, "has_text_frame", False)
            and "Complete focused action" in shape.text
        ]

        assert connectors
        assert core_text_shapes
        if step_count == 5:
            assert all(
                _overlap_area(connector, text_shape) == 0
                for connector in connectors
                for text_shape in core_text_shapes
            )


def test_visual_variant_selection_is_deterministic_and_index_sensitive() -> None:
    deck = Deck.model_validate(
        {
            "deck_id": "variant_demo",
            "title": "Variant Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [_template_slide_payload("three_column", 1)],
        }
    )
    slide = deck.slides[0]

    first = _visual_variant_for_slide(deck, slide, 1)
    second = _visual_variant_for_slide(deck, slide, 1)
    shifted = _visual_variant_for_slide(deck, slide, 2, variant_count=2)

    assert first == second
    assert 0 <= first < VISUAL_VARIANT_COUNTS["three_column"]
    assert shifted != _visual_variant_for_slide(deck, slide, 1, variant_count=2)
    for layout in [
        "title_slide",
        "two_column",
        "three_column",
        "four_cards",
        "metric_cards",
        "process_flow",
        "risk_matrix",
        "key_takeaway",
        "closing_slide",
    ]:
        assert VISUAL_VARIANT_COUNTS[layout] >= 2


def test_card_grid_layouts_have_distinct_visual_structures(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "card_grid_variety_demo",
            "title": "Card Grid Variety Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                _template_slide_payload("three_column", 1),
                _template_slide_payload("four_cards", 2),
                _template_slide_payload("metric_cards", 3),
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "card_grid_variety.pptx")
    presentation = Presentation(output_path)
    signatures = [
        _large_blank_rect_signature(slide)
        for slide in presentation.slides
    ]

    assert len(set(signatures)) == len(signatures)


def test_comparison_matrix_renders_aligned_rows(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "comparison_matrix_rows_demo",
            "title": "Comparison Matrix Rows Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [_template_slide_payload("comparison_matrix", 1)],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "comparison_rows.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "Compare" in visible_texts
    assert "判断" in visible_texts or "Point" in visible_texts
    assert "Normal AI" in visible_texts
    assert "AI Agent" in visible_texts
    assert "Input / Output" not in visible_texts
    assert "State" not in visible_texts
    assert "Option A" not in visible_texts
    assert "Option B" not in visible_texts
    assert "Decision 1" not in visible_texts
    assert "Decision 2" not in visible_texts
    assert any("workflow ownership" in text for text in visible_texts)


def test_metric_cards_renders_four_metrics_as_four_separate_cards(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    slide_payload = _template_slide_payload("metric_cards", 1)
    slide_payload["elements"].append(
        {
            "element_id": "s1_body_4",
            "type": "text",
            "bbox": {"x": 4.8, "y": 2.6, "width": 3.0, "height": 1.0},
            "text": "用户信任度\n观察解释清晰度和可控性感知",
        }
    )
    deck = Deck.model_validate(
        {
            "deck_id": "metric_cards_four_demo",
            "title": "Metric Cards Four Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [slide_payload],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "metric_cards_four.pptx")
    rendered_slide = Presentation(output_path).slides[0]
    visible_texts = _visible_texts(Presentation(output_path))
    card_backgrounds = [
        shape
        for shape in rendered_slide.shapes
        if not getattr(shape, "text", "").strip()
        and shape.width >= Inches(5.0)
        and shape.height >= Inches(1.5)
    ]

    assert "用户信任度" in visible_texts
    assert "观察解释清晰度和可控性感知" in visible_texts
    assert {"01", "02", "03", "04"} <= set(visible_texts)
    assert len(card_backgrounds) >= 4


@pytest.mark.parametrize("risk_count", [3, 4])
def test_risk_matrix_renders_three_to_four_risks(tmp_path: Path, risk_count: int) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    slide_payload = _template_slide_payload("risk_matrix", 1)
    slide_payload["elements"] = slide_payload["elements"][:1] + slide_payload["elements"][1 : risk_count + 1]
    deck = Deck.model_validate(
        {
            "deck_id": f"risk_matrix_{risk_count}_demo",
            "title": "Risk Matrix Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [slide_payload],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / f"risk_{risk_count}.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "Issue" in visible_texts
    assert "Effect" in visible_texts
    assert "Action" in visible_texts
    assert any("Hallucination" in text for text in visible_texts)


def test_risk_matrix_renders_fallback_mitigation_when_missing(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    slide_payload = _template_slide_payload("risk_matrix", 1)
    slide_payload["elements"][1]["text"] = "权限过大\n误操作影响用户数据"
    deck = Deck.model_validate(
        {
            "deck_id": "risk_matrix_missing_mitigation_demo",
            "title": "Risk Matrix Missing Mitigation Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [slide_payload],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "risk_missing_mitigation.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "权限过大" in visible_texts
    assert "误操作影响用户数据" in visible_texts
    assert "设置权限边界、人工确认和操作日志。" in visible_texts


def test_closing_slide_renders_action_items_as_heading_body_pairs(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "closing_action_pairs_demo",
            "title": "Closing Action Pairs Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Closing Action Pairs Demo",
                    "layout": "closing_slide",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "让 Agent 产品先可控，再扩张",
                        },
                        {
                            "element_id": "action_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.6},
                            "text": "01 定义边界",
                        },
                        {
                            "element_id": "action_2",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.0, "width": 5.0, "height": 0.6},
                            "text": "02 设计闭环\n把任务拆成输入、执行、校验、回滚和交付。",
                        },
                        {
                            "element_id": "action_3",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.6, "width": 5.0, "height": 0.6},
                            "text": "03 量化风险",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "closing_action_pairs.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "定义边界" in visible_texts
    assert "明确边界。" in visible_texts
    assert "设计闭环" in visible_texts
    assert "把任务拆成输入、执行、校验、回滚和交付。" in visible_texts
    assert "量化风险" in visible_texts
    assert any(text in visible_texts for text in {"记录失败样本。", "建立评估指标。"})
    assert "01 定义边界" not in visible_texts


def test_key_takeaway_renders_fallback_explanations(tmp_path: Path) -> None:
    theme = load_theme(EXAMPLES_DIR / "theme.json")
    deck = Deck.model_validate(
        {
            "deck_id": "key_takeaway_fallback_demo",
            "title": "Key Takeaway Fallback Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "Key Takeaway Fallback Demo",
                    "layout": "key_takeaway",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "Focus on workflow ownership",
                        },
                        {
                            "element_id": "takeaway_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.6},
                            "text": "Start with one journey\nUse evaluation before scaling",
                        },
                        {
                            "element_id": "takeaway_2",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.0, "width": 5.0, "height": 0.6},
                            "text": "Keep human checkpoints",
                        },
                    ],
                }
            ],
        }
    )

    output_path = render_deck_to_pptx(deck, theme, tmp_path / "key_takeaway_fallback.pptx")
    visible_texts = _visible_texts(Presentation(output_path))

    assert "Keep human checkpoints" in visible_texts
    assert any(
        text in visible_texts
        for text in {
            "Define the boundary and next checkpoint.",
            "Design one confirmation point for the workflow.",
            "Record failure samples and review checkpoints.",
            "Choose one launch metric and review it weekly.",
        }
    )
    assert "Turn this point into a concrete next action." not in visible_texts


def test_sanitize_deck_ir_for_render_removes_placeholder_risk_rows() -> None:
    deck = Deck.model_validate(
        {
            "deck_id": "sanitize_risk_rows_demo",
            "title": "Sanitize Risk Rows Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "风险矩阵",
                    "layout": "risk_matrix",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "风险矩阵",
                        },
                        {
                            "element_id": "row_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.8},
                            "text": "risk\nimpact\nmitigation",
                        },
                        {
                            "element_id": "row_2",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.3, "width": 5.0, "height": 0.8},
                            "text": "权限越界\n误操作用户数据\n限制高风险权限并要求人工确认",
                        },
                        {
                            "element_id": "row_3",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 3.2, "width": 5.0, "height": 0.8},
                            "text": "日志缺失\n问题无法追溯\n记录输入、工具调用和输出版本",
                        },
                    ],
                }
            ],
        }
    )

    sanitized, warnings = sanitize_deck_ir_for_render(deck)
    body_texts = [element.text for element in sanitized.slides[0].elements[1:]]

    assert all("risk\nimpact\nmitigation" not in text.lower() for text in body_texts)
    assert any("dropped placeholder-only risk matrix row" in warning for warning in warnings)


def test_sanitize_deck_ir_for_render_strips_risk_label_prefixes() -> None:
    deck = Deck.model_validate(
        {
            "deck_id": "sanitize_risk_prefix_demo",
            "title": "Sanitize Risk Prefix Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_025",
                    "title": "风险治理",
                    "layout": "risk_matrix",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "风险治理",
                        },
                        {
                            "element_id": "row_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.8},
                            "text": "risk：权限越界\nimpact：误操作用户数据\nmitigation：限制高风险权限并要求人工确认",
                        },
                        {
                            "element_id": "row_2",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.3, "width": 5.0, "height": 0.8},
                            "text": "Risk: 日志缺失\nImpact: 问题无法追溯\nMitigation: 记录输入、工具调用和输出版本",
                        },
                    ],
                }
            ],
        }
    )

    sanitized, warnings = sanitize_deck_ir_for_render(deck)
    body_text = "\n".join(element.text for element in sanitized.slides[0].elements[1:])

    assert "risk：" not in body_text.lower()
    assert "impact：" not in body_text.lower()
    assert "mitigation：" not in body_text.lower()
    assert "risk:" not in body_text.lower()
    assert "impact:" not in body_text.lower()
    assert "mitigation:" not in body_text.lower()
    assert "权限越界" in body_text
    assert "误操作用户数据" in body_text
    assert "限制高风险权限并要求人工确认" in body_text
    assert any(
        "slide_025: stripped risk/impact/mitigation label prefixes from risk_matrix elements" in warning
        for warning in warnings
    )


def test_render_long_deck_report_records_risk_label_prefix_sanitization(tmp_path: Path) -> None:
    deck = Deck.model_validate(
        {
            "deck_id": "render_risk_prefix_demo",
            "title": "Render Risk Prefix Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_025",
                    "title": "风险治理",
                    "layout": "risk_matrix",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "风险治理",
                        },
                        {
                            "element_id": "row_1",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.8},
                            "text": "risk：权限越界\nimpact：误操作用户数据\nmitigation：限制高风险权限并要求人工确认",
                        },
                        {
                            "element_id": "row_2",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 2.3, "width": 5.0, "height": 0.8},
                            "text": "模型幻觉\n用户据此做出错误判断\n设置事实校验和人工确认",
                        },
                    ],
                }
            ],
        }
    )
    input_path = tmp_path / "generated_long_deck_ir.json"
    output_path = tmp_path / "generated_long_deck.pptx"
    report_path = tmp_path / "long_deck_render_report.json"
    input_path.write_text(json.dumps(deck.model_dump(mode="json"), ensure_ascii=False), encoding="utf-8")

    report = render_long_deck_ir_to_pptx(
        input_path,
        output_path,
        report_path,
        theme_path=EXAMPLES_DIR / "theme.json",
        assets_dir=EXAMPLES_DIR,
    )

    saved_report = json.loads(report_path.read_text(encoding="utf-8"))
    assert report.status == "succeeded"
    assert output_path.exists()
    assert any("stripped risk/impact/mitigation label prefixes" in warning for warning in report.warnings)
    assert any("stripped risk/impact/mitigation label prefixes" in warning for warning in saved_report["warnings"])


def test_sanitize_deck_ir_for_render_falls_back_when_comparison_rows_are_placeholders() -> None:
    deck = Deck.model_validate(
        {
            "deck_id": "sanitize_comparison_rows_demo",
            "title": "Sanitize Comparison Rows Demo",
            "canvas_width_in": 13.333,
            "canvas_height_in": 7.5,
            "slides": [
                {
                    "slide_id": "slide_001",
                    "title": "责任对比",
                    "layout": "comparison_matrix",
                    "elements": [
                        {
                            "element_id": "title",
                            "type": "text",
                            "bbox": {"x": 0.5, "y": 0.5, "width": 6.0, "height": 0.6},
                            "text": "责任对比",
                        },
                        {
                            "element_id": "left",
                            "type": "text",
                            "bbox": {"x": 0.8, "y": 1.4, "width": 5.0, "height": 0.8},
                            "text": "基准侧\n判断点 1\n判断点 2",
                        },
                        {
                            "element_id": "right",
                            "type": "text",
                            "bbox": {"x": 6.1, "y": 1.4, "width": 5.0, "height": 0.8},
                            "text": "Agent 侧\n判断点 3",
                        },
                    ],
                }
            ],
        }
    )

    sanitized, warnings = sanitize_deck_ir_for_render(deck)

    assert sanitized.slides[0].layout == "two_column"
    assert any("comparison_matrix had fewer than two real comparison rows" in warning for warning in warnings)
    assert all(
        placeholder not in "\n".join(element.text for element in sanitized.slides[0].elements)
        for placeholder in {"基准侧", "Agent 侧", "判断点 1", "判断点 2", "判断点 3"}
    )
