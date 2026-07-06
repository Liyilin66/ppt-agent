import importlib.util
import json
import sys
from pathlib import Path

from pptx import Presentation

from ppt_agent.job_store import JobStore
from ppt_agent.ppt_master_execution import (
    PPT_MASTER_EXECUTION_PLAN_FILENAME,
    prepare_ppt_master_execution,
)
from ppt_agent.ppt_master_project import bootstrap_ppt_master_visual_project


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_prepare_execution_module():
    script_path = _repo_root() / "scripts" / "prepare_ppt_master_execution.py"
    spec = importlib.util.spec_from_file_location("prepare_ppt_master_execution", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Could not load script module from {script_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["prepare_ppt_master_execution"] = module
    spec.loader.exec_module(module)
    return module


def _mock_ppt_master_root(tmp_path: Path, *, include_skill: bool = True) -> Path:
    root = tmp_path / "ppt-master"
    skill_dir = root / "skills" / "ppt-master"
    skill_dir.mkdir(parents=True, exist_ok=True)
    if include_skill:
        (skill_dir / "SKILL.md").write_text("# PPT Master Skill\n", encoding="utf-8")
    (skill_dir / "scripts").mkdir(parents=True, exist_ok=True)
    (root / "requirements.txt").write_text("# test requirements\n", encoding="utf-8")
    return root


def _build_package(job_dir: Path) -> Path:
    package_dir = job_dir / "ppt_master_package"
    package_dir.mkdir(parents=True, exist_ok=True)
    (package_dir / "source.md").write_text("# Presentation Request\n", encoding="utf-8")
    (package_dir / "run_prompt.md").write_text("# PPT Master Local Job Prompt\n", encoding="utf-8")
    return package_dir


def _build_output(job_dir: Path, *, slide_count: int = 3) -> Path:
    output_dir = job_dir / "ppt_master_output"
    output_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    while len(presentation.slides) < slide_count:
        presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(output_dir / "generated_by_ppt_master.pptx")
    (output_dir / "generation_notes.md").write_text(
        "Native DrawingML shapes (directly editable)\n",
        encoding="utf-8",
    )
    return output_dir


def test_prepare_execution_waits_for_external_run_when_package_and_ppt_master_exist(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)

    plan = prepare_ppt_master_execution("job-123", job_dir, ppt_master_root=root)

    assert plan.status == "waiting_for_external_ppt_master_run"
    assert plan.ppt_master_root == root.resolve()
    assert plan.source_path == (job_dir / "ppt_master_package" / "source.md").resolve()
    assert plan.run_prompt_path == (job_dir / "ppt_master_package" / "run_prompt.md").resolve()
    assert plan.output_dir == (job_dir / "ppt_master_output").resolve()
    assert plan.expected_pptx_path.name == "generated_by_ppt_master.pptx"
    assert (job_dir / PPT_MASTER_EXECUTION_PLAN_FILENAME).exists()
    assert any("register_ppt_master_output.py" in step for step in plan.suggested_steps)


def test_prepare_execution_includes_bootstrapped_visual_project(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)
    project = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)

    plan = prepare_ppt_master_execution("job-123", job_dir, ppt_master_root=root)

    assert plan.status == "waiting_for_external_ppt_master_run"
    assert plan.project_dir == project.project_dir
    assert any("PROJECT_INSTRUCTIONS.md" in step for step in plan.suggested_steps)


def test_prepare_execution_detects_existing_output(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)
    _build_output(job_dir)

    plan = prepare_ppt_master_execution("job-123", job_dir, ppt_master_root=root)

    assert plan.status == "output_detected"
    assert plan.expected_pptx_path.is_file()


def test_prepare_execution_reports_missing_package(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"

    plan = prepare_ppt_master_execution("job-123", job_dir, ppt_master_root=root)

    assert plan.status == "missing_package"
    assert any("source document" in warning for warning in plan.warnings)
    assert any("run prompt" in warning for warning in plan.warnings)


def test_prepare_execution_reports_unavailable_ppt_master(tmp_path: Path) -> None:
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)

    plan = prepare_ppt_master_execution("job-123", job_dir, ppt_master_root=tmp_path / "missing")

    assert plan.status == "ppt_master_unavailable"
    assert plan.ppt_master_root == (tmp_path / "missing").resolve()
    assert plan.warnings


def test_prepare_execution_cli_writes_plan(tmp_path: Path, capsys) -> None:
    module = _load_prepare_execution_module()
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)

    exit_code = module.main(
        [
            "--job-id",
            "job-123",
            "--job-dir",
            str(job_dir),
            "--ppt-master-dir",
            str(root),
            "--db-path",
            str(tmp_path / "jobs.sqlite3"),
        ]
    )
    captured = capsys.readouterr()
    payload = json.loads((job_dir / PPT_MASTER_EXECUTION_PLAN_FILENAME).read_text(encoding="utf-8"))

    assert exit_code == 0
    assert "status: waiting_for_external_ppt_master_run" in captured.out
    assert payload["status"] == "waiting_for_external_ppt_master_run"
    assert payload["expected_pptx_path"].endswith("generated_by_ppt_master.pptx")


def test_prepare_execution_cli_registers_existing_output_when_job_exists(tmp_path: Path, capsys) -> None:
    module = _load_prepare_execution_module()
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "data" / "jobs" / "job-123"
    _build_package(job_dir)
    _build_output(job_dir, slide_count=2)
    db_path = tmp_path / "jobs.sqlite3"
    store = JobStore(db_path)
    job = store.create_job(job_type="long_deck")
    job_dir_for_real_id = tmp_path / "data" / "jobs" / job.job_id
    job_dir.rename(job_dir_for_real_id)

    exit_code = module.main(
        [
            "--job-id",
            job.job_id,
            "--job-dir",
            str(job_dir_for_real_id),
            "--ppt-master-dir",
            str(root),
            "--db-path",
            str(db_path),
        ]
    )
    captured = capsys.readouterr()
    artifacts = JobStore(db_path).list_artifacts(job.job_id)

    assert exit_code == 0
    assert "status: output_detected" in captured.out
    assert "registered_output_artifacts:" in captured.out
    assert {artifact.name for artifact in artifacts} >= {
        "ppt_master_generated_pptx",
        "ppt_master_generation_notes",
        "ppt_master_output_manifest",
    }
