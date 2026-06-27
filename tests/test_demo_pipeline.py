import json
from pathlib import Path

from pptx import Presentation

from scripts.run_demo_pipeline import run_demo_pipeline


def test_run_demo_pipeline_generates_expected_outputs(tmp_path: Path) -> None:
    repo_root = Path(__file__).resolve().parents[1]
    examples_dir = repo_root / "examples"

    outputs = run_demo_pipeline(output_dir=tmp_path, examples_dir=examples_dir)

    expected_files = {
        "qa_report": "qa_report.json",
        "sample_deck": "sample_deck.pptx",
        "patch_result": "patch_result.json",
        "patched_qa_report": "patched_qa_report.json",
        "patched_sample_deck": "patched_sample_deck.pptx",
    }

    assert set(outputs) == set(expected_files)
    for key, filename in expected_files.items():
        path = outputs[key]
        assert path == tmp_path / filename
        assert path.exists()
        assert path.stat().st_size > 0

    qa_report = json.loads((tmp_path / "qa_report.json").read_text(encoding="utf-8"))
    patch_result = json.loads((tmp_path / "patch_result.json").read_text(encoding="utf-8"))
    patched_qa_report = json.loads((tmp_path / "patched_qa_report.json").read_text(encoding="utf-8"))

    assert qa_report["deck_id"] == "sample_clean_business_deck"
    assert patch_result["applied_count"] == 3
    assert patch_result["issues"] == []
    assert patch_result["deck"]["slides"][0]["elements"][0]["text"] == "Updated Q3 Operating Review"
    assert patched_qa_report["deck_id"] == "sample_clean_business_deck"

    assert len(Presentation(tmp_path / "sample_deck.pptx").slides) == 3
    assert len(Presentation(tmp_path / "patched_sample_deck.pptx").slides) == 3
