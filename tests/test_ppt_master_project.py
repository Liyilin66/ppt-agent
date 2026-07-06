import importlib.util
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

import ppt_agent.ppt_master_runner as runner
from ppt_agent.ppt_master_project import (
    PROJECT_INSTRUCTIONS_FILENAME,
    PROJECT_MANIFEST_FILENAME,
    PPT_MASTER_VISUAL_PROJECT_MANIFEST_FILENAME,
    bootstrap_ppt_master_visual_project,
)
from ppt_agent.ppt_master_runner import run_ppt_master_local_export


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_bootstrap_cli_module():
    script_path = _repo_root() / "scripts" / "bootstrap_ppt_master_project.py"
    spec = importlib.util.spec_from_file_location("bootstrap_ppt_master_project", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Could not load script module from {script_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["bootstrap_ppt_master_project"] = module
    spec.loader.exec_module(module)
    return module


def _mock_ppt_master_root(tmp_path: Path, *, include_skill: bool = True) -> Path:
    root = tmp_path / "ppt-master"
    skill_dir = root / "skills" / "ppt-master"
    skill_dir.mkdir(parents=True, exist_ok=True)
    if include_skill:
        (skill_dir / "SKILL.md").write_text("# PPT Master Skill\n", encoding="utf-8")
    (skill_dir / "scripts").mkdir(parents=True, exist_ok=True)
    (root / "requirements.txt").write_text("# test requirements\n", encoding="utf-8")
    return root


def _build_package(job_dir: Path) -> Path:
    package_dir = job_dir / "ppt_master_package"
    package_dir.mkdir(parents=True, exist_ok=True)
    (package_dir / "source.md").write_text("# Presentation Request\n", encoding="utf-8")
    (package_dir / "run_prompt.md").write_text("# PPT Master Local Job Prompt\n", encoding="utf-8")
    return package_dir


def _write_pptx(path: Path, *, slide_count: int = 2) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    while len(presentation.slides) < slide_count:
        presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)


def test_bootstrap_reports_missing_package(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"

    project = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)

    assert project.status == "missing_package"
    assert (job_dir / PPT_MASTER_VISUAL_PROJECT_MANIFEST_FILENAME).exists()
    assert any("source.md" in warning for warning in project.warnings)
    assert not project.project_dir.exists()


def test_bootstrap_reports_unavailable_ppt_master(tmp_path: Path) -> None:
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)

    project = bootstrap_ppt_master_visual_project(
        "job-123",
        job_dir,
        ppt_master_root=tmp_path / "missing-ppt-master",
    )

    assert project.status == "ppt_master_unavailable"
    assert project.warnings
    assert (job_dir / PPT_MASTER_VISUAL_PROJECT_MANIFEST_FILENAME).exists()
    assert not project.project_dir.exists()


def test_bootstrap_creates_visual_project_scaffold(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    package_dir = _build_package(job_dir)

    project = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)
    manifest = json.loads(project.project_manifest_path.read_text(encoding="utf-8"))
    instructions = project.project_instructions_path.read_text(encoding="utf-8")

    assert project.status == "created"
    assert project.project_source_path.read_text(encoding="utf-8") == (package_dir / "source.md").read_text(
        encoding="utf-8"
    )
    assert project.project_prompt_path.read_text(encoding="utf-8") == (package_dir / "run_prompt.md").read_text(
        encoding="utf-8"
    )
    for directory_name in ["inputs", "svg_output", "svg_final", "exports", "logs"]:
        assert (project.project_dir / directory_name).is_dir()
    assert (project.project_dir / PROJECT_INSTRUCTIONS_FILENAME).is_file()
    assert (project.project_dir / PROJECT_MANIFEST_FILENAME).is_file()
    assert (job_dir / PPT_MASTER_VISUAL_PROJECT_MANIFEST_FILENAME).is_file()
    assert manifest["status"] == "created"
    assert str(root / "skills" / "ppt-master" / "SKILL.md") in instructions
    assert "inputs/source.md" in instructions
    assert "inputs/run_prompt.md" in instructions
    assert "svg_output" in instructions
    assert "svg_final" in instructions
    assert "generated_by_ppt_master.pptx" in instructions
    assert "has not generated SVG slides" in instructions


def test_bootstrap_already_exists_without_overwriting(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)
    first = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)
    marker = first.project_dir / "logs" / "keep.txt"
    marker.write_text("do not remove\n", encoding="utf-8")

    second = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)

    assert second.status == "already_exists"
    assert marker.read_text(encoding="utf-8") == "do not remove\n"


def test_bootstrap_cli_creates_scaffold(tmp_path: Path, capsys) -> None:
    module = _load_bootstrap_cli_module()
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)

    exit_code = module.main(
        [
            "--job-id",
            "job-123",
            "--job-dir",
            str(job_dir),
            "--ppt-master-dir",
            str(root),
            "--db-path",
            str(tmp_path / "jobs.sqlite3"),
        ]
    )
    captured = capsys.readouterr()

    assert exit_code == 0
    assert "status: created" in captured.out
    assert "PROJECT_INSTRUCTIONS.md" in captured.out
    assert (job_dir / "ppt_master_output" / "ppt_master_visual_project" / PROJECT_INSTRUCTIONS_FILENAME).exists()


def test_runner_recognizes_bootstrap_project_without_svg_inputs(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)
    project = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)

    result = run_ppt_master_local_export("job-123", job_dir, ppt_master_root=root, register_output=False)

    assert result.status == "requires_external_ai_generation"
    assert result.project_dir == project.project_dir
    assert any("No .svg files found" in warning for warning in result.warnings)


def test_runner_exports_from_bootstrap_project_after_svg_appears(tmp_path: Path, monkeypatch) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)
    project = bootstrap_ppt_master_visual_project("job-123", job_dir, ppt_master_root=root)
    (project.expected_svg_output_dir / "P01_cover.svg").write_text(
        '<svg viewBox="0 0 1280 720" xmlns="http://www.w3.org/2000/svg"><text>Test</text></svg>',
        encoding="utf-8",
    )
    commands: list[list[str]] = []

    def fake_run(command, cwd, check, capture_output, text, timeout):  # noqa: ANN001
        commands.append(command)
        if command[1].endswith("svg_to_pptx.py"):
            output_path = Path(command[command.index("-o") + 1])
            _write_pptx(output_path, slide_count=2)
        return subprocess.CompletedProcess(command, 0, stdout="ok\n", stderr="")

    monkeypatch.setattr(runner.subprocess, "run", fake_run)

    result = run_ppt_master_local_export("job-123", job_dir, ppt_master_root=root, register_output=False)

    assert result.status == "succeeded"
    assert result.project_dir == project.project_dir
    assert result.slide_count == 2
    assert {Path(command[1]).name for command in commands} == {
        "svg_quality_checker.py",
        "finalize_svg.py",
        "svg_to_pptx.py",
    }
