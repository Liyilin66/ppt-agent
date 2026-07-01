import importlib.util
import json
import sys
from pathlib import Path

from pptx import Presentation

from ppt_agent.job_store import JobStore
from ppt_agent.ppt_master_output import (
    PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
    PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
    detect_ppt_master_output,
    register_ppt_master_output_artifacts,
)


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_register_module():
    script_path = _repo_root() / "scripts" / "register_ppt_master_output.py"
    spec = importlib.util.spec_from_file_location("register_ppt_master_output", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Could not load script module from {script_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["register_ppt_master_output"] = module
    spec.loader.exec_module(module)
    return module


def _build_mock_output_dir(base_dir: Path, *, slide_count: int = 3) -> Path:
    output_dir = base_dir / "data" / "jobs" / "job-123" / "ppt_master_output"
    output_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    while len(presentation.slides) < slide_count:
        layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(layout)
        textbox = slide.shapes.add_textbox(50, 50, 400, 80)
        textbox.text_frame.text = f"Slide {len(presentation.slides)}"
    if len(presentation.slides) > slide_count:
        while len(presentation.slides) > slide_count:
            slide_id = presentation.slides._sldIdLst[-1].rId  # type: ignore[attr-defined]
            presentation.part.drop_rel(slide_id)
            del presentation.slides._sldIdLst[-1]  # type: ignore[attr-defined]

    pptx_path = output_dir / "generated_by_ppt_master.pptx"
    presentation.save(pptx_path)

    (output_dir / "generation_notes.md").write_text(
        "# Notes\n\nMode: Native DrawingML shapes (directly editable)\n",
        encoding="utf-8",
    )
    project_dir = output_dir / "agent_pm_recovery_ppt169_20260701"
    (project_dir / "svg_output").mkdir(parents=True, exist_ok=True)
    return output_dir


def test_detect_ppt_master_output_accepts_valid_output_dir(tmp_path: Path) -> None:
    output_dir = _build_mock_output_dir(tmp_path, slide_count=4)

    manifest = detect_ppt_master_output(output_dir)

    assert manifest.detected is True
    assert manifest.output_dir == output_dir.resolve()
    assert manifest.pptx_path == (output_dir / "generated_by_ppt_master.pptx").resolve()
    assert manifest.notes_path == (output_dir / "generation_notes.md").resolve()
    assert manifest.project_dir == (output_dir / "agent_pm_recovery_ppt169_20260701").resolve()
    assert manifest.slide_count == 4
    assert manifest.is_editable_claimed is True
    assert manifest.generation_status == "succeeded"
    assert manifest.job_id == "job-123"


def test_detect_ppt_master_output_reports_missing_pptx(tmp_path: Path) -> None:
    output_dir = tmp_path / "data" / "jobs" / "job-123" / "ppt_master_output"
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "generation_notes.md").write_text("# Notes\n", encoding="utf-8")

    manifest = detect_ppt_master_output(output_dir)

    assert manifest.detected is False
    assert manifest.pptx_path is None
    assert manifest.generation_status == "missing_pptx"
    assert any("PPTX output not found" in warning for warning in manifest.warnings)


def test_register_ppt_master_output_artifacts_creates_expected_records(tmp_path: Path) -> None:
    output_dir = _build_mock_output_dir(tmp_path, slide_count=2)
    store = JobStore(tmp_path / "jobs.sqlite3")
    job = store.create_job(job_type="long_deck")

    result = register_ppt_master_output_artifacts(store, job_id=job.job_id, output_dir=output_dir)
    artifact_names = {artifact.name for artifact in store.list_artifacts(job.job_id)}

    assert result.manifest.detected is True
    assert result.pptx_artifact is not None
    assert result.notes_artifact is not None
    assert result.manifest_artifact.name == PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT
    assert artifact_names >= {
        PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
        PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
        PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    }
    manifest_payload = json.loads(result.manifest_path.read_text(encoding="utf-8"))
    assert manifest_payload["slide_count"] == 2


def test_register_ppt_master_output_cli_writes_manifest_and_registers_artifacts(
    tmp_path: Path,
    capsys,
) -> None:
    module = _load_register_module()
    db_path = tmp_path / "jobs.sqlite3"
    store = JobStore(db_path)
    job = store.create_job(job_type="long_deck")
    output_dir = _build_mock_output_dir(tmp_path / "workspace")
    expected_output_dir = tmp_path / "data" / "jobs" / job.job_id / "ppt_master_output"
    expected_output_dir.parent.mkdir(parents=True, exist_ok=True)
    output_dir.rename(expected_output_dir)
    output_dir = expected_output_dir

    exit_code = module.main(
        [
            "--job-id",
            job.job_id,
            "--output-dir",
            str(output_dir),
            "--db-path",
            str(db_path),
        ]
    )
    captured = capsys.readouterr()
    refreshed_store = JobStore(db_path)
    artifacts = refreshed_store.list_artifacts(job.job_id)

    assert exit_code == 0
    assert "ppt_master_output_manifest:" in captured.out
    assert (output_dir / "ppt_master_output_manifest.json").exists()
    assert {artifact.name for artifact in artifacts} >= {
        PPT_MASTER_OUTPUT_PPTX_ARTIFACT,
        PPT_MASTER_OUTPUT_NOTES_ARTIFACT,
        PPT_MASTER_OUTPUT_MANIFEST_ARTIFACT,
    }


def test_register_ppt_master_output_cli_reports_missing_pptx(tmp_path: Path, capsys) -> None:
    module = _load_register_module()
    db_path = tmp_path / "jobs.sqlite3"
    store = JobStore(db_path)
    job = store.create_job(job_type="long_deck")
    output_dir = tmp_path / "data" / "jobs" / job.job_id / "ppt_master_output"
    output_dir.mkdir(parents=True, exist_ok=True)

    exit_code = module.main(
        [
            "--job-id",
            job.job_id,
            "--output-dir",
            str(output_dir),
            "--db-path",
            str(db_path),
        ]
    )
    captured = capsys.readouterr()

    assert exit_code == 2
    assert "Could not register PPT Master output" in captured.err
    assert "generated_by_ppt_master.pptx" in captured.err


def test_register_ppt_master_output_cli_rejects_short_deck_job(tmp_path: Path, capsys) -> None:
    module = _load_register_module()
    db_path = tmp_path / "jobs.sqlite3"
    store = JobStore(db_path)
    job = store.create_job(job_type="short_deck")
    output_dir = _build_mock_output_dir(tmp_path / "workspace", slide_count=2)
    expected_output_dir = tmp_path / "data" / "jobs" / job.job_id / "ppt_master_output"
    expected_output_dir.parent.mkdir(parents=True, exist_ok=True)
    output_dir.rename(expected_output_dir)
    output_dir = expected_output_dir

    exit_code = module.main(
        [
            "--job-id",
            job.job_id,
            "--output-dir",
            str(output_dir),
            "--db-path",
            str(db_path),
        ]
    )
    captured = capsys.readouterr()

    assert exit_code == 2
    assert "is not a long_deck job" in captured.err
