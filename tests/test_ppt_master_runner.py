import importlib.util
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

import ppt_agent.ppt_master_runner as runner
from ppt_agent.job_store import JobStore
from ppt_agent.ppt_master_runner import (
    PPT_MASTER_RUNNER_RESULT_ARTIFACT,
    PPT_MASTER_RUNNER_RESULT_FILENAME,
    run_ppt_master_local_export,
)


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_runner_cli_module():
    script_path = _repo_root() / "scripts" / "run_ppt_master_local_export.py"
    spec = importlib.util.spec_from_file_location("run_ppt_master_local_export", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Could not load script module from {script_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["run_ppt_master_local_export"] = module
    spec.loader.exec_module(module)
    return module


def _mock_ppt_master_root(tmp_path: Path) -> Path:
    root = tmp_path / "ppt-master"
    scripts_dir = root / "skills" / "ppt-master" / "scripts"
    scripts_dir.mkdir(parents=True, exist_ok=True)
    (root / "skills" / "ppt-master" / "SKILL.md").write_text("# PPT Master Skill\n", encoding="utf-8")
    (root / "requirements.txt").write_text("# test requirements\n", encoding="utf-8")
    return root


def _build_package(job_dir: Path) -> None:
    package_dir = job_dir / "ppt_master_package"
    package_dir.mkdir(parents=True, exist_ok=True)
    (package_dir / "source.md").write_text("# Presentation Request\n", encoding="utf-8")
    (package_dir / "run_prompt.md").write_text("# PPT Master Local Job Prompt\n", encoding="utf-8")


def _build_project(job_dir: Path) -> Path:
    project_dir = job_dir / "ppt_master_output" / "agent_pm_recovery_ppt169_20260701"
    svg_output = project_dir / "svg_output"
    svg_output.mkdir(parents=True, exist_ok=True)
    (svg_output / "P01_cover.svg").write_text(
        '<svg viewBox="0 0 1280 720" xmlns="http://www.w3.org/2000/svg"><text>Test</text></svg>',
        encoding="utf-8",
    )
    (project_dir / "spec_lock.md").write_text("# Execution Lock\n", encoding="utf-8")
    return project_dir


def _write_pptx(path: Path, *, slide_count: int = 3) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    while len(presentation.slides) < slide_count:
        presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)


def test_runner_reports_missing_package(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"

    result = run_ppt_master_local_export("job-123", job_dir, ppt_master_root=root, register_output=False)

    assert result.status == "missing_package"
    assert (job_dir / PPT_MASTER_RUNNER_RESULT_FILENAME).exists()
    assert any("ppt_master_package" in warning for warning in result.warnings)


def test_runner_requires_external_generation_when_package_exists_but_project_missing(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)

    result = run_ppt_master_local_export("job-123", job_dir, ppt_master_root=root, register_output=False)

    assert result.status == "requires_external_ai_generation"
    assert result.project_dir is None
    assert "AI IDE" in result.message


def test_runner_reports_unavailable_ppt_master(tmp_path: Path) -> None:
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)

    result = run_ppt_master_local_export(
        "job-123",
        job_dir,
        ppt_master_root=tmp_path / "missing-ppt-master",
        register_output=False,
    )

    assert result.status == "ppt_master_unavailable"
    assert result.warnings


def test_runner_builds_commands_and_registers_successful_mock_export(
    tmp_path: Path,
    monkeypatch,
) -> None:
    root = _mock_ppt_master_root(tmp_path)
    store = JobStore(tmp_path / "jobs.sqlite3")
    job = store.create_job(job_type="long_deck")
    job_dir = tmp_path / "jobs" / job.job_id
    _build_package(job_dir)
    project_dir = _build_project(job_dir)
    commands: list[list[str]] = []

    def fake_run(command, cwd, check, capture_output, text, timeout):  # noqa: ANN001
        commands.append(command)
        if command[1].endswith("svg_to_pptx.py"):
            output_path = Path(command[command.index("-o") + 1])
            _write_pptx(output_path, slide_count=3)
        return subprocess.CompletedProcess(command, 0, stdout="ok\n", stderr="")

    monkeypatch.setattr(runner.subprocess, "run", fake_run)

    result = run_ppt_master_local_export(
        job.job_id,
        job_dir,
        ppt_master_root=root,
        project_dir=project_dir,
        store=store,
    )
    artifact_names = {artifact.name for artifact in store.list_artifacts(job.job_id)}

    assert result.status == "succeeded"
    assert result.project_dir == project_dir.resolve()
    assert result.slide_count == 3
    assert result.registered is True
    assert {Path(command[1]).name for command in commands} == {
        "svg_quality_checker.py",
        "finalize_svg.py",
        "svg_to_pptx.py",
    }
    assert any("--only" in command and "native" in command for command in commands)
    assert artifact_names >= {
        "ppt_master_generated_pptx",
        "ppt_master_generation_notes",
        "ppt_master_output_manifest",
    }


def test_runner_result_artifact_can_be_registered(tmp_path: Path) -> None:
    root = _mock_ppt_master_root(tmp_path)
    store = JobStore(tmp_path / "jobs.sqlite3")
    job = store.create_job(job_type="long_deck")
    job_dir = tmp_path / "jobs" / job.job_id

    result = run_ppt_master_local_export(job.job_id, job_dir, ppt_master_root=root, store=store)
    artifact = runner.register_ppt_master_runner_result_artifact(store, job_id=job.job_id, job_dir=job_dir)

    assert result.status == "missing_package"
    assert artifact.name == PPT_MASTER_RUNNER_RESULT_ARTIFACT
    assert artifact.kind == "json"


def test_runner_cli_prints_requires_external_generation(tmp_path: Path, capsys) -> None:
    module = _load_runner_cli_module()
    root = _mock_ppt_master_root(tmp_path)
    job_dir = tmp_path / "jobs" / "job-123"
    _build_package(job_dir)

    exit_code = module.main(
        [
            "--job-id",
            "job-123",
            "--job-dir",
            str(job_dir),
            "--ppt-master-dir",
            str(root),
            "--db-path",
            str(tmp_path / "jobs.sqlite3"),
        ]
    )
    captured = capsys.readouterr()

    assert exit_code == 0
    assert "status: requires_external_ai_generation" in captured.out
    assert "generate the visual project first" in captured.out
